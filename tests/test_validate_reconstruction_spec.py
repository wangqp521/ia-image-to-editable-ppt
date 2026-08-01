from __future__ import annotations

import contextlib
import copy
import hashlib
import importlib.util
import io
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from tests.fixture_specs import make_asset_fallback_spec, make_minimal_spec


SCRIPT_PATH = Path(__file__).resolve().parents[1] / "scripts" / "validate_reconstruction_spec.py"
SPEC = importlib.util.spec_from_file_location("validate_reconstruction_spec", SCRIPT_PATH)
if SPEC is None or SPEC.loader is None:
    raise RuntimeError(f"Cannot load {SCRIPT_PATH}")
MODULE = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(MODULE)

PPTX_VALIDATOR_PATH = Path(__file__).resolve().parents[1] / "scripts" / "validate_pptx.py"
PPTX_SPEC = importlib.util.spec_from_file_location("validate_pptx_for_spec_tests", PPTX_VALIDATOR_PATH)
if PPTX_SPEC is None or PPTX_SPEC.loader is None:
    raise RuntimeError(f"Cannot load {PPTX_VALIDATOR_PATH}")
PPTX_VALIDATOR = importlib.util.module_from_spec(PPTX_SPEC)
PPTX_SPEC.loader.exec_module(PPTX_VALIDATOR)

VISUAL_DIFF_PATH = Path(__file__).resolve().parents[1] / "scripts" / "create_visual_diff.py"
VISUAL_DIFF_SPEC = importlib.util.spec_from_file_location("visual_diff_for_spec_tests", VISUAL_DIFF_PATH)
if VISUAL_DIFF_SPEC is None or VISUAL_DIFF_SPEC.loader is None:
    raise RuntimeError(f"Cannot load {VISUAL_DIFF_PATH}")
VISUAL_DIFF = importlib.util.module_from_spec(VISUAL_DIFF_SPEC)
VISUAL_DIFF_SPEC.loader.exec_module(VISUAL_DIFF)

COORDINATE_OVERLAY_PATH = Path(__file__).resolve().parents[1] / "scripts" / "create_coordinate_overlay.py"
COORDINATE_OVERLAY_SPEC = importlib.util.spec_from_file_location(
    "coordinate_overlay_for_spec_tests",
    COORDINATE_OVERLAY_PATH,
)
if COORDINATE_OVERLAY_SPEC is None or COORDINATE_OVERLAY_SPEC.loader is None:
    raise RuntimeError(f"Cannot load {COORDINATE_OVERLAY_PATH}")
COORDINATE_OVERLAY = importlib.util.module_from_spec(COORDINATE_OVERLAY_SPEC)
COORDINATE_OVERLAY_SPEC.loader.exec_module(COORDINATE_OVERLAY)

REFERENCE_ROOT: tempfile.TemporaryDirectory[str] | None = None
REFERENCE_PATH: Path | None = None
COORDINATE_EVIDENCE: dict | None = None


def setUpModule() -> None:
    global REFERENCE_ROOT, REFERENCE_PATH, COORDINATE_EVIDENCE
    REFERENCE_ROOT = tempfile.TemporaryDirectory()
    REFERENCE_PATH = Path(REFERENCE_ROOT.name) / "source.png"
    Image.new("RGB", (1600, 900), "white").save(REFERENCE_PATH)
    overlay_path = Path(REFERENCE_ROOT.name) / "coordinate-overlay.png"
    report = COORDINATE_OVERLAY.create_coordinate_overlay(
        REFERENCE_PATH,
        overlay_path,
    )
    COORDINATE_EVIDENCE = {
        "path": str(overlay_path.resolve()),
        "sha256": hashlib.sha256(overlay_path.read_bytes()).hexdigest(),
        "source_sha256": report["source"]["sha256"],
        "manifest_sha256": report["coordinate_overlay_manifest_sha256"],
        "grid": report["grid"],
        "inspection": "passed",
    }


def tearDownModule() -> None:
    global REFERENCE_ROOT, REFERENCE_PATH, COORDINATE_EVIDENCE
    if REFERENCE_ROOT is not None:
        REFERENCE_ROOT.cleanup()
    REFERENCE_ROOT = None
    REFERENCE_PATH = None
    COORDINATE_EVIDENCE = None


def image_identity(path: Path) -> dict:
    return {
        "path": str(path.resolve()),
        "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
    }


def valid_spec() -> dict:
    text = "标题"
    if REFERENCE_PATH is None or COORDINATE_EVIDENCE is None:
        raise RuntimeError("reference fixture is not initialized")
    reference = image_identity(REFERENCE_PATH)
    return {
        "schema_version": 2,
        "page_id": "page-001",
        "verification_profile": "rapid",
        "delivery_status": "pending",
        "session_reuse": {
            "mode": "fresh_reconstruction",
            "reason": "new_session",
            "artifacts": [],
        },
        "content_reference": dict(reference),
        "clean_visual_reference": dict(reference),
        "canvas": {
            "source_size": [1600, 900],
            "visual_size": [1600, 900],
            "page_frame_bbox": [0, 0, 1600, 900],
            "slide_size_emu": [12192000, 6858000],
            "mapping_mode": "direct_16_9",
            "background": "#FFFFFF",
        },
        "activated_modules": [
            "page_layout",
            "typography",
            "representation_plan",
            "background",
        ],
        "modules": {
            "page_layout": {
                "anchors": [],
                "relationships": [],
                "layout_invariants": [],
                "density_targets": {},
                "coordinate_overlay_evidence": copy.deepcopy(COORDINATE_EVIDENCE),
            },
            "typography": {
                "slide_coordinate_unit": "EMU",
                "items": [
                    {
                        "element_id": "title",
                        "text": text,
                        "source_font_guess": "Noto Sans CJK SC",
                        "selected_font": "Noto Sans CJK SC",
                        "fallback_reason": None,
                        "fallback_trace": None,
                        "runs": [
                            {
                                "start": 0,
                                "end": len(text),
                                "font_size": 24,
                                "font_weight": 700,
                                "color": "#000000",
                                "letter_spacing": 0,
                                "italic": False,
                                "underline": False,
                                "strike": False,
                                "baseline": 0,
                            }
                        ],
                        "paragraphs": [
                            {
                                "start": 0,
                                "end": len(text),
                                "alignment": "left",
                                "line_spacing": 1.0,
                                "space_before": 0,
                                "space_after": 0,
                                "indent": 0,
                                "list": {"is_list": False, "level": 0, "bullet": None},
                            }
                        ],
                        "text_box": {
                            "x": 228600,
                            "y": 228600,
                            "w": 6096000,
                            "h": 457200,
                            "margins": {"left": 0, "right": 0, "top": 0, "bottom": 0},
                            "alignment": "left",
                            "vertical_alignment": "top",
                            "wrap": False,
                            "overflow": False,
                            "soft_breaks": [],
                            "paragraph_breaks": [],
                        },
                        "internal_font_declaration": "Noto Sans CJK SC",
                        "font_declaration_verified": False,
                    }
                ],
            },
            "representation_plan": {
                "items": [
                    {
                        "source_fact_id": "fact-title-001",
                        "semantic_role": "title",
                        "source_bbox": [30, 30, 800, 60],
                        "required": True,
                        "selected_mode": "native",
                        "required_editability": "full",
                        "fallback_policy": "forbid",
                        "bound_element_ids": ["title"],
                        "reason": "title is represented as editable native text",
                        "coverage_status": "covered",
                        "evidence": ["representation-plan.json"],
                    }
                ]
            },
            "background": {
                "items": [
                    {
                        "background_id": "background-001",
                        "role": "base",
                        "source_bbox": [0, 0, 1600, 900],
                        "selected_mode": "native",
                        "bound_element_id": "background-base",
                        "source_provenance": {
                            "kind": "native_measurement",
                            "source_path": reference["path"],
                            "source_sha256": reference["sha256"],
                        },
                        "reason": "measured solid page background",
                        "evidence": [reference["path"]],
                        "contains_foreground_semantics": False,
                    }
                ]
            },
        },
        "regions": [
            {
                "region_id": "header",
                "source_bbox": [0, 0, 1600, 120],
                "slide_bbox": [0, 0, 12192000, 914400],
                "layer": 1,
                "padding": {"left": 0, "right": 0, "top": 0, "bottom": 0},
                "element_ids": ["title", "background-base"],
            }
        ],
        "elements": [
            {
                "element_id": "title",
                "kind": "text",
                "source_bbox": [30, 30, 800, 60],
                "slide_bbox": [228600, 228600, 6096000, 457200],
                "layer": 10,
                "editable": True,
                "confidence": "high",
                "style": {"fill": "noFill"},
                "content": {"text": text},
            },
            {
                "element_id": "background-base",
                "kind": "shape",
                "source_bbox": [0, 0, 1600, 900],
                "slide_bbox": [0, 0, 12192000, 6858000],
                "layer": 0,
                "editable": True,
                "confidence": "high",
                "style": {
                    "shape_type": "rectangle",
                    "fill": {"type": "solid", "color": "#FFFFFF", "opacity": 1},
                    "effects": "none",
                    "rotation": 0,
                },
                "content": {},
            },
        ],
        "reading_order": ["background-base", "title"],
        "visual_gate": {"status": "pending", "evidence": [], "tripwire": None},
        "editability_gate": {"status": "pending", "evidence": []},
    }


def valid_list_spec() -> dict:
    candidate = valid_spec()
    text = "第一项第二项"
    candidate["elements"][0]["content"]["text"] = text
    item = candidate["modules"]["typography"]["items"][0]
    item["text"] = text
    item["runs"][0]["end"] = len(text)
    item["paragraphs"] = [
        {
            "start": 0,
            "end": 3,
            "alignment": "left",
            "line_spacing": 1.0,
            "space_before": 0,
            "space_after": 0,
            "margin_left": 342900,
            "indent": -228600,
            "list": {
                "is_list": True,
                "level": 0,
                "bullet_type": "char",
                "bullet": "•",
                "bullet_font": "follow_text",
                "bullet_size_mode": "follow_text",
                "bullet_size_value": None,
                "bullet_color": "follow_text",
            },
        },
        {
            "start": 3,
            "end": len(text),
            "alignment": "left",
            "line_spacing": 1.0,
            "space_before": 0,
            "space_after": 0,
            "margin_left": 342900,
            "indent": -228600,
            "list": {
                "is_list": True,
                "level": 0,
                "bullet_type": "char",
                "bullet": "•",
                "bullet_font": "follow_text",
                "bullet_size_mode": "follow_text",
                "bullet_size_value": None,
                "bullet_color": "follow_text",
            },
        },
    ]
    item["text_box"]["paragraph_breaks"] = [3]
    return candidate


def append_native_element(
    candidate: dict,
    *,
    element_id: str,
    kind: str,
    style: dict,
) -> None:
    source_bbox = [30, 120, 200, 100]
    slide_bbox = [228600, 914400, 1524000, 762000]
    candidate["elements"].append(
        {
            "element_id": element_id,
            "kind": kind,
            "source_bbox": source_bbox,
            "slide_bbox": slide_bbox,
            "layer": 2,
            "editable": True,
            "confidence": "high",
            "style": style,
            "content": {},
        }
    )
    candidate["regions"][0]["element_ids"].append(element_id)
    candidate["reading_order"].append(element_id)
    candidate["modules"]["representation_plan"]["items"].append(
        {
            "source_fact_id": f"fact-{element_id}",
            "semantic_role": kind,
            "source_bbox": source_bbox,
            "required": True,
            "selected_mode": "native",
            "required_editability": "full",
            "fallback_policy": "forbid",
            "bound_element_ids": [element_id],
            "reason": "renderer contract regression fixture",
            "coverage_status": "covered",
            "evidence": ["representation-plan.json"],
        }
    )


class ValidateReconstructionSpecTests(unittest.TestCase):
    def _artifact(self, path: Path, payload: bytes) -> dict:
        path.write_bytes(payload)
        return {"path": str(path.resolve()), "sha256": hashlib.sha256(payload).hexdigest()}

    def assertAuthoringMatchesPrebuild(self, candidate: dict) -> tuple[dict, dict]:
        authoring = MODULE.validate_spec(candidate, stage="authoring")
        prebuild = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertEqual("authoring", authoring["stage"])
        self.assertEqual("prebuild", prebuild["stage"])
        self.assertEqual(
            {key: value for key, value in authoring.items() if key != "stage"},
            {key: value for key, value in prebuild.items() if key != "stage"},
        )
        return authoring, prebuild

    def test_authoring_and_prebuild_are_equivalent_for_valid_spec(self):
        authoring, prebuild = self.assertAuthoringMatchesPrebuild(valid_spec())

        self.assertTrue(authoring["valid"], authoring)
        self.assertTrue(prebuild["valid"], prebuild)

    def test_authoring_and_prebuild_are_equivalent_for_vertical_alignment_mismatch(self):
        candidate = valid_spec()
        candidate["elements"][0]["style"]["vertical_alignment"] = "middle"

        authoring, prebuild = self.assertAuthoringMatchesPrebuild(candidate)

        expected = {
            "code": "UNSUPPORTED_CAPABILITY",
            "path": "elements.title.style.vertical_alignment",
            "detail": "text style vertical_alignment must match typography text_box",
        }
        self.assertFalse(authoring["valid"])
        self.assertIn(expected, authoring["errors"])
        self.assertIn(expected, prebuild["errors"])

    def test_authoring_and_prebuild_are_equivalent_for_legacy_text_run(self):
        candidate = valid_spec()
        run = candidate["modules"]["typography"]["items"][0]["runs"][0]
        for field in ("italic", "underline", "strike", "baseline"):
            del run[field]
        run["decoration"] = "none"

        authoring, prebuild = self.assertAuthoringMatchesPrebuild(candidate)

        self.assertFalse(authoring["valid"])
        codes = {issue["code"] for issue in prebuild["errors"]}
        self.assertIn("SPEC_TEXT_RUN_STYLE_INVALID", codes)
        self.assertIn("SPEC_TYPOGRAPHY_FIELD_UNKNOWN", codes)

    def test_authoring_and_prebuild_are_equivalent_for_representation_gap(self):
        candidate = valid_spec()
        candidate["modules"]["representation_plan"]["items"] = []

        authoring, prebuild = self.assertAuthoringMatchesPrebuild(candidate)

        self.assertFalse(authoring["valid"])
        self.assertIn(
            "REPRESENTATION_INCOMPLETE",
            {issue["code"] for issue in prebuild["errors"]},
        )

    def test_validate_spec_rejects_unknown_stage_with_complete_choices(self):
        with self.assertRaisesRegex(
            ValueError,
            "stage must be authoring, prebuild, or final",
        ):
            MODULE.validate_spec(valid_spec(), stage="draft")

    def test_cli_accepts_authoring_and_writes_authoring_report(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            spec_path = root / "page-reconstruction.json"
            report_path = root / "authoring-validation.json"
            spec_path.write_text(json.dumps(valid_spec()), encoding="utf-8")
            with contextlib.redirect_stdout(io.StringIO()):
                exit_code = MODULE.main(
                    [
                        str(spec_path),
                        "--stage",
                        "authoring",
                        "--output",
                        str(report_path),
                    ]
                )

            self.assertEqual(0, exit_code)
            self.assertTrue(report_path.is_file())
            report = json.loads(report_path.read_text(encoding="utf-8"))
            self.assertTrue(report["valid"], report)
            self.assertEqual("authoring", report["stage"])

    def test_cli_rejects_unknown_stage_with_choice_diagnostic(self):
        stderr = io.StringIO()
        with contextlib.redirect_stderr(stderr):
            with self.assertRaises(SystemExit) as raised:
                MODULE.main(["page-reconstruction.json", "--stage", "draft"])

        self.assertEqual(2, raised.exception.code)
        diagnostic = stderr.getvalue()
        self.assertIn("argument --stage: invalid choice: 'draft'", diagnostic)
        for choice in ("authoring", "prebuild", "final"):
            self.assertIn(choice, diagnostic)

    def test_prebuild_accepts_all_explicit_verification_profiles(self):
        for profile in ("rapid", "reviewed", "strict"):
            with self.subTest(profile=profile):
                candidate = valid_spec()
                candidate["verification_profile"] = profile
                result = MODULE.validate_spec(candidate, stage="prebuild")
                self.assertNotIn(
                    "SPEC_VERIFICATION_PROFILE_INVALID",
                    {item["code"] for item in result["errors"]},
                )

    def test_prebuild_requires_coordinate_overlay_for_every_profile(self):
        for profile in ("rapid", "reviewed", "strict"):
            with self.subTest(profile=profile):
                candidate = valid_spec()
                candidate["verification_profile"] = profile
                del candidate["modules"]["page_layout"]["coordinate_overlay_evidence"]

                result = MODULE.validate_spec(candidate, stage="prebuild")

                self.assertEqual(
                    result["errors"],
                    [
                        {
                            "code": "UNSUPPORTED_CAPABILITY",
                            "path": "modules.page_layout",
                            "detail": "missing fields: coordinate_overlay_evidence",
                        }
                    ],
                )

    def test_prebuild_rejects_stale_or_uninspected_coordinate_overlay(self):
        stale = valid_spec()
        stale["modules"]["page_layout"]["coordinate_overlay_evidence"][
            "manifest_sha256"
        ] = "0" * 64
        uninspected = valid_spec()
        uninspected["modules"]["page_layout"]["coordinate_overlay_evidence"][
            "inspection"
        ] = "pending"

        stale_result = MODULE.validate_spec(stale, stage="prebuild")
        uninspected_result = MODULE.validate_spec(uninspected, stage="prebuild")

        self.assertIn(
            "SPEC_COORDINATE_OVERLAY_EVIDENCE_STALE",
            {item["code"] for item in stale_result["errors"]},
        )
        self.assertIn(
            "SPEC_PREBUILD_VISUAL_INSPECTION_NOT_PASSED",
            {item["code"] for item in uninspected_result["errors"]},
        )

    def test_source_change_invalidates_old_coordinate_overlay(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source_path = root / "source.png"
            Image.new("RGB", (1600, 900), "white").save(source_path)
            identity = image_identity(source_path)
            candidate["content_reference"] = dict(identity)
            candidate["clean_visual_reference"] = dict(identity)
            self._refresh_prebuild_visual_evidence(candidate, root)

            Image.new("RGB", (1600, 900), "black").save(source_path)
            changed_identity = image_identity(source_path)
            candidate["content_reference"] = dict(changed_identity)
            candidate["clean_visual_reference"] = dict(changed_identity)

            result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_COORDINATE_OVERLAY_EVIDENCE_STALE",
            {item["code"] for item in result["errors"]},
        )

    def test_prebuild_rejects_unknown_verification_profile(self):
        candidate = valid_spec()
        candidate["verification_profile"] = "automatic"
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn(
            "SPEC_VERIFICATION_PROFILE_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_prebuild_rejects_delivery_status_from_another_profile(self):
        candidate = valid_spec()
        candidate["delivery_status"] = "reviewed_passed"
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn(
            "SPEC_DELIVERY_STATUS_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_legacy_spec_without_profile_uses_strict_contract(self):
        candidate = valid_spec()
        del candidate["verification_profile"]
        del candidate["delivery_status"]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertTrue(result["valid"], result)
        self.assertEqual("strict", result["verification_profile"])

    def _attach_final_gates(
        self,
        candidate: dict,
        root: Path,
        validator_payload: dict | None,
        *,
        profile: str = "strict",
        page_size_pt: list[float] | None = None,
    ) -> None:
        candidate["verification_profile"] = profile
        candidate["delivery_status"] = {
            "rapid": "rapid_validated",
            "reviewed": "reviewed_passed",
            "strict": "strict_gate_passed",
        }[profile]
        candidate["modules"]["typography"]["items"][0]["font_declaration_verified"] = True
        presentation = Presentation()
        presentation.slide_width = 12_192_000
        presentation.slide_height = 6_858_000
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height
        )
        background.name = "ia:background-base"
        item = candidate["modules"]["typography"]["items"][0]
        box = item["text_box"]
        text_box = slide.shapes.add_textbox(box["x"], box["y"], box["w"], box["h"])
        text_box.name = "ia:title"
        paragraph = text_box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = item["text"]
        run.font.bold = True
        run.font.size = Pt(24)
        pptx_path = root / "page.pptx"
        presentation.save(pptx_path)
        pptx = image_identity(pptx_path)

        preview_path = root / "preview.png"
        Image.new("RGB", (1920, 1080), "white").save(preview_path)
        preview = image_identity(preview_path)
        pdf_path = root / "page.pdf"
        pdf_path.write_bytes(b"%PDF-1.4\n% render fixture\n")
        font_report_path = root / "pdffonts.json"
        font_report_path.write_text(
            json.dumps({"resolved_fonts": ["NotoSansCJKsc-Regular"]}),
            encoding="utf-8",
        )
        fontconfig_path = root / "fontconfig.xml"
        fontconfig_path.write_text("<fontconfig/>", encoding="utf-8")
        soffice_path = root / "soffice"
        pdftoppm_path = root / "pdftoppm"
        pdffonts_path = root / "pdffonts"
        for executable in (soffice_path, pdftoppm_path, pdffonts_path):
            executable.write_text("#!/bin/sh\nexit 0\n", encoding="utf-8")
            executable.chmod(0o755)
        runtime_payload = {
            "valid": True,
            "errors": [],
            "renderer_backend": "libreoffice",
            "preview_size": [1920, 1080],
            "executables": {
                "soffice": {
                    "path": str(soffice_path.resolve()),
                    "version": "LibreOffice 26.2.3.2",
                    "sha256": image_identity(soffice_path)["sha256"],
                },
                "pdftoppm": {
                    "path": str(pdftoppm_path.resolve()),
                    "version": "pdftoppm 26.07.0",
                    "sha256": image_identity(pdftoppm_path)["sha256"],
                },
                "pdffonts": {
                    "path": str(pdffonts_path.resolve()),
                    "version": "pdffonts 26.07.0",
                    "sha256": image_identity(pdffonts_path)["sha256"],
                },
            },
            "fontconfig": {
                "path": str(fontconfig_path.resolve()),
                "sha256": image_identity(fontconfig_path)["sha256"],
            },
        }
        runtime_path = root / "preflight-runtime.json"
        runtime_path.write_text(json.dumps(runtime_payload), encoding="utf-8")
        candidate["runtime_preflight"] = image_identity(runtime_path)
        render_payload = {
            "schema_version": 1,
            "pptx": pptx,
            "renderer": {
                "backend": "libreoffice",
                "path": str(soffice_path.resolve()),
                "version": runtime_payload["executables"]["soffice"]["version"],
                "executable_sha256": runtime_payload["executables"]["soffice"]["sha256"],
                "fontconfig_path": str(fontconfig_path.resolve()),
                "fontconfig_sha256": runtime_payload["fontconfig"]["sha256"],
                "isolated_profile": True,
            },
            "pdf": {
                **image_identity(pdf_path),
                "pages": 1,
                "page_size_pt": page_size_pt or [960, 540],
            },
            "font_report": {
                **image_identity(font_report_path),
                "resolved_fonts": ["NotoSansCJKsc-Regular"],
            },
            "rasterizer": {
                "path": str(pdftoppm_path.resolve()),
                "version": runtime_payload["executables"]["pdftoppm"]["version"],
                "executable_sha256": runtime_payload["executables"]["pdftoppm"]["sha256"],
                "output_size": [1920, 1080],
            },
            "preview": {**preview, "size": [1920, 1080]},
        }
        render_report_path = root / "render-report.json"
        render_report_path.write_text(json.dumps(render_payload), encoding="utf-8")
        visual_report = VISUAL_DIFF.build_visual_diff_from_render_report(
            Path(candidate["clean_visual_reference"]["path"]),
            render_report_path,
            root / "visual-diff",
            regions=candidate["regions"],
            profile=profile,
        )
        report_path = Path(visual_report["report"])
        report = image_identity(report_path)
        actual_validator = PPTX_VALIDATOR.validate_pptx(
            pptx_path,
            expected_slides=1,
            reconstruction_spec=candidate,
        )
        if validator_payload is not None:
            actual_validator.update(dict(validator_payload))
        validator_payload = actual_validator
        validator_payload.setdefault("pptx_sha256", pptx["sha256"])
        validator = self._artifact(
            root / "validator.json",
            json.dumps(validator_payload, ensure_ascii=False).encode("utf-8"),
        )
        coverage = {
            "canvas_and_regions": "checked",
            "objects_and_geometry": "checked",
            "text_and_typography": "not_applicable",
            "tables_and_matrices": "not_applicable",
            "graphics_connectors_charts": "not_applicable",
            "pictures_crop_layers": "not_applicable",
            "high_risk_regions": "not_applicable",
        }
        kinds = {
            element.get("kind")
            for element in candidate.get("elements", [])
            if isinstance(element, dict)
        }
        activated = set(candidate.get("activated_modules", []))
        if kinds & {"text", "special_text"} or activated & {"typography", "special_text"}:
            coverage["text_and_typography"] = "checked"
        if kinds & {"table", "matrix"}:
            coverage["tables_and_matrices"] = "checked"
        if kinds & {"shape", "line", "status", "diagram", "chart"} or activated & {"graphics", "diagram", "chart"}:
            coverage["graphics_connectors_charts"] = "checked"
        if kinds & {"icon", "picture"} or activated & {"icons", "picture_framing"}:
            coverage["pictures_crop_layers"] = "checked"
        high_risk = candidate.get("modules", {}).get("high_risk")
        if "high_risk" in activated and isinstance(high_risk, dict) and high_risk.get("items"):
            coverage["high_risk_regions"] = "checked"
        candidate["visual_gate"] = {
            "status": "passed",
            "review_round": 1,
            "evidence": [visual_report["evidence"]["overlay"]["path"]],
            "tripwire": {
                "available": False,
                "triggered": None,
                "reason": "no_approved_baseline",
            },
            "pptx": pptx,
            "preview": preview,
            "report": report,
            "render_report": image_identity(render_report_path),
            "reviewer": {
                "mode": "independent_read_only_subagent",
                "page_id": candidate["page_id"],
                "decision": "passed",
                "source_sha256": candidate["clean_visual_reference"]["sha256"],
                "preview_sha256": preview["sha256"],
                "coverage": coverage,
                "findings": [],
                "p2_disclosures": [],
            },
            "review": {
                "whole_page": "passed",
                "title": "passed",
                "body": "passed",
                "footer": "passed",
                "high_risk_regions": [],
            },
        }
        candidate["editability_gate"] = {
            "status": "passed",
            "evidence": [validator["path"]],
            "pptx": pptx,
            "validator": validator,
            "review": {
                "text_and_data": "passed",
                "native_text_structure": "passed",
                "basic_structure": "passed",
                "full_slide_picture_risk": "passed",
            },
        }
        if profile == "rapid":
            candidate["visual_gate"]["status"] = "not_independently_reviewed"
            candidate["visual_gate"].pop("review_round")
            candidate["visual_gate"].pop("reviewer")
            candidate["visual_gate"].pop("review")

    def _valid_final_spec(
        self,
        profile: str,
        *,
        page_size_pt: list[float] | None = None,
        visual_regions: list[dict] | None = None,
        review_round: int = 1,
    ) -> dict:
        """Build a final fixture from production reports, never claimed booleans."""
        from tests.test_review_admission import (
            AdmissionFixture,
            _load_script,
            _write_json,
        )
        from lib import review_contracts

        candidate = valid_spec()
        rendered_text = "测试标题"
        candidate["elements"][0]["content"]["text"] = rendered_text
        typography_item = candidate["modules"]["typography"]["items"][0]
        typography_item["text"] = rendered_text
        typography_item["runs"][0]["end"] = len(rendered_text)
        typography_item["paragraphs"][0]["end"] = len(rendered_text)
        temporary = tempfile.TemporaryDirectory()
        self.addCleanup(temporary.cleanup)
        root = Path(temporary.name)
        (root / "legacy-final").mkdir()
        self._attach_final_gates(
            candidate,
            root / "legacy-final",
            None,
            profile=profile,
        )
        source_spec = copy.deepcopy(candidate)

        class FinalEvidenceFixture(AdmissionFixture):
            SPEC_FACTORY = staticmethod(lambda _root: copy.deepcopy(source_spec))

            def setUp(inner_self) -> None:
                inner_self._temporary = tempfile.TemporaryDirectory()
                inner_self.root = Path(inner_self._temporary.name)
                inner_self.output = inner_self.root / "admission"
                inner_self.invocations = inner_self.root / "invocations"
                inner_self._generation = 0
                inner_self._rejection_attempt = 0
                inner_self._response_attempt = 0
                inner_self._round_two_attempt = 0
                inner_self._make_fixture(profile=profile)

        fixture = FinalEvidenceFixture()
        fixture.setUp()
        self.addCleanup(fixture.tearDown)

        if page_size_pt is not None:
            fixture.render_report["pdf"]["page_size_pt"] = list(page_size_pt)
            _write_json(fixture.render_path, fixture.render_report)
            geometry = _load_script(
                "create_rendered_text_geometry.py",
                f"task9_final_text_geometry_{id(fixture)}",
            )
            fixture.text_path = fixture.fixture / "task9-rendered-text-geometry.json"
            fixture.text_report = geometry.create_rendered_text_geometry(
                fixture.spec_path,
                fixture.pptx,
                fixture.build_path,
                fixture.render_path,
                fixture.runtime_path,
                fixture.text_path,
            )
            self.assertTrue(fixture.text_report["valid"], fixture.text_report)

        if visual_regions is not None or page_size_pt is not None:
            visual = _load_script(
                "create_visual_diff.py",
                f"task9_final_visual_diff_{id(fixture)}",
            )
            fixture.visual_dir = fixture.fixture / "task9-visual"
            fixture.visual_report = visual.build_visual_diff_from_render_report(
                fixture.source,
                fixture.render_path,
                fixture.visual_dir,
                regions=(
                    visual_regions
                    if visual_regions is not None
                    else fixture.spec["regions"]
                ),
                profile=profile,
            )
            fixture.visual_path = fixture.visual_dir / "visual-diff.json"
            fixture.overlay = Path(
                fixture.visual_report["evidence"]["overlay"]["path"]
            )

        prior = None
        if profile != "rapid" and review_round == 2:
            self.assertEqual("strict", profile)
            prior = fixture.validated_changes_required_response()
            fixture.add_high_risk_mapping(prior)
            # Task 7's round-two test helper carries one compatibility-only
            # nested identity that the production geometry producer never
            # emits.  Final evidence must equal a fresh production report.
            fixture.text_report["inputs"].pop("input_spec_sha256", None)
            _write_json(fixture.text_path, fixture.text_report)
            candidate["activated_modules"].append("high_risk")
            candidate["modules"]["high_risk"] = copy.deepcopy(
                fixture.spec["modules"]["high_risk"]
            )

        visual_gate = candidate["visual_gate"]
        editability_gate = candidate["editability_gate"]
        visual_gate.update(
            {
                "pptx": image_identity(fixture.pptx),
                "preview": image_identity(fixture.preview),
                "report": image_identity(fixture.visual_path),
                "render_report": image_identity(fixture.render_path),
                "background_contract": image_identity(fixture.background_path),
                "rendered_text_geometry": image_identity(fixture.text_path),
                "evidence": [str(fixture.overlay.resolve())],
            }
        )
        candidate["runtime_preflight"] = image_identity(fixture.runtime_path)
        editability_gate.update(
            {
                "pptx": image_identity(fixture.pptx),
                "validator": image_identity(fixture.structure_path),
                "evidence": [str(fixture.structure_path.resolve())],
            }
        )

        if profile == "rapid":
            return candidate

        freezer = _load_script(
            "freeze_reconstruction_spec.py",
            f"task9_freeze_pre_review_{id(fixture)}",
        )
        working_spec_path = root / "page-reconstruction.json"
        _write_json(working_spec_path, candidate)
        pre_review_spec_path = root / "pre-review-spec-snapshot.json"
        snapshot_report = freezer.freeze_spec(
            working_spec_path,
            pre_review_spec_path,
            "pre-review",
        )
        self.assertEqual(
            str(pre_review_spec_path.resolve()),
            snapshot_report["snapshot"]["path"],
        )
        admission_inputs = review_contracts.AdmissionInputs(
            spec=pre_review_spec_path,
            pptx=fixture.pptx,
            build_report=fixture.build_path,
            structure_report=fixture.structure_path,
            render_report=fixture.render_path,
            text_geometry=fixture.text_path,
            background_report=fixture.background_path,
            visual_diff=fixture.visual_path,
            review_round=review_round,
            prior_admission=(prior["admission_path"] if prior is not None else None),
            prior_invocation=(prior["invocation_path"] if prior is not None else None),
            prior_response_validation=(
                prior["validation_path"] if prior is not None else None
            ),
        )
        admission = review_contracts.issue_admission(
            admission_inputs,
            fixture.output,
        )

        assert isinstance(admission, dict)
        invocation = review_contracts.record_invocation(
            fixture.admission_path,
            fixture.invocations,
        )
        invocation_path = fixture.invocations / (
            f"{admission['page_id']}-round-{admission['review_round']}-invocation.json"
        )
        response = fixture.valid_response(admission)
        response["coverage"] = copy.deepcopy(
            visual_gate["reviewer"]["coverage"]
        )
        if review_round == 2:
            response["coverage"]["high_risk_regions"] = "checked"
        response_path = fixture.root / "review-response.json"
        response_path.write_text(
            json.dumps(response, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
            encoding="utf-8",
        )
        response_validation_path = fixture.root / "review-response-validation.json"
        response_validation = review_contracts.validate_response(
            fixture.admission_path,
            invocation_path,
            response_path,
            response_validation_path,
        )
        self.assertTrue(response_validation["valid"], response_validation)
        self.assertEqual(invocation["admission_id"], admission["admission_id"])
        # issue() canonically rewrites every bound JSON payload.  Refresh all
        # identities after admission so the legacy final checks and the new
        # immutable review chain point at the same production bytes.
        visual_gate.update(
            {
                "report": image_identity(fixture.visual_path),
                "render_report": image_identity(fixture.render_path),
                "background_contract": image_identity(fixture.background_path),
                "rendered_text_geometry": image_identity(fixture.text_path),
                "review_round": admission["review_round"],
                "review_admission": image_identity(fixture.admission_path),
                "review_invocation": image_identity(invocation_path),
                "review_response_validation": image_identity(
                    response_validation_path
                ),
            }
        )
        candidate["runtime_preflight"] = image_identity(fixture.runtime_path)
        editability_gate["validator"] = image_identity(fixture.structure_path)
        visual_gate["reviewer"] = {
            "mode": "independent_read_only_subagent",
            **response,
        }
        return candidate

    def test_final_accepts_production_pre_review_evidence_for_each_profile(self) -> None:
        for profile in ("rapid", "reviewed", "strict"):
            with self.subTest(profile=profile):
                spec = self._valid_final_spec(profile)
                temporary = tempfile.TemporaryDirectory()
                self.addCleanup(temporary.cleanup)
                spec_path = Path(temporary.name) / "page-reconstruction.json"
                report_path = Path(temporary.name) / "final-validation.json"
                spec_path.write_text(
                    json.dumps(spec, ensure_ascii=False, indent=2) + "\n",
                    encoding="utf-8",
                )
                with contextlib.redirect_stdout(io.StringIO()):
                    exit_code = MODULE.main(
                        [
                            str(spec_path),
                            "--stage",
                            "final",
                            "--output",
                            str(report_path),
                        ]
                    )
                result = json.loads(report_path.read_text(encoding="utf-8"))
                self.assertEqual(0, exit_code, result)
                self.assertTrue(result["valid"], result)

    def test_final_rejects_pre_review_state_drift_for_each_gate_family(self) -> None:
        def add_high_risk(spec: dict) -> None:
            spec["activated_modules"].append("high_risk")
            spec["modules"]["high_risk"] = {
                "items": [
                    {
                        "risk_id": "post-review-drift",
                        "source": "manual:post-review-drift",
                        "scope": "whole_page",
                        "category": "objects_and_geometry",
                        "expected": "unchanged",
                        "strategy": "verify",
                        "result": "passed",
                        "evidence": [spec["visual_gate"]["evidence"][0]],
                        "confidence": "high",
                        "severity": "P1",
                        "verification": {},
                    }
                ]
            }

        mutations = {
            "tripwire": lambda spec: spec["visual_gate"].__setitem__(
                "tripwire",
                {"available": True, "triggered": False},
            ),
            "editability": lambda spec: spec["editability_gate"]["review"].__setitem__(
                "basic_structure", "changes_required"
            ),
            "high_risk": add_high_risk,
            "visual_evidence": lambda spec: spec["visual_gate"]["evidence"].append(
                spec["visual_gate"]["evidence"][0]
            ),
        }
        for label, mutate in mutations.items():
            with self.subTest(label=label):
                spec = self._valid_final_spec("reviewed")
                mutate(spec)
                result = MODULE.validate_spec(spec, stage="final")
                self.assertIn(
                    "REVIEW_ADMISSION_STALE",
                    {item["code"] for item in result["errors"]},
                    result,
                )

    def test_final_rejects_replaced_build_snapshot(self) -> None:
        spec = self._valid_final_spec("rapid")
        text_path = Path(spec["visual_gate"]["rendered_text_geometry"]["path"])
        text_report = json.loads(text_path.read_text(encoding="utf-8"))
        build_snapshot = Path(text_report["input_paths"]["spec"])
        build_snapshot.write_text(
            build_snapshot.read_text(encoding="utf-8") + " ",
            encoding="utf-8",
        )

        result = MODULE.validate_spec(spec, stage="final")

        self.assertTrue(
            {"TEXT_GEOMETRY_IDENTITY_MISMATCH", "BACKGROUND_ASSET_INVALID"}
            & {item["code"] for item in result["errors"]},
            result,
        )

    def test_final_rejects_pre_review_snapshot_path_identity_mismatch(self) -> None:
        spec = self._valid_final_spec("reviewed")
        admission_path = Path(spec["visual_gate"]["review_admission"]["path"])
        admission = json.loads(admission_path.read_text(encoding="utf-8"))
        original = Path(admission["artifacts"]["spec"]["path"])
        copied = original.with_name("copied-pre-review-spec.json")
        copied.write_bytes(original.read_bytes())
        admission["artifacts"]["spec"]["path"] = str(copied.resolve())
        admission_path.write_text(
            json.dumps(admission, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
            encoding="utf-8",
        )
        spec["visual_gate"]["review_admission"] = image_identity(admission_path)

        result = MODULE.validate_spec(spec, stage="final")

        self.assertIn(
            "REVIEW_ADMISSION_STALE",
            {item["code"] for item in result["errors"]},
            result,
        )

    def test_final_rejects_missing_rendered_text_geometry(self) -> None:
        spec = self._valid_final_spec(profile="reviewed")
        del spec["visual_gate"]["rendered_text_geometry"]

        result = MODULE.validate_spec(spec, stage="final")

        self.assertIn(
            "TEXT_GEOMETRY_MISSING", {item["code"] for item in result["errors"]}
        )

    def test_final_rejects_reviewer_not_bound_to_validated_admission(self) -> None:
        spec = self._valid_final_spec(profile="reviewed")
        spec["visual_gate"]["reviewer"]["admission_id"] = "f" * 64

        result = MODULE.validate_spec(spec, stage="final")

        self.assertIn(
            "REVIEW_RESPONSE_INVALID", {item["code"] for item in result["errors"]}
        )

    def test_final_rejects_visual_gate_round_not_bound_to_admission(self) -> None:
        spec = self._valid_final_spec(profile="reviewed")
        spec["visual_gate"]["review_round"] = 2

        result = MODULE.validate_spec(spec, stage="final")

        self.assertIn(
            "REVIEW_RESPONSE_INVALID", {item["code"] for item in result["errors"]}
        )

    def test_final_allows_delivery_fields_but_not_renderable_changes(self) -> None:
        spec = self._valid_final_spec(profile="reviewed")
        self.assertTrue(MODULE.validate_spec(spec, stage="final")["valid"])
        spec["elements"][0]["slide_bbox"][2] += 12700

        result = MODULE.validate_spec(spec, stage="final")

        self.assertIn(
            "TEXT_GEOMETRY_IDENTITY_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rapid_requires_pre_review_gates_and_forbids_reviewer_artifacts(self) -> None:
        for missing, code in (
            ("background_contract", "BACKGROUND_DECLARATION_MISSING"),
            ("rendered_text_geometry", "TEXT_GEOMETRY_MISSING"),
        ):
            with self.subTest(missing=missing):
                spec = self._valid_final_spec(profile="rapid")
                del spec["visual_gate"][missing]
                result = MODULE.validate_spec(spec, stage="final")
                self.assertIn(code, {item["code"] for item in result["errors"]})

        reviewed = self._valid_final_spec(profile="reviewed")
        for field in (
            "review_admission",
            "review_invocation",
            "review_response_validation",
        ):
            with self.subTest(forbidden=field):
                spec = self._valid_final_spec(profile="rapid")
                spec["visual_gate"][field] = copy.deepcopy(
                    reviewed["visual_gate"][field]
                )
                result = MODULE.validate_spec(spec, stage="final")
                self.assertIn(
                    "SPEC_RAPID_REVIEWER_FORBIDDEN",
                    {item["code"] for item in result["errors"]},
                )

    def test_final_reviewed_and_strict_require_complete_review_artifact_chain(self) -> None:
        for profile in ("reviewed", "strict"):
            for missing, code in (
                ("review_admission", "REVIEW_ADMISSION_NOT_ISSUED"),
                ("review_invocation", "REVIEW_RESPONSE_INVALID"),
                ("review_response_validation", "REVIEW_RESPONSE_INVALID"),
            ):
                with self.subTest(profile=profile, missing=missing):
                    spec = self._valid_final_spec(profile=profile)
                    del spec["visual_gate"][missing]
                    result = MODULE.validate_spec(spec, stage="final")
                    self.assertIn(code, {item["code"] for item in result["errors"]})

    def test_final_accepts_rapid_without_reviewer(self):
        candidate = self._valid_final_spec(profile="rapid")
        result = MODULE.validate_spec(candidate, stage="final")
        self.assertTrue(result["valid"], result)

    def test_final_accepts_libreoffice_hundredth_mm_page_size(self):
        candidate = self._valid_final_spec(
            profile="rapid",
            page_size_pt=[960.009448818898, 540.0],
        )
        result = MODULE.validate_spec(candidate, stage="final")
        self.assertTrue(result["valid"], result)

    def test_final_rejects_page_size_outside_render_tolerance(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
                page_size_pt=[961.01, 540.0],
            )
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RENDER_REPORT_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_missing_render_report(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            del candidate["visual_gate"]["render_report"]
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RENDER_REPORT_MISSING",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_render_report_pptx_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            report_path = Path(candidate["visual_gate"]["render_report"]["path"])
            payload = json.loads(report_path.read_text(encoding="utf-8"))
            payload["pptx"]["sha256"] = "f" * 64
            report_path.write_text(json.dumps(payload), encoding="utf-8")
            candidate["visual_gate"]["render_report"] = image_identity(report_path)
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RENDER_PPTX_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_render_report_preview_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            report_path = Path(candidate["visual_gate"]["render_report"]["path"])
            payload = json.loads(report_path.read_text(encoding="utf-8"))
            payload["preview"]["sha256"] = "f" * 64
            report_path.write_text(json.dumps(payload), encoding="utf-8")
            candidate["visual_gate"]["render_report"] = image_identity(report_path)
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RENDER_PREVIEW_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_render_report_runtime_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            report_path = Path(candidate["visual_gate"]["render_report"]["path"])
            payload = json.loads(report_path.read_text(encoding="utf-8"))
            payload["renderer"]["version"] = "LibreOffice 99.0.0"
            report_path.write_text(json.dumps(payload), encoding="utf-8")
            candidate["visual_gate"]["render_report"] = image_identity(report_path)
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RENDER_RUNTIME_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_visual_diff_render_report_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            visual_path = Path(candidate["visual_gate"]["report"]["path"])
            payload = json.loads(visual_path.read_text(encoding="utf-8"))
            payload["render_report"]["sha256"] = "f" * 64
            visual_path.write_text(json.dumps(payload), encoding="utf-8")
            candidate["visual_gate"]["report"] = image_identity(visual_path)
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RENDER_REPORT_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_font_trace_pdf_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            candidate["modules"]["typography"]["items"][0]["fallback_trace"] = {
                "requested_font": "Noto Sans CJK SC",
                "resolved_fonts": ["NotoSansCJKsc-Regular"],
                "pdf_sha256": "f" * 64,
            }
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_FONT_TRACE_RENDER_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rapid_rejects_independent_passed_status(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            candidate["visual_gate"]["status"] = "passed"
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RAPID_VISUAL_STATUS_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rapid_rejects_delivery_status_from_other_profile(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            candidate["delivery_status"] = "reviewed_passed"
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_DELIVERY_STATUS_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rapid_rejects_reviewer_claim(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            candidate["visual_gate"]["reviewer"] = {"mode": "independent_read_only_subagent"}
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_RAPID_REVIEWER_FORBIDDEN",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rapid_rejects_triggered_tripwire(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            candidate["visual_gate"]["tripwire"] = {
                "available": True,
                "triggered": True,
                "reason": "below_minimum_similarity",
            }
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_VISUAL_TRIPWIRE_TRIGGERED",
            {item["code"] for item in result["errors"]},
        )

    def test_final_reviewed_requires_independent_reviewer(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="reviewed",
            )
            del candidate["visual_gate"]["reviewer"]
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_INDEPENDENT_VISUAL_REVIEW_REQUIRED",
            {item["code"] for item in result["errors"]},
        )

    def test_final_accepts_reviewed_with_current_reviewer(self):
        candidate = self._valid_final_spec(profile="reviewed")
        result = MODULE.validate_spec(candidate, stage="final")
        self.assertTrue(result["valid"], result)

    def test_final_reviewed_accepts_partial_region_evidence(self):
        candidate = self._valid_final_spec(
            profile="reviewed",
            visual_regions=[],
        )
        result = MODULE.validate_spec(candidate, stage="final")
        self.assertTrue(result["valid"], result)

    def test_final_rejects_visual_diff_from_another_profile(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                profile="rapid",
            )
            report_path = Path(candidate["visual_gate"]["report"]["path"])
            report = json.loads(report_path.read_text(encoding="utf-8"))
            report["verification_profile"] = "reviewed"
            report_path.write_text(json.dumps(report), encoding="utf-8")
            candidate["visual_gate"]["report"] = image_identity(report_path)
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn(
            "SPEC_VISUAL_DIFF_PROFILE_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def _add_valid_icon_contract(self, candidate: dict, root: Path) -> None:
        source_path = root / "source.png"
        source_image = Image.new("RGB", (1600, 900), (0, 0, 0))
        source_image.paste((241, 90, 34), (104, 104, 128, 128))
        source_image.save(source_path)
        source = {
            "path": str(source_path.resolve()),
            "sha256": hashlib.sha256(source_path.read_bytes()).hexdigest(),
        }
        icons_dir = root / "assets" / "icons"
        icons_dir.mkdir(parents=True)
        asset_path = icons_dir / "status.png"
        icon = Image.new("RGBA", (32, 32), (0, 0, 0, 0))
        icon.paste((241, 90, 34, 255), (4, 4, 28, 28))
        icon.save(asset_path)
        asset = {
            "path": str(asset_path.resolve()),
            "sha256": hashlib.sha256(asset_path.read_bytes()).hexdigest(),
        }
        alpha_mask_sha256 = hashlib.sha256(icon.getchannel("A").tobytes()).hexdigest()
        candidate["content_reference"] = source
        candidate["clean_visual_reference"] = source
        candidate["activated_modules"].append("icons")
        candidate["elements"].append(
            {
                "element_id": "status-icon",
                "kind": "icon",
                "source_bbox": [100, 100, 32, 32],
                "slide_bbox": [762000, 762000, 243840, 243840],
                "layer": 3,
                "editable": False,
                "confidence": "high",
                "style": {"rotation": 0, "opacity": 1},
                "content": {
                    "asset": {
                        "path": asset["path"],
                        "asset_sha256": asset["sha256"],
                        "pixel_size": [32, 32],
                    },
                    "mode": "none",
                    "crop": {"left": 0, "top": 0, "right": 0, "bottom": 0},
                },
            }
        )
        candidate["regions"][0]["element_ids"].append("status-icon")
        candidate["reading_order"].append("status-icon")
        candidate["modules"]["representation_plan"]["items"].append(
            {
                "source_fact_id": "fact-status-icon",
                "semantic_role": "icon",
                "source_bbox": [100, 100, 32, 32],
                "required": True,
                "selected_mode": "asset",
                "required_editability": "none",
                "fallback_policy": "allow_minimal_asset",
                "bound_element_ids": ["status-icon"],
                "reason": "icon is represented by its alpha-isolated local asset",
                "coverage_status": "covered",
                "evidence": [str(asset_path.resolve())],
            }
        )
        candidate["modules"]["icons"] = {
            "schema_version": 2,
            "page_id": "page-001",
            "slide_coordinate_unit": "EMU",
            "clean_visual_reference": source["path"],
            "clean_visual_sha256": source["sha256"],
            "icons": [
                {
                    "icon_id": "status-icon",
                    "element_id": "status-icon",
                    "category": "simple_symbol",
                    "instance_count": 1,
                    "repeat_group": None,
                    "semantic_scope": "icon_only",
                    "source_bbox": [100, 100, 32, 32],
                    "slide_bbox": [762000, 762000, 243840, 243840],
                    "layer": 3,
                    "source_path": source["path"],
                    "source_sha256": source["sha256"],
                    "crop_mode": "alpha_isolation",
                    "padding": 0,
                    "background_handling": "border_connected_background_to_alpha",
                    "asset_path": asset["path"],
                    "asset_sha256": asset["sha256"],
                    "alpha_mask_sha256": alpha_mask_sha256,
                    "final_width": 32,
                    "final_height": 32,
                    "sharpness": "source_preserved",
                    "validation": "passed",
                    "native_redraw": False,
                    "selectable_picture_verified": False,
                    "object_type": "picture",
                }
            ],
        }
        self._refresh_prebuild_visual_evidence(candidate, root)

    def _refresh_prebuild_visual_evidence(self, candidate: dict, root: Path) -> None:
        work_dir = root / "work"
        work_dir.mkdir(parents=True, exist_ok=True)
        coordinate_path = work_dir / "coordinate-overlay.png"
        coordinate = COORDINATE_OVERLAY.create_coordinate_overlay(
            Path(candidate["clean_visual_reference"]["path"]),
            coordinate_path,
        )
        candidate["modules"]["page_layout"]["coordinate_overlay_evidence"] = {
            "path": str(coordinate_path.resolve()),
            "sha256": hashlib.sha256(coordinate_path.read_bytes()).hexdigest(),
            "source_sha256": coordinate["source"]["sha256"],
            "manifest_sha256": coordinate["coordinate_overlay_manifest_sha256"],
            "grid": coordinate["grid"],
            "inspection": "passed",
        }
    def _replace_icon_asset(self, candidate: dict, image: Image.Image) -> None:
        icon = candidate["modules"]["icons"]["icons"][0]
        asset_path = Path(icon["asset_path"])
        image.save(asset_path)
        icon["asset_sha256"] = hashlib.sha256(asset_path.read_bytes()).hexdigest()
        icon["alpha_mask_sha256"] = (
            hashlib.sha256(image.getchannel("A").tobytes()).hexdigest()
            if "A" in image.getbands()
            else None
        )

    def test_valid_prebuild_spec_passes(self):
        with tempfile.TemporaryDirectory() as directory:
            result = MODULE.validate_spec(
                make_minimal_spec(Path(directory)), stage="prebuild"
            )
        self.assertTrue(result["valid"], result)
        self.assertEqual([], result["errors"])

    def test_prebuild_rejects_native_shape_fill_string_at_renderer_gate(self):
        candidate = valid_spec()
        append_native_element(
            candidate,
            element_id="native-shape",
            kind="shape",
            style={
                "shape_type": "roundRect",
                "adjustments": [0.2],
                "fill": "#112233",
            },
        )

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertFalse(result["valid"])
        self.assertIn(
            {
                "code": "UNSUPPORTED_CAPABILITY",
                "path": "elements.native-shape.style.fill",
                "detail": "contract must be an object",
            },
            result["errors"],
        )

    def test_prebuild_rejects_line_stroke_string_at_renderer_gate(self):
        candidate = valid_spec()
        append_native_element(
            candidate,
            element_id="native-line",
            kind="line",
            style={"line": "#112233"},
        )

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertFalse(result["valid"])
        self.assertIn(
            {
                "code": "UNSUPPORTED_CAPABILITY",
                "path": "elements.native-line.style.line",
                "detail": "contract must be an object",
            },
            result["errors"],
        )

    def test_prebuild_rejects_text_renderer_only_fill_payload(self):
        candidate = valid_spec()
        candidate["elements"][0]["style"]["fill"] = "#112233"

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertFalse(result["valid"])
        self.assertIn(
            {
                "code": "UNSUPPORTED_CAPABILITY",
                "path": "elements.title.style.fill",
                "detail": "only noFill is supported",
            },
            result["errors"],
        )

    def test_prebuild_rejects_picture_renderer_only_crop_payload(self):
        with tempfile.TemporaryDirectory() as directory:
            candidate = make_asset_fallback_spec(Path(directory))
            candidate["elements"][-1]["content"]["crop"] = "full"

            result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertFalse(result["valid"])
        self.assertIn(
            {
                "code": "UNSUPPORTED_CAPABILITY",
                "path": "elements.artwork-picture.content.crop",
                "detail": "crop must be an object",
            },
            result["errors"],
        )

    def test_prebuild_requires_explicit_modern_text_run_fields(self):
        required = {
            "italic": False,
            "underline": False,
            "strike": False,
            "baseline": 0,
        }
        for missing in required:
            with self.subTest(missing=missing):
                candidate = valid_spec()
                run = candidate["modules"]["typography"]["items"][0]["runs"][0]
                run.update(required)
                del run[missing]

                result = MODULE.validate_spec(candidate, stage="prebuild")

                self.assertIn(
                    "SPEC_TEXT_RUN_STYLE_INVALID",
                    {item["code"] for item in result["errors"]},
                )

    def test_prebuild_rejects_invalid_text_run_flag_and_baseline_values(self):
        cases = (
            ("italic", 1),
            ("underline", "yes"),
            ("strike", None),
            ("baseline", 1.5),
            ("baseline", -100001),
            ("baseline", 100001),
        )
        for field, value in cases:
            with self.subTest(field=field, value=value):
                candidate = valid_spec()
                run = candidate["modules"]["typography"]["items"][0]["runs"][0]
                run.update(
                    {
                        "italic": False,
                        "underline": False,
                        "strike": False,
                        "baseline": 0,
                    }
                )
                run[field] = value

                result = MODULE.validate_spec(candidate, stage="prebuild")

                self.assertIn(
                    "SPEC_TEXT_RUN_STYLE_INVALID",
                    {item["code"] for item in result["errors"]},
                )

    def test_prebuild_rejects_unknown_fields_at_every_typography_contract_level(self):
        cases = (
            ("item", "modules.typography.items[0].future_effect", lambda item: item),
            (
                "text_box",
                "modules.typography.items[0].text_box.future_effect",
                lambda item: item["text_box"],
            ),
            (
                "paragraph",
                "modules.typography.items[0].paragraphs[0].future_effect",
                lambda item: item["paragraphs"][0],
            ),
            (
                "run",
                "modules.typography.items[0].runs[0].future_effect",
                lambda item: item["runs"][0],
            ),
            (
                "list",
                "modules.typography.items[0].paragraphs[0].list.future_effect",
                lambda item: item["paragraphs"][0]["list"],
            ),
        )
        for case, expected_path, select_contract in cases:
            with self.subTest(case=case):
                candidate = valid_spec()
                typography_item = candidate["modules"]["typography"]["items"][0]
                select_contract(typography_item)["future_effect"] = "glow"

                result = MODULE.validate_spec(candidate, stage="prebuild")

                self.assertIn(
                    {
                        "code": "SPEC_TYPOGRAPHY_FIELD_UNKNOWN",
                        "path": expected_path,
                        "detail": "unknown field: future_effect",
                    },
                    result["errors"],
                )

    def test_prebuild_rejects_legacy_decoration_on_modern_run(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["items"][0]["runs"][0][
            "decoration"
        ] = "none"

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            {
                "code": "SPEC_TYPOGRAPHY_FIELD_UNKNOWN",
                "path": "modules.typography.items[0].runs[0].decoration",
                "detail": "unknown field: decoration",
            },
            result["errors"],
        )

    def test_final_keeps_legacy_decoration_only_runs_compatible(self):
        candidate = valid_spec()
        run = candidate["modules"]["typography"]["items"][0]["runs"][0]
        for field in ("italic", "underline", "strike", "baseline"):
            del run[field]
        run["decoration"] = "none"

        result = MODULE.validate_spec(candidate, stage="final")

        self.assertNotIn(
            "SPEC_TEXT_RUN_STYLE_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_prebuild_requires_representation_plan_activation(self):
        candidate = valid_spec()
        candidate["activated_modules"].remove("representation_plan")
        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_ACTIVATED_MODULES_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_current_schema_requires_background_module_activation(self):
        missing_activation = valid_spec()
        missing_activation["activated_modules"].remove("background")
        missing_module = valid_spec()
        missing_module["modules"].pop("background")

        activation_result = MODULE.validate_spec(
            missing_activation, stage="authoring"
        )
        module_result = MODULE.validate_spec(missing_module, stage="prebuild")

        self.assertIn(
            "SPEC_ACTIVATED_MODULES_INVALID",
            {item["code"] for item in activation_result["errors"]},
        )
        self.assertIn(
            "UNSUPPORTED_CAPABILITY",
            {item["code"] for item in module_result["errors"]},
        )

    def test_icon_renderer_still_rejects_background_picture_mode(self) -> None:
        from tests.test_build_pptx_from_spec import make_icon_spec

        with tempfile.TemporaryDirectory() as directory:
            spec = make_icon_spec(Path(directory))
            spec["modules"]["background"]["items"][0][
                "bound_element_id"
            ] = spec["modules"]["icons"]["icons"][0]["element_id"]
            spec["modules"]["background"]["items"][0][
                "selected_mode"
            ] = "background_picture"

            result = MODULE.validate_spec(spec, stage="prebuild")

        self.assertIn(
            "BACKGROUND_ICON_BINDING_FORBIDDEN",
            {item["code"] for item in result["errors"]},
        )

    def test_final_keeps_legacy_spec_compatible_without_representation_plan(self):
        candidate = valid_spec()
        candidate["activated_modules"].remove("representation_plan")
        candidate["modules"].pop("representation_plan")

        result = MODULE.validate_spec(candidate, stage="final")

        codes = {item["code"] for item in result["errors"]}
        self.assertNotIn("SPEC_ACTIVATED_MODULES_INVALID", codes)
        self.assertNotIn("REPRESENTATION_INCOMPLETE", codes)

    def test_prebuild_rejects_empty_representation_plan(self):
        candidate = valid_spec()
        candidate["modules"]["representation_plan"] = {"items": []}

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "REPRESENTATION_INCOMPLETE",
            {item["code"] for item in result["errors"]},
        )

    def test_prebuild_report_hash_changes_with_representation_field(self):
        with tempfile.TemporaryDirectory() as directory:
            candidate = make_minimal_spec(Path(directory))
            changed = copy.deepcopy(candidate)
            changed["modules"]["representation_plan"]["items"][0]["reason"] = "changed evidence rationale"

            first = MODULE.validate_spec(candidate, stage="prebuild")
            second = MODULE.validate_spec(changed, stage="prebuild")

        self.assertTrue(first["valid"], first)
        self.assertTrue(second["valid"], second)
        self.assertNotEqual(first["spec_sha256"], second["spec_sha256"])

    def test_prebuild_rejects_missing_or_changed_reference(self):
        missing = valid_spec()
        missing["content_reference"] = {
            "path": "/tmp/ia-reference-does-not-exist.png",
            "sha256": "0" * 64,
        }
        changed = valid_spec()
        changed["clean_visual_reference"]["sha256"] = "f" * 64

        missing_result = MODULE.validate_spec(missing, stage="prebuild")
        changed_result = MODULE.validate_spec(changed, stage="prebuild")

        self.assertIn(
            "SPEC_REFERENCE_NOT_FOUND",
            {item["code"] for item in missing_result["errors"]},
        )
        self.assertIn(
            "SPEC_REFERENCE_HASH_MISMATCH",
            {item["code"] for item in changed_result["errors"]},
        )

    def test_activated_module_must_be_nonempty(self):
        candidate = valid_spec()
        candidate["activated_modules"].append("graphics")
        candidate["modules"]["graphics"] = {}

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_ACTIVATED_MODULE_EMPTY",
            {item["code"] for item in result["errors"]},
        )

    def test_module_element_references_must_exist(self):
        candidate = valid_spec()
        candidate["activated_modules"].append("graphics")
        candidate["modules"]["graphics"] = {
            "items": [{"element_id": "missing-element"}]
        }

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_MODULE_ELEMENT_REFERENCE_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_reading_order_must_cover_every_element(self):
        candidate = valid_spec()
        extra = dict(candidate["elements"][0])
        extra["element_id"] = "subtitle"
        candidate["elements"].append(extra)
        candidate["regions"][0]["element_ids"].append("subtitle")

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_READING_ORDER_COVERAGE_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_regions_must_cover_every_element(self):
        candidate = valid_spec()
        extra = dict(candidate["elements"][0])
        extra["element_id"] = "subtitle"
        candidate["elements"].append(extra)
        candidate["reading_order"].append("subtitle")

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_REGION_COVERAGE_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_icon_asset_path_does_not_derive_from_clean_reference(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            self._add_valid_icon_contract(candidate, root)
            original = Path(candidate["clean_visual_reference"]["path"])
            reference_dir = root / "references"
            reference_dir.mkdir()
            moved = reference_dir / "source.png"
            original.replace(moved)
            identity = image_identity(moved)
            candidate["content_reference"] = dict(identity)
            candidate["clean_visual_reference"] = dict(identity)
            icons = candidate["modules"]["icons"]
            icons["clean_visual_reference"] = identity["path"]
            icons["clean_visual_sha256"] = identity["sha256"]
            icon = icons["icons"][0]
            icon["source_path"] = identity["path"]
            icon["source_sha256"] = identity["sha256"]
            self._refresh_prebuild_visual_evidence(candidate, root)

            result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertTrue(result["valid"], result)

    def test_icon_asset_real_dimensions_must_match_declaration(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            icon = Image.new("RGBA", (16, 16), (0, 0, 0, 0))
            icon.paste((241, 90, 34, 255), (4, 4, 12, 12))
            self._replace_icon_asset(candidate, icon)

            result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_ICON_ASSET_DIMENSIONS_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_low_risk_typography_can_omit_manual_trials_and_metrics(self):
        candidate = valid_spec()

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertTrue(result["valid"], result)

    def test_typography_does_not_require_candidates(self):
        candidate = valid_spec()
        item = candidate["modules"]["typography"]["items"][0]
        self.assertNotIn("candidates", item)

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertTrue(result["valid"], result)

    def test_removed_font_trial_fields_are_rejected(self):
        candidate = valid_spec()
        item = candidate["modules"]["typography"]["items"][0]
        item["candidate_trials"] = [{"font": "Noto Sans CJK SC"}]

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_REMOVED_FONT_WORKFLOW_FIELD",
            {entry["code"] for entry in result["errors"]},
        )

    def test_unknown_source_font_requires_noto_sans(self):
        candidate = valid_spec()
        item = candidate["modules"]["typography"]["items"][0]
        item["source_font_guess"] = "unknown"
        item["selected_font"] = "STKaiti"

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_UNCERTAIN_FONT_FALLBACK_INVALID",
            {entry["code"] for entry in result["errors"]},
        )

    def test_unknown_source_font_requires_fallback_reason(self):
        candidate = valid_spec()
        item = candidate["modules"]["typography"]["items"][0]
        item["source_font_guess"] = "unknown"
        item["selected_font"] = "Noto Sans CJK SC"
        item["fallback_reason"] = None

        result = MODULE.validate_spec(candidate, stage="prebuild")

        self.assertIn(
            "SPEC_UNCERTAIN_FONT_FALLBACK_INVALID",
            {entry["code"] for entry in result["errors"]},
        )

    def test_typography_requires_emu(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["slide_coordinate_unit"] = "px"
        result = MODULE.validate_spec(candidate, stage="prebuild")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_TYPOGRAPHY_UNIT_INVALID", codes)

    def test_slide_bbox_rejects_pixel_scale_values(self):
        candidate = valid_spec()
        candidate["elements"][0]["slide_bbox"] = [24, 20, 1056, 64]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_SLIDE_BBOX_UNIT_SUSPECT", codes)

    def test_slide_bbox_rejects_mixed_pixel_positions(self):
        candidate = valid_spec()
        candidate["elements"][0]["slide_bbox"][:2] = [30, 30]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_SLIDE_BBOX_MAPPING_INVALID", {item["code"] for item in result["errors"]})

    def test_region_rejects_pixel_scale_and_out_of_bounds(self):
        candidate = valid_spec()
        candidate["regions"][0]["slide_bbox"] = [0, 0, 1600, 120]
        candidate["regions"][0]["source_bbox"] = [2000, 0, 100, 100]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_REGION_BBOX_OUT_OF_BOUNDS", codes)
        self.assertIn("SPEC_SLIDE_BBOX_UNIT_SUSPECT", codes)

    def test_typography_text_box_must_match_element_emu_bbox(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["items"][0]["text_box"].update(
            {"x": 30, "y": 30, "w": 800, "h": 60}
        )
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_TEXT_BOX_MAPPING_INVALID", {item["code"] for item in result["errors"]})

    def test_unknown_activated_module_is_rejected(self):
        candidate = valid_spec()
        candidate["activated_modules"] = ["page_layout", "typogrphy"]
        candidate["modules"]["typogrphy"] = {}
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ACTIVATED_MODULE_UNKNOWN", {item["code"] for item in result["errors"]})

    def test_activated_icons_require_complete_contract(self):
        candidate = valid_spec()
        candidate["activated_modules"].append("icons")
        candidate["modules"]["icons"] = {}
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ICONS_FIELD_MISSING", {item["code"] for item in result["errors"]})

    def test_complete_icon_contract_passes_prebuild(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertTrue(result["valid"], result)

    def test_icon_page_does_not_require_crop_review_or_inspection(self):
        for profile in ("rapid", "reviewed", "strict"):
            with self.subTest(profile=profile), tempfile.TemporaryDirectory() as directory:
                candidate = valid_spec()
                candidate["verification_profile"] = profile
                self._add_valid_icon_contract(candidate, Path(directory))

                result = MODULE.validate_spec(candidate, stage="prebuild")

                self.assertTrue(result["valid"], result)

    def test_final_icon_contract_requires_selectability_verification(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            result = MODULE.validate_spec(candidate, stage="final")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_ICON_SELECTABILITY_NOT_VERIFIED", codes)

    def test_non_alpha_icon_crop_modes_are_rejected(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            for crop_mode in ("background_preserved", "tight_rect"):
                with self.subTest(crop_mode=crop_mode):
                    candidate["modules"]["icons"]["icons"][0]["crop_mode"] = crop_mode
                    result = MODULE.validate_spec(candidate, stage="prebuild")
                    self.assertIn(
                        "SPEC_ICON_CROP_MODE_INVALID",
                        {item["code"] for item in result["errors"]},
                    )

    def test_alpha_isolation_rejects_mismatched_alpha_hash(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            candidate["modules"]["icons"]["icons"][0]["alpha_mask_sha256"] = "f" * 64
            result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ICON_ALPHA_MASK_MISMATCH", {item["code"] for item in result["errors"]})

    def test_alpha_isolation_rejects_changed_rgb_even_under_transparent_pixel(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            icon = Image.new("RGBA", (32, 32), (0, 0, 0, 0))
            icon.paste((241, 90, 34, 255), (4, 4, 28, 28))
            icon.putpixel((0, 0), (1, 2, 3, 0))
            self._replace_icon_asset(candidate, icon)
            result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ICON_RGB_MISMATCH", {item["code"] for item in result["errors"]})

    def test_alpha_isolation_rejects_opaque_asset(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            self._replace_icon_asset(
                candidate,
                Image.new("RGBA", (32, 32), (241, 90, 34, 255)),
            )
            result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ICON_ALPHA_CONTENT_INVALID", {item["code"] for item in result["errors"]})

    def test_alpha_isolation_rejects_non_rgba_png(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            icon = Image.new("LA", (32, 32), (0, 0))
            icon.paste((128, 255), (4, 4, 28, 28))
            self._replace_icon_asset(candidate, icon)
            result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ICON_ALPHA_CONTENT_INVALID", {item["code"] for item in result["errors"]})

    def test_alpha_isolation_rejects_foreground_touching_edge(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._add_valid_icon_contract(candidate, Path(directory))
            icon = Image.new("RGBA", (32, 32), (0, 0, 0, 0))
            icon.paste((241, 90, 34, 255), (0, 4, 28, 28))
            self._replace_icon_asset(candidate, icon)
            result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_ICON_FOREGROUND_TOUCHES_EDGE", {item["code"] for item in result["errors"]})

    def test_same_session_reuse_requires_verified_artifacts(self):
        candidate = valid_spec()
        candidate["session_reuse"] = {"mode": "same_session_reuse", "reason": "continue", "artifacts": []}
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_SESSION_ARTIFACTS_INVALID", {item["code"] for item in result["errors"]})

    def test_element_requires_style_and_content(self):
        candidate = valid_spec()
        del candidate["elements"][0]["style"]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_ELEMENT_FIELD_MISSING", codes)

    def test_text_runs_must_cover_full_text(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["items"][0]["runs"][0]["end"] = 1
        result = MODULE.validate_spec(candidate, stage="prebuild")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_TEXT_RUN_COVERAGE_INVALID", codes)

    def test_text_run_requires_font_weight(self):
        candidate = valid_spec()
        del candidate["modules"]["typography"]["items"][0]["runs"][0]["font_weight"]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_TEXT_RUN_STYLE_INVALID", {item["code"] for item in result["errors"]})

    def test_text_run_rejects_boolean_font_weight(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["items"][0]["runs"][0]["font_weight"] = True
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_TEXT_RUN_STYLE_INVALID", {item["code"] for item in result["errors"]})

    def test_text_run_rejects_out_of_range_font_weight(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["items"][0]["runs"][0]["font_weight"] = 1200
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_TEXT_RUN_STYLE_INVALID", {item["code"] for item in result["errors"]})

    def test_paragraph_requires_list_contract(self):
        candidate = valid_spec()
        del candidate["modules"]["typography"]["items"][0]["paragraphs"][0]["list"]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_PARAGRAPH_LIST_INVALID", {item["code"] for item in result["errors"]})

    def test_native_list_requires_complete_bullet_contract(self):
        candidate = valid_list_spec()
        del candidate["modules"]["typography"]["items"][0]["paragraphs"][0]["list"]["bullet_type"]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_NATIVE_LIST_CONTRACT_INVALID", {item["code"] for item in result["errors"]})

    def test_native_list_requires_margin_and_hanging_indent(self):
        candidate = valid_list_spec()
        del candidate["modules"]["typography"]["items"][0]["paragraphs"][0]["margin_left"]
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_NATIVE_LIST_INDENT_INVALID", {item["code"] for item in result["errors"]})

    def test_text_box_paragraph_breaks_match_paragraph_boundaries(self):
        candidate = valid_list_spec()
        candidate["modules"]["typography"]["items"][0]["text_box"]["paragraph_breaks"] = []
        result = MODULE.validate_spec(candidate, stage="prebuild")
        self.assertIn("SPEC_PARAGRAPH_BREAKS_INVALID", {item["code"] for item in result["errors"]})

    def test_final_stage_requires_both_gates_passed(self):
        candidate = valid_spec()
        candidate["verification_profile"] = "strict"
        candidate["delivery_status"] = "strict_gate_passed"
        result = MODULE.validate_spec(candidate, stage="final")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_VISUAL_GATE_NOT_PASSED", codes)
        self.assertIn("SPEC_EDITABILITY_GATE_NOT_PASSED", codes)

    def test_final_stage_requires_current_artifact_identities(self):
        candidate = valid_spec()
        candidate["visual_gate"] = {"status": "passed", "evidence": ["overlay.png"]}
        candidate["editability_gate"] = {"status": "passed", "evidence": ["validator.json"]}
        candidate["modules"]["typography"]["items"][0]["font_declaration_verified"] = True
        result = MODULE.validate_spec(candidate, stage="final")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_GATE_ARTIFACT_MISSING", codes)

    def test_final_stage_requires_complete_gate_reviews(self):
        candidate = valid_spec()
        candidate["verification_profile"] = "strict"
        candidate["delivery_status"] = "strict_gate_passed"
        candidate["visual_gate"] = {"status": "passed", "evidence": "overlay.png"}
        candidate["editability_gate"] = {"status": "passed", "evidence": ["validator.json"]}
        candidate["modules"]["typography"]["items"][0]["font_declaration_verified"] = True
        result = MODULE.validate_spec(candidate, stage="final")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_GATE_EVIDENCE_INVALID", codes)
        self.assertIn("SPEC_VISUAL_REVIEW_INVALID", codes)
        self.assertIn("SPEC_EDITABILITY_REVIEW_INVALID", codes)

    def test_final_stage_accepts_matching_current_artifacts(self):
        candidate = self._valid_final_spec(profile="strict")
        result = MODULE.validate_spec(candidate, stage="final")
        self.assertTrue(result["valid"], result)

    def test_final_rejects_non_pptx_payload(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            pptx_path = Path(candidate["visual_gate"]["pptx"]["path"])
            pptx_path.write_bytes(b"not-a-pptx")
            pptx_identity = image_identity(pptx_path)
            candidate["visual_gate"]["pptx"] = dict(pptx_identity)
            candidate["editability_gate"]["pptx"] = dict(pptx_identity)
            validator_path = Path(candidate["editability_gate"]["validator"]["path"])
            validator = json.loads(validator_path.read_text(encoding="utf-8"))
            validator["pptx_sha256"] = pptx_identity["sha256"]
            validator_path.write_text(json.dumps(validator), encoding="utf-8")
            candidate["editability_gate"]["validator"] = image_identity(validator_path)

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_CURRENT_PPTX_VALIDATION_FAILED",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_non_image_preview(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            preview_path = Path(candidate["visual_gate"]["preview"]["path"])
            preview_path.write_bytes(b"not-an-image")
            preview = image_identity(preview_path)
            candidate["visual_gate"]["preview"] = preview
            candidate["visual_gate"]["reviewer"]["preview_sha256"] = preview["sha256"]

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_PREVIEW_IMAGE_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_invalid_visual_diff_json(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            report_path = Path(candidate["visual_gate"]["report"]["path"])
            report_path.write_bytes(b"not-json")
            candidate["visual_gate"]["report"] = image_identity(report_path)

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_DIFF_REPORT_INVALID",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_skipped_visual_evidence(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            report_path = Path(candidate["visual_gate"]["report"]["path"])
            report = json.loads(report_path.read_text(encoding="utf-8"))
            report["region_summary"]["skipped"] = 1
            report_path.write_text(json.dumps(report), encoding="utf-8")
            candidate["visual_gate"]["report"] = image_identity(report_path)

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_DIFF_EVIDENCE_INCOMPLETE",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_visual_diff_source_hash_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            report_path = Path(candidate["visual_gate"]["report"]["path"])
            report = json.loads(report_path.read_text(encoding="utf-8"))
            report["reference"]["sha256"] = "f" * 64
            report_path.write_text(json.dumps(report), encoding="utf-8")
            candidate["visual_gate"]["report"] = image_identity(report_path)

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_DIFF_SOURCE_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_reviewer_page_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            candidate["visual_gate"]["reviewer"]["page_id"] = "page-wrong"

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_REVIEW_PAGE_MISMATCH",
            {item["code"] for item in result["errors"]},
        )

    def test_final_rejects_missing_independent_visual_reviewer(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            del candidate["visual_gate"]["reviewer"]

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_INDEPENDENT_VISUAL_REVIEW_REQUIRED",
            {entry["code"] for entry in result["errors"]},
        )

    def test_final_requires_visual_review_round(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            del candidate["visual_gate"]["review_round"]
            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_REVIEW_ROUND_INVALID",
            {entry["code"] for entry in result["errors"]},
        )

    def test_final_rejects_visual_review_round_outside_one_to_two(self):
        for review_round in (0, 3, 4, True, 1.5):
            with self.subTest(review_round=review_round):
                candidate = valid_spec()
                with tempfile.TemporaryDirectory() as directory:
                    self._attach_final_gates(
                        candidate,
                        Path(directory),
                        {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                    )
                    candidate["visual_gate"]["review_round"] = review_round
                    result = MODULE.validate_spec(candidate, stage="final")

                self.assertIn(
                    "SPEC_VISUAL_REVIEW_ROUND_INVALID",
                    {entry["code"] for entry in result["errors"]},
                )

    def test_final_accepts_visual_review_rounds_one_to_two(self):
        for review_round in (1, 2):
            with self.subTest(review_round=review_round):
                candidate = self._valid_final_spec(
                    profile="strict",
                    review_round=review_round,
                )
                result = MODULE.validate_spec(candidate, stage="final")

                self.assertTrue(result["valid"], result)

    def test_final_requires_complete_visual_review_coverage(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            del candidate["visual_gate"]["reviewer"]["coverage"]
            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_REVIEW_COVERAGE_INVALID",
            {entry["code"] for entry in result["errors"]},
        )

    def test_final_rejects_unknown_or_invalid_visual_coverage(self):
        mutations = (
            lambda coverage: coverage.update({"unknown": "checked"}),
            lambda coverage: coverage.update({"canvas_and_regions": "skipped"}),
            lambda coverage: coverage.update({"canvas_and_regions": "not_reviewable"}),
        )
        for mutate in mutations:
            candidate = valid_spec()
            with tempfile.TemporaryDirectory() as directory:
                self._attach_final_gates(
                    candidate,
                    Path(directory),
                    {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                )
                mutate(candidate["visual_gate"]["reviewer"]["coverage"])
                result = MODULE.validate_spec(candidate, stage="final")

            self.assertIn(
                "SPEC_VISUAL_REVIEW_COVERAGE_INVALID",
                {entry["code"] for entry in result["errors"]},
            )

    def test_final_rejects_unhashable_visual_coverage_values_without_crashing(self):
        for invalid in ([], {}, 1, None):
            with self.subTest(invalid=invalid):
                candidate = valid_spec()
                with tempfile.TemporaryDirectory() as directory:
                    self._attach_final_gates(
                        candidate,
                        Path(directory),
                        {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                    )
                    candidate["visual_gate"]["reviewer"]["coverage"][
                        "canvas_and_regions"
                    ] = invalid
                    result = MODULE.validate_spec(candidate, stage="final")

                self.assertIn(
                    "SPEC_VISUAL_REVIEW_COVERAGE_INVALID",
                    {entry["code"] for entry in result["errors"]},
                )

    def test_final_requires_applicable_visual_coverage_to_be_checked(self):
        for category in (
            "canvas_and_regions",
            "objects_and_geometry",
            "text_and_typography",
        ):
            with self.subTest(category=category):
                candidate = valid_spec()
                with tempfile.TemporaryDirectory() as directory:
                    self._attach_final_gates(
                        candidate,
                        Path(directory),
                        {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                    )
                    candidate["visual_gate"]["reviewer"]["coverage"][category] = (
                        "not_applicable"
                    )
                    result = MODULE.validate_spec(candidate, stage="final")

                self.assertIn(
                    "SPEC_VISUAL_REVIEW_COVERAGE_INVALID",
                    {entry["code"] for entry in result["errors"]},
                )

    def test_final_requires_complete_finding_fields(self):
        finding = {
            "severity": "P2",
            "category": "pictures_crop_layers",
            "location": "右上角图片",
            "source_fact": "原图边缘为圆角",
            "observed_difference": "预览圆角半径略小",
            "evidence": "region-picture.png",
        }
        for field in tuple(finding):
            with self.subTest(field=field):
                candidate = valid_spec()
                with tempfile.TemporaryDirectory() as directory:
                    self._attach_final_gates(
                        candidate,
                        Path(directory),
                        {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                    )
                    incomplete = dict(finding)
                    del incomplete[field]
                    candidate["visual_gate"]["reviewer"]["findings"] = [incomplete]
                    result = MODULE.validate_spec(candidate, stage="final")

                self.assertIn(
                    "SPEC_INDEPENDENT_VISUAL_REVIEW_INVALID",
                    {entry["code"] for entry in result["errors"]},
                )

    def test_final_rejects_non_string_finding_enums_without_crashing(self):
        base = {
            "severity": "P2",
            "category": "pictures_crop_layers",
            "location": "右上角图片",
            "source_fact": "原图边缘为圆角",
            "observed_difference": "预览圆角半径略小",
            "evidence": "region-picture.png",
        }
        for field in ("severity", "category"):
            for invalid in ([], {}, 1, None):
                with self.subTest(field=field, invalid=invalid):
                    candidate = valid_spec()
                    with tempfile.TemporaryDirectory() as directory:
                        self._attach_final_gates(
                            candidate,
                            Path(directory),
                            {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                        )
                        finding = dict(base)
                        finding[field] = invalid
                        candidate["visual_gate"]["reviewer"]["findings"] = [finding]
                        result = MODULE.validate_spec(candidate, stage="final")

                    self.assertIn(
                        "SPEC_INDEPENDENT_VISUAL_REVIEW_INVALID",
                        {entry["code"] for entry in result["errors"]},
                    )

    def test_final_rejects_decision_and_findings_inconsistency(self):
        blocking = {
            "severity": "P1",
            "category": "text_and_typography",
            "location": "标题",
            "source_fact": "原图标题完整",
            "observed_difference": "预览标题被截断",
            "evidence": "region-title.png",
        }
        scenarios = (
            ("passed", [blocking]),
            ("not_reviewable", [blocking]),
            ("changes_required", []),
        )
        for decision, findings in scenarios:
            with self.subTest(decision=decision):
                candidate = valid_spec()
                with tempfile.TemporaryDirectory() as directory:
                    self._attach_final_gates(
                        candidate,
                        Path(directory),
                        {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                    )
                    reviewer = candidate["visual_gate"]["reviewer"]
                    reviewer["decision"] = decision
                    reviewer["findings"] = findings
                    result = MODULE.validate_spec(candidate, stage="final")

                self.assertIn(
                    "SPEC_INDEPENDENT_VISUAL_REVIEW_INVALID",
                    {entry["code"] for entry in result["errors"]},
                )

    def test_final_rejects_non_string_reviewer_decision_without_crashing(self):
        for invalid in ([], {}, 1, None):
            with self.subTest(invalid=invalid):
                candidate = valid_spec()
                with tempfile.TemporaryDirectory() as directory:
                    self._attach_final_gates(
                        candidate,
                        Path(directory),
                        {"valid": True, "errors": [], "native_list_contracts_checked": 0},
                    )
                    candidate["visual_gate"]["reviewer"]["decision"] = invalid
                    result = MODULE.validate_spec(candidate, stage="final")

                self.assertIn(
                    "SPEC_INDEPENDENT_VISUAL_REVIEW_INVALID",
                    {entry["code"] for entry in result["errors"]},
                )

    def test_final_rejects_validator_report_for_different_pptx(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {
                    "valid": True,
                    "errors": [],
                    "native_list_contracts_checked": 0,
                    "pptx_sha256": "f" * 64,
                },
            )
            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VALIDATOR_PPTX_MISMATCH",
            {entry["code"] for entry in result["errors"]},
        )

    def test_final_rejects_open_visual_p1(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            candidate["visual_gate"]["reviewer"]["findings"] = [
                {
                    "severity": "P1",
                    "category": "text_and_typography",
                    "location": "右侧文本框",
                    "source_fact": "原图文本完整显示",
                    "observed_difference": "预览末行被截断",
                    "evidence": "region-right.png",
                }
            ]

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_OPEN_BLOCKING_DIFFERENCE",
            {entry["code"] for entry in result["errors"]},
        )

    def test_final_rejects_reviewed_preview_hash_mismatch(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            candidate["visual_gate"]["reviewer"]["preview_sha256"] = "f" * 64

            result = MODULE.validate_spec(candidate, stage="final")

        self.assertIn(
            "SPEC_VISUAL_REVIEW_PREVIEW_MISMATCH",
            {entry["code"] for entry in result["errors"]},
        )

    def test_final_stage_rejects_failed_validator_report(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": False, "errors": ["NATIVE_LIST_TEXTBOX_MISSING"], "native_list_contracts_checked": 0},
            )
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn("SPEC_VALIDATOR_REPORT_FAILED", {item["code"] for item in result["errors"]})

    def test_final_stage_requires_native_list_contract_evidence(self):
        candidate = valid_list_spec()
        with tempfile.TemporaryDirectory() as directory:
            self._attach_final_gates(
                candidate,
                Path(directory),
                {"valid": True, "errors": [], "native_list_contracts_checked": 0},
            )
            result = MODULE.validate_spec(candidate, stage="final")
        self.assertIn("SPEC_NATIVE_LIST_VALIDATION_MISSING", {item["code"] for item in result["errors"]})

    def test_final_stage_rejects_open_p1(self):
        candidate = valid_spec()
        candidate["visual_gate"] = {"status": "passed", "evidence": ["overlay.png"]}
        candidate["editability_gate"] = {"status": "passed", "evidence": ["validator.json"]}
        candidate["activated_modules"].append("high_risk")
        candidate["modules"]["high_risk"] = {
            "items": [
                {
                    "risk_id": "layout-diff",
                    "severity": "P1",
                    "result": "changes_required",
                }
            ]
        }
        result = MODULE.validate_spec(candidate, stage="final")
        codes = {item["code"] for item in result["errors"]}
        self.assertIn("SPEC_OPEN_BLOCKING_DIFFERENCE", codes)
        self.assertIn("SPEC_HIGH_RISK_ITEM_INVALID", codes)

    def test_cli_writes_json_and_returns_failure(self):
        candidate = valid_spec()
        candidate["modules"]["typography"]["slide_coordinate_unit"] = "px"
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "page-reconstruction.json"
            path.write_text(json.dumps(candidate), encoding="utf-8")
            with contextlib.redirect_stdout(io.StringIO()):
                exit_code = MODULE.main([str(path), "--stage", "prebuild"])
        self.assertEqual(1, exit_code)

    def test_cli_output_matches_stdout_json(self):
        candidate = valid_spec()
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            spec_path = root / "page-reconstruction.json"
            report_path = root / "reports" / "prebuild.json"
            spec_path.write_text(json.dumps(candidate), encoding="utf-8")
            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = MODULE.main(
                    [
                        str(spec_path),
                        "--stage",
                        "prebuild",
                        "--output",
                        str(report_path),
                    ]
                )

            self.assertEqual(0, exit_code)
            self.assertTrue(report_path.is_file())
            self.assertEqual(
                json.loads(stdout.getvalue()),
                json.loads(report_path.read_text(encoding="utf-8")),
            )


if __name__ == "__main__":
    unittest.main()
