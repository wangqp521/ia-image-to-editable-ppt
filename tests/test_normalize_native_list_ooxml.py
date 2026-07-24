from __future__ import annotations

import copy
import hashlib
import importlib.util
import json
import tempfile
import unittest
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement


ROOT = Path(__file__).resolve().parents[1]
SCRIPT_PATH = ROOT / "scripts" / "normalize_native_list_ooxml.py"
VALIDATOR_PATH = ROOT / "scripts" / "validate_pptx.py"
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def _load_module(path: Path, name: str):
    module_spec = importlib.util.spec_from_file_location(name, path)
    if module_spec is None or module_spec.loader is None:
        raise RuntimeError(f"Cannot load {path}")
    module = importlib.util.module_from_spec(module_spec)
    module_spec.loader.exec_module(module)
    return module


VALIDATOR = _load_module(VALIDATOR_PATH, "ia_validate_pptx_for_normalizer_test")
NORMALIZER = (
    _load_module(SCRIPT_PATH, "ia_native_list_normalizer_test")
    if SCRIPT_PATH.is_file()
    else None
)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


class NormalizeNativeListOoxmlTest(unittest.TestCase):
    X = 900_000
    Y = 900_000
    W = 3_500_000
    H = 1_000_000
    MARGIN_LEFT = 342_900
    INDENT = -228_600
    ITEMS = ("第一项", "第二项")

    def setUp(self) -> None:
        if NORMALIZER is None:
            self.fail(f"normalizer script is missing: {SCRIPT_PATH}")
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        self.input_pptx = self.root / "input.pptx"
        self.output_pptx = self.root / "output.pptx"
        self.spec_path = self.root / "page-reconstruction.json"
        self.report_path = self.root / "normalization-report.json"

    def tearDown(self) -> None:
        if hasattr(self, "temp_dir"):
            self.temp_dir.cleanup()

    def _spec(self) -> dict:
        paragraphs = []
        cursor = 0
        for value in self.ITEMS:
            end = cursor + len(value)
            paragraphs.append(
                {
                    "start": cursor,
                    "end": end,
                    "margin_left": self.MARGIN_LEFT,
                    "indent": self.INDENT,
                    "list": {
                        "is_list": True,
                        "level": 0,
                        "bullet_type": "char",
                        "bullet": "•",
                        "bullet_font": "follow_text",
                        "bullet_size_mode": "follow_text",
                        "bullet_size_value": None,
                        "bullet_color": "follow_text",
                    },
                }
            )
            cursor = end
        return {
            "canvas": {"slide_size_emu": [12_192_000, 6_858_000]},
            "modules": {
                "typography": {
                    "items": [
                        {
                            "element_id": "list-01",
                            "text": "".join(self.ITEMS),
                            "paragraphs": paragraphs,
                            "text_box": {
                                "x": self.X,
                                "y": self.Y,
                                "w": self.W,
                                "h": self.H,
                                "paragraph_breaks": [len(self.ITEMS[0])],
                            },
                        }
                    ]
                }
            },
        }

    def _add_bullet_style(self, paragraph, *, follow_text: bool) -> None:
        properties = paragraph._p.get_or_add_pPr()
        properties.set("lvl", "0")
        properties.set("marL", str(self.MARGIN_LEFT))
        properties.set("indent", str(self.INDENT))
        if follow_text:
            for tag in ("a:buClrTx", "a:buSzTx", "a:buFontTx"):
                properties.append(OxmlElement(tag))
        else:
            color = OxmlElement("a:buClr")
            rgb = OxmlElement("a:srgbClr")
            rgb.set("val", "303030")
            color.append(rgb)
            size = OxmlElement("a:buSzPct")
            size.set("val", "100000")
            font = OxmlElement("a:buFont")
            font.set("typeface", "Noto Sans CJK SC")
            for child in (color, size, font):
                properties.append(child)
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        properties.append(bullet)

    def _build_pptx(self, *, follow_text: bool) -> None:
        presentation = Presentation()
        presentation.slide_width = 12_192_000
        presentation.slide_height = 6_858_000
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(self.X, self.Y, self.W, self.H)
        text_box.name = "ia:list-01"
        frame = text_box.text_frame
        for index, value in enumerate(self.ITEMS):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = value
            self._add_bullet_style(paragraph, follow_text=follow_text)
        presentation.save(self.input_pptx)

    def _write_spec(self, spec: dict) -> None:
        self.spec_path.write_text(
            json.dumps(spec, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )

    def _paragraph_properties(self, pptx: Path) -> list[ET.Element]:
        with zipfile.ZipFile(pptx) as archive:
            root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
        return root.findall(".//a:p/a:pPr", NS)

    def test_follow_text_styles_replace_explicit_style_nodes(self) -> None:
        self._build_pptx(follow_text=False)
        self._write_spec(self._spec())

        result = NORMALIZER.normalize_pptx(
            self.input_pptx,
            self.spec_path,
            self.output_pptx,
            self.report_path,
        )

        self.assertEqual(result["paragraphs_changed"], 2)
        for properties in self._paragraph_properties(self.output_pptx):
            self.assertIsNotNone(properties.find("a:buFontTx", NS))
            self.assertIsNone(properties.find("a:buFont", NS))
            self.assertIsNotNone(properties.find("a:buSzTx", NS))
            self.assertIsNone(properties.find("a:buSzPct", NS))
            self.assertIsNotNone(properties.find("a:buClrTx", NS))
            self.assertIsNone(properties.find("a:buClr", NS))
            self.assertIsNotNone(properties.find("a:buChar", NS))

    def test_explicit_styles_are_serialized_without_follow_text_nodes(self) -> None:
        self._build_pptx(follow_text=True)
        spec = self._spec()
        for paragraph in spec["modules"]["typography"]["items"][0]["paragraphs"]:
            paragraph["list"].update(
                {
                    "bullet_font": "Arial",
                    "bullet_size_mode": "percent",
                    "bullet_size_value": 90,
                    "bullet_color": "#336699",
                }
            )
        self._write_spec(spec)

        NORMALIZER.normalize_pptx(
            self.input_pptx,
            self.spec_path,
            self.output_pptx,
            self.report_path,
        )

        for properties in self._paragraph_properties(self.output_pptx):
            font = properties.find("a:buFont", NS)
            size = properties.find("a:buSzPct", NS)
            rgb = properties.find("a:buClr/a:srgbClr", NS)
            self.assertEqual(font.get("typeface"), "Arial")
            self.assertEqual(size.get("val"), "90000")
            self.assertEqual(rgb.get("val"), "336699")
            self.assertIsNone(properties.find("a:buFontTx", NS))
            self.assertIsNone(properties.find("a:buSzTx", NS))
            self.assertIsNone(properties.find("a:buClrTx", NS))

    def test_already_normalized_pptx_is_copied_byte_for_byte(self) -> None:
        self._build_pptx(follow_text=True)
        self._write_spec(self._spec())

        result = NORMALIZER.normalize_pptx(
            self.input_pptx,
            self.spec_path,
            self.output_pptx,
            self.report_path,
        )

        self.assertEqual(result["paragraphs_changed"], 0)
        self.assertEqual(_sha256(self.input_pptx), _sha256(self.output_pptx))

    def test_paragraph_count_mismatch_fails_without_output_pptx(self) -> None:
        self._build_pptx(follow_text=False)
        spec = self._spec()
        extra = copy.deepcopy(spec["modules"]["typography"]["items"][0]["paragraphs"][-1])
        extra["start"] = 6
        extra["end"] = 9
        spec["modules"]["typography"]["items"][0]["paragraphs"].append(extra)
        self._write_spec(spec)

        with self.assertRaises(NORMALIZER.NormalizeError) as caught:
            NORMALIZER.normalize_pptx(
                self.input_pptx,
                self.spec_path,
                self.output_pptx,
                self.report_path,
            )

        self.assertEqual(caught.exception.code, "NORMALIZE_PARAGRAPH_COUNT_MISMATCH")
        self.assertFalse(self.output_pptx.exists())

    def test_normalized_output_passes_native_list_structure_gate(self) -> None:
        self._build_pptx(follow_text=False)
        spec = self._spec()
        self._write_spec(spec)
        NORMALIZER.normalize_pptx(
            self.input_pptx,
            self.spec_path,
            self.output_pptx,
            self.report_path,
        )

        result = VALIDATOR.validate_pptx(self.output_pptx, 1, spec)

        self.assertTrue(result["valid"], result)
        self.assertEqual(result["native_list_paragraphs"], 2)
        self.assertEqual(result["native_list_contracts_checked"], 1)


if __name__ == "__main__":
    unittest.main()
