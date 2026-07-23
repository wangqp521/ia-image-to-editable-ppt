"""Regression tests for the editable-object and full-slide-picture checks."""

from __future__ import annotations

import importlib.util
import hashlib
import json
import os
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement


VALIDATOR_PATH = Path(__file__).parents[1] / "scripts" / "validate_pptx.py"
MERGER_PATH = Path(__file__).parents[1] / "scripts" / "merge_pptx.py"
SPEC = importlib.util.spec_from_file_location("ia_validate_pptx_test", VALIDATOR_PATH)
assert SPEC is not None and SPEC.loader is not None
VALIDATOR = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(VALIDATOR)


class ValidatePptxPictureCoverageTest(unittest.TestCase):
    def _presentation(self, directory: Path, name: str, *, full_picture: bool, text: bool) -> Path:
        image_path = directory / "source.png"
        Image.new("RGB", (1600, 900), "#27506b").save(image_path)

        presentation = Presentation()
        presentation.slide_width = 12_192_000
        presentation.slide_height = 6_858_000
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if full_picture:
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
        else:
            slide.shapes.add_picture(str(image_path), 914_400, 914_400, 914_400, 514_350)
        if text:
            text_box = slide.shapes.add_textbox(914_400, 914_400, 4_572_000, 914_400)
            text_box.name = "ia:title"
            self.assertEqual(text_box.shape_type, MSO_SHAPE_TYPE.TEXT_BOX)
            text_box.text_frame.paragraphs[0].text = "可编辑标题"
        output = directory / name
        presentation.save(output)
        return output

    def _merge_spec(
        self,
        directory: Path,
        pptx_path: Path,
        page_id: str,
        *,
        visual_status: str = "passed",
        profile: str = "strict",
    ) -> Path:
        source_path = directory / "source.png"
        preview_path = directory / f"{page_id}-preview.png"
        Image.new("RGB", (1600, 900), "#27506B").save(preview_path)

        def identity(path: Path) -> dict:
            return {
                "path": str(path.resolve()),
                "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            }

        source = identity(source_path)
        preview = identity(preview_path)
        pptx = identity(pptx_path)
        decision = "passed" if visual_status == "passed" else "changes_required"
        delivery_status = {
            "rapid": "rapid_validated" if visual_status == "passed" else "rapid_validation_failed",
            "reviewed": "reviewed_passed" if visual_status == "passed" else "reviewed_failed",
            "strict": "strict_gate_passed" if visual_status == "passed" else "strict_gate_failed",
        }[profile]
        spec = {
            "page_id": page_id,
            "verification_profile": profile,
            "delivery_status": delivery_status,
            "clean_visual_reference": source,
            "elements": [
                {
                    "element_id": "title",
                    "kind": "text",
                    "slide_bbox": [914_400, 914_400, 4_572_000, 914_400],
                    "editable": True,
                    "content": {"text": "可编辑标题"},
                }
            ],
            "visual_gate": {
                "status": (
                    "not_independently_reviewed"
                    if profile == "rapid" and visual_status == "passed"
                    else visual_status
                ),
                "pptx": pptx,
                "preview": preview,
            },
            "editability_gate": {"status": "passed", "pptx": pptx},
        }
        if profile != "rapid":
            spec["visual_gate"]["reviewer"] = {
                "page_id": page_id,
                "decision": decision,
                "source_sha256": source["sha256"],
                "preview_sha256": preview["sha256"],
            }
        validation = VALIDATOR.validate_pptx(
            pptx_path, expected_slides=1, reconstruction_spec=spec
        )
        validator_path = directory / f"{page_id}-validator.json"
        validator_path.write_text(json.dumps(validation), encoding="utf-8")
        spec["editability_gate"]["validator"] = identity(validator_path)
        spec_path = directory / f"{page_id}.json"
        spec_path.write_text(json.dumps(spec), encoding="utf-8")
        return spec_path

    def test_full_slide_background_with_text_is_valid_but_risky(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self._presentation(
                Path(temp_dir), "background-with-text.pptx", full_picture=True, text=True
            )
            result = VALIDATOR.validate_pptx(output, expected_slides=1)

        self.assertTrue(result["valid"])
        self.assertTrue(result["full_slide_picture_risk"])
        self.assertNotIn("FULL_SLIDE_PICTURE_ONLY", result["errors"])
        self.assertIn("FULL_SLIDE_PICTURE_WITH_EDITABLE_OBJECTS", result["warnings"])

    def test_validation_result_binds_current_pptx_sha256(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self._presentation(
                Path(temp_dir), "hash-bound.pptx", full_picture=True, text=True
            )
            expected = hashlib.sha256(output.read_bytes()).hexdigest()
            result = VALIDATOR.validate_pptx(output, expected_slides=1)

        self.assertEqual(expected, result["pptx_sha256"])

    def test_unreadable_pptx_returns_structured_error(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self._presentation(
                Path(temp_dir), "unreadable.pptx", full_picture=True, text=True
            )
            os.chmod(output, 0)
            try:
                result = VALIDATOR.validate_pptx(output, expected_slides=1)
            finally:
                os.chmod(output, 0o600)

        self.assertFalse(result["valid"])
        self.assertIn("PPTX_ZIP_INVALID", result["errors"])

    def test_picture_only_slide_is_rejected_for_lacking_editable_objects(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self._presentation(
                Path(temp_dir), "picture-only.pptx", full_picture=False, text=False
            )
            result = VALIDATOR.validate_pptx(output, expected_slides=1)

        self.assertFalse(result["valid"])
        self.assertIn("NO_EDITABLE_OBJECTS", result["errors"])

    def test_full_slide_picture_without_editable_objects_is_rejected(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self._presentation(
                Path(temp_dir), "full-picture-only.pptx", full_picture=True, text=False
            )
            result = VALIDATOR.validate_pptx(output, expected_slides=1)

        self.assertFalse(result["valid"])
        self.assertIn("FULL_SLIDE_PICTURE_ONLY", result["errors"])
        self.assertIn("NO_EDITABLE_OBJECTS", result["errors"])

    def test_summary_cli_omits_verbose_object_arrays(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = self._presentation(
                Path(temp_dir), "background-with-text.pptx", full_picture=True, text=True
            )
            completed = subprocess.run(
                [sys.executable, str(VALIDATOR_PATH), str(output), "--summary"],
                check=False,
                capture_output=True,
                text=True,
            )

        self.assertEqual(completed.returncode, 0, completed.stderr)
        result = json.loads(completed.stdout)
        self.assertNotIn("text_objects", result)
        self.assertNotIn("native_shape_objects", result)
        self.assertNotIn("picture_objects", result)
        self.assertNotIn("structure_objects", result)
        self.assertNotIn("picture_objects", result["slides"][0])
        self.assertEqual(result["slides"][0]["editable_object_count"], 1)

    def test_merge_accepts_validated_full_slide_background_pages(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            first = self._presentation(
                directory, "page-001.pptx", full_picture=True, text=True
            )
            second = self._presentation(
                directory, "page-002.pptx", full_picture=True, text=True
            )
            first_spec = self._merge_spec(directory, first, "page-001")
            second_spec = self._merge_spec(directory, second, "page-002")
            merged = directory / "merged.pptx"
            completed = subprocess.run(
                [
                    sys.executable,
                    str(MERGER_PATH),
                    "--input",
                    str(first),
                    "--spec",
                    str(first_spec),
                    "--input",
                    str(second),
                    "--spec",
                    str(second_spec),
                    "--output",
                    str(merged),
                ],
                check=False,
                capture_output=True,
                text=True,
            )

        self.assertEqual(completed.returncode, 0, completed.stderr)
        result = json.loads(completed.stdout)
        self.assertTrue(result["ok"])
        self.assertEqual(result["result"]["slide_count"], 2)
        self.assertTrue(result["result"]["validation"]["valid"])

    def test_merge_requires_one_spec_per_input(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            page = self._presentation(
                directory, "page-001.pptx", full_picture=True, text=True
            )
            completed = subprocess.run(
                [
                    sys.executable,
                    str(MERGER_PATH),
                    "--input",
                    str(page),
                    "--output",
                    str(directory / "merged.pptx"),
                ],
                check=False,
                capture_output=True,
                text=True,
            )

        self.assertEqual(completed.returncode, 2)
        self.assertEqual("SPEC_COUNT_MISMATCH", json.loads(completed.stderr)["code"])

    def test_merge_rejects_duplicate_page_id(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            first = self._presentation(directory, "one.pptx", full_picture=True, text=True)
            second = self._presentation(directory, "two.pptx", full_picture=True, text=True)
            first_spec = self._merge_spec(directory, first, "page-001")
            second_spec = self._merge_spec(directory, second, "page-001")
            completed = subprocess.run(
                [
                    sys.executable, str(MERGER_PATH),
                    "--input", str(first), "--spec", str(first_spec),
                    "--input", str(second), "--spec", str(second_spec),
                    "--output", str(directory / "merged.pptx"),
                ],
                check=False, capture_output=True, text=True,
            )

        self.assertEqual(2, completed.returncode)
        self.assertEqual("PAGE_ID_DUPLICATE", json.loads(completed.stderr)["code"])

    def test_merge_rejects_input_hash_not_bound_by_spec(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            page = self._presentation(directory, "page.pptx", full_picture=True, text=True)
            spec_path = self._merge_spec(directory, page, "page-001")
            spec = json.loads(spec_path.read_text(encoding="utf-8"))
            spec["visual_gate"]["pptx"]["sha256"] = "f" * 64
            spec["editability_gate"]["pptx"]["sha256"] = "f" * 64
            spec_path.write_text(json.dumps(spec), encoding="utf-8")
            completed = subprocess.run(
                [
                    sys.executable, str(MERGER_PATH),
                    "--input", str(page), "--spec", str(spec_path),
                    "--output", str(directory / "merged.pptx"),
                ],
                check=False, capture_output=True, text=True,
            )

        self.assertEqual(2, completed.returncode)
        self.assertEqual("INPUT_SPEC_HASH_MISMATCH", json.loads(completed.stderr)["code"])

    def test_merge_preserves_failed_visual_status_in_output_label(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            page = self._presentation(directory, "page.pptx", full_picture=True, text=True)
            spec_path = self._merge_spec(
                directory, page, "page-001", visual_status="changes_required"
            )
            completed = subprocess.run(
                [
                    sys.executable, str(MERGER_PATH),
                    "--input", str(page), "--spec", str(spec_path),
                    "--output", str(directory / "merged.pptx"),
                ],
                check=False, capture_output=True, text=True,
            )

        self.assertEqual(0, completed.returncode, completed.stderr)
        self.assertIn("delivery_label", json.loads(completed.stdout)["result"])
        self.assertEqual(
            "完整视觉门禁未通过版",
            json.loads(completed.stdout)["result"]["delivery_label"],
        )

    def test_merge_accepts_rapid_without_reviewer(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            page = self._presentation(directory, "page.pptx", full_picture=True, text=True)
            spec_path = self._merge_spec(directory, page, "page-001", profile="rapid")
            completed = subprocess.run(
                [
                    sys.executable, str(MERGER_PATH),
                    "--input", str(page), "--spec", str(spec_path),
                    "--output", str(directory / "merged.pptx"),
                ],
                check=False, capture_output=True, text=True,
            )

        self.assertEqual(0, completed.returncode, completed.stderr)
        result = json.loads(completed.stdout)["result"]
        self.assertEqual("rapid", result["verification_profile"])
        self.assertEqual("快速校验版", result["delivery_label"])

    def test_merge_reports_reviewed_delivery_label(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            page = self._presentation(directory, "page.pptx", full_picture=True, text=True)
            spec_path = self._merge_spec(directory, page, "page-001", profile="reviewed")
            completed = subprocess.run(
                [
                    sys.executable, str(MERGER_PATH),
                    "--input", str(page), "--spec", str(spec_path),
                    "--output", str(directory / "merged.pptx"),
                ],
                check=False, capture_output=True, text=True,
            )

        self.assertEqual(0, completed.returncode, completed.stderr)
        result = json.loads(completed.stdout)["result"]
        self.assertEqual("reviewed", result["verification_profile"])
        self.assertEqual("独立复核通过版", result["delivery_label"])

    def test_merge_rejects_mixed_verification_profiles(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            first = self._presentation(directory, "one.pptx", full_picture=True, text=True)
            second = self._presentation(directory, "two.pptx", full_picture=True, text=True)
            first_spec = self._merge_spec(directory, first, "page-001", profile="rapid")
            second_spec = self._merge_spec(directory, second, "page-002", profile="strict")
            completed = subprocess.run(
                [
                    sys.executable, str(MERGER_PATH),
                    "--input", str(first), "--spec", str(first_spec),
                    "--input", str(second), "--spec", str(second_spec),
                    "--output", str(directory / "merged.pptx"),
                ],
                check=False, capture_output=True, text=True,
            )

        self.assertEqual(2, completed.returncode)
        self.assertEqual(
            "VERIFICATION_PROFILE_MISMATCH",
            json.loads(completed.stderr)["code"],
        )


class ValidatePptxNativeListTest(unittest.TestCase):
    X = 900_000
    Y = 900_000
    W = 3_500_000
    H = 1_000_000
    MARGIN_LEFT = 342_900
    INDENT = -228_600
    ITEMS = ("第一项", "第二项")

    def test_auto_number_bullet_contract_is_extracted(self) -> None:
        properties = OxmlElement("a:pPr")
        bullet = OxmlElement("a:buAutoNum")
        bullet.set("type", "arabicPeriod")
        properties.append(bullet)

        result = VALIDATOR._native_bullet_contract([properties], 1)

        self.assertTrue(result["is_list"])
        self.assertEqual(result["level"], 1)
        self.assertEqual(result["bullet_type"], "auto_number")
        self.assertEqual(result["bullet"], "arabicPeriod")

    def test_picture_bullet_contract_is_extracted(self) -> None:
        properties = OxmlElement("a:pPr")
        properties.append(OxmlElement("a:buBlip"))

        result = VALIDATOR._native_bullet_contract([properties], 0)

        self.assertTrue(result["is_list"])
        self.assertEqual(result["bullet_type"], "picture")
        self.assertEqual(result["bullet"], "blip")

    def test_explicit_bullet_style_is_extracted(self) -> None:
        properties = OxmlElement("a:pPr")
        font = OxmlElement("a:buFont")
        font.set("typeface", "Arial")
        size = OxmlElement("a:buSzPct")
        size.set("val", "95000")
        color = OxmlElement("a:buClr")
        rgb = OxmlElement("a:srgbClr")
        rgb.set("val", "336699")
        color.append(rgb)
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        for child in (font, size, color, bullet):
            properties.append(child)

        result = VALIDATOR._native_bullet_contract([properties], 0)

        self.assertEqual(result["bullet_font"], "Arial")
        self.assertEqual(result["bullet_size_mode"], "percent")
        self.assertEqual(result["bullet_size_value"], 95)
        self.assertEqual(result["bullet_color"], "#336699")

    def test_bullet_none_overrides_inherited_list(self) -> None:
        local = OxmlElement("a:pPr")
        local.append(OxmlElement("a:buNone"))
        inherited = OxmlElement("a:lvl1pPr")
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        inherited.append(bullet)

        result = VALIDATOR._native_bullet_contract([local, inherited], 0)

        self.assertEqual(result, {"is_list": False, "level": 0, "bullet": None})

    def _list_contract(self) -> dict:
        text = "".join(self.ITEMS)
        paragraphs = []
        cursor = 0
        for value in self.ITEMS:
            end = cursor + len(value)
            paragraphs.append(
                {
                    "start": cursor,
                    "end": end,
                    "alignment": "left",
                    "line_spacing": 1.0,
                    "space_before": 0,
                    "space_after": 0,
                    "margin_left": self.MARGIN_LEFT,
                    "indent": self.INDENT,
                    "list": {
                        "is_list": True,
                        "level": 0,
                        "bullet_type": "char",
                        "bullet": "•",
                        "bullet_font": "follow_text",
                        "bullet_size_mode": "follow_text",
                        "bullet_size_value": None,
                        "bullet_color": "follow_text",
                    },
                }
            )
            cursor = end
        return {
            "canvas": {"slide_size_emu": [12_192_000, 6_858_000]},
            "modules": {
                "typography": {
                    "items": [
                        {
                            "element_id": "list-01",
                            "text": text,
                            "paragraphs": paragraphs,
                            "text_box": {
                                "x": self.X,
                                "y": self.Y,
                                "w": self.W,
                                "h": self.H,
                                "paragraph_breaks": [len(self.ITEMS[0])],
                            },
                        }
                    ]
                }
            },
        }

    def _presentation(self) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = 12_192_000
        presentation.slide_height = 6_858_000
        presentation.slides.add_slide(presentation.slide_layouts[6])
        return presentation

    def _set_native_bullet(self, paragraph, *, bullet: str = "•", indent: int | None = None) -> None:
        properties = paragraph._p.get_or_add_pPr()
        properties.set("lvl", "0")
        properties.set("marL", str(self.MARGIN_LEFT))
        properties.set("indent", str(self.INDENT if indent is None else indent))
        bullet_element = OxmlElement("a:buChar")
        bullet_element.set("char", bullet)
        properties.append(bullet_element)

    def _add_native_list_textbox(self, slide) -> None:
        text_box = slide.shapes.add_textbox(self.X, self.Y, self.W, self.H)
        text_box.name = "ia:list-01"
        frame = text_box.text_frame
        for index, value in enumerate(self.ITEMS):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = value
            self._set_native_bullet(paragraph)

    def _add_split_native_list(self, slide) -> None:
        for index, value in enumerate(self.ITEMS):
            text_box = slide.shapes.add_textbox(
                self.X,
                self.Y + index * 450_000,
                self.W,
                400_000,
            )
            paragraph = text_box.text_frame.paragraphs[0]
            paragraph.text = value
            self._set_native_bullet(paragraph)

    def _add_fake_list(self, slide) -> None:
        for index, value in enumerate(self.ITEMS):
            y = self.Y + index * 450_000
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, self.X, y + 130_000, 80_000, 80_000)
            dot.line.fill.background()
            text_box = slide.shapes.add_textbox(self.X + 200_000, y, self.W, 400_000)
            text_box.text_frame.paragraphs[0].text = value

    def test_same_textbox_native_list_matches_spec(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "native-list.pptx"
            presentation = self._presentation()
            self._add_native_list_textbox(presentation.slides[0])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._list_contract())

        self.assertTrue(result["valid"], result)
        self.assertEqual(result["native_list_paragraphs"], 2)

    def test_fake_dot_shapes_are_rejected_against_list_spec(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "fake-list.pptx"
            presentation = self._presentation()
            self._add_fake_list(presentation.slides[0])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._list_contract())

        self.assertFalse(result["valid"])
        self.assertIn("NATIVE_LIST_TEXTBOX_MISSING", result["errors"])

    def test_split_native_list_textboxes_are_rejected(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "split-list.pptx"
            presentation = self._presentation()
            self._add_split_native_list(presentation.slides[0])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._list_contract())

        self.assertFalse(result["valid"])
        self.assertIn("NATIVE_LIST_TEXTBOX_MISSING", result["errors"])

    def test_native_list_bullet_character_must_match_spec(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "wrong-bullet.pptx"
            presentation = self._presentation()
            text_box = presentation.slides[0].shapes.add_textbox(self.X, self.Y, self.W, self.H)
            text_box.name = "ia:list-01"
            for index, value in enumerate(self.ITEMS):
                paragraph = text_box.text_frame.paragraphs[0] if index == 0 else text_box.text_frame.add_paragraph()
                paragraph.text = value
                self._set_native_bullet(paragraph, bullet="▪")
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._list_contract())

        self.assertFalse(result["valid"])
        self.assertIn("NATIVE_LIST_STRUCTURE_MISMATCH", result["errors"])

    def test_native_list_indent_must_match_spec(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "wrong-indent.pptx"
            presentation = self._presentation()
            text_box = presentation.slides[0].shapes.add_textbox(self.X, self.Y, self.W, self.H)
            text_box.name = "ia:list-01"
            for index, value in enumerate(self.ITEMS):
                paragraph = text_box.text_frame.paragraphs[0] if index == 0 else text_box.text_frame.add_paragraph()
                paragraph.text = value
                self._set_native_bullet(paragraph, indent=-100_000)
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._list_contract())

        self.assertFalse(result["valid"])
        self.assertIn("NATIVE_LIST_INDENT_MISMATCH", result["errors"])

    def test_fake_list_remains_generic_valid_without_spec(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "generic-fake-list.pptx"
            presentation = self._presentation()
            self._add_fake_list(presentation.slides[0])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1)

        self.assertTrue(result["valid"], result)
        self.assertEqual(result["native_list_paragraphs"], 0)

    def test_cli_spec_rejects_fake_list(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            directory = Path(temp_dir)
            output = directory / "fake-list.pptx"
            spec_path = directory / "page-reconstruction.json"
            presentation = self._presentation()
            self._add_fake_list(presentation.slides[0])
            presentation.save(output)
            spec_path.write_text(json.dumps(self._list_contract(), ensure_ascii=False), encoding="utf-8")
            completed = subprocess.run(
                [sys.executable, str(VALIDATOR_PATH), str(output), "--spec", str(spec_path)],
                check=False,
                capture_output=True,
                text=True,
            )

        self.assertEqual(completed.returncode, 2, completed.stdout)
        self.assertIn("NATIVE_LIST_TEXTBOX_MISSING", json.loads(completed.stdout)["errors"])


class ValidatePptxTypographyRunContractTest(unittest.TestCase):
    X = 900_000
    Y = 900_000
    W = 3_500_000
    H = 700_000
    TEXT = "重点正文"

    @staticmethod
    def _set_run_font(run, font_name: str) -> None:
        run.font.name = font_name
        east_asian = OxmlElement("a:ea")
        east_asian.set("typeface", font_name)
        run._r.get_or_add_rPr().append(east_asian)

    def _contract(self, runs: list[dict], selected_font: str | None = None) -> dict:
        return {
            "elements": [
                {
                    "element_id": "title-01",
                    "kind": "text",
                    "content": {"text": self.TEXT},
                }
            ],
            "modules": {
                "typography": {
                    "items": [
                        {
                            "element_id": "title-01",
                            "text": self.TEXT,
                            "runs": runs,
                            **(
                                {
                                    "selected_font": selected_font,
                                    "internal_font_declaration": selected_font,
                                }
                                if selected_font
                                else {}
                            ),
                            "text_box": {
                                "x": self.X,
                                "y": self.Y,
                                "w": self.W,
                                "h": self.H,
                            },
                        }
                    ]
                }
            },
        }

    def _presentation(
        self,
        run_values: list[tuple[str, bool]],
        font_name: str | None = None,
    ) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = 12_192_000
        presentation.slide_height = 6_858_000
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(self.X, self.Y, self.W, self.H)
        text_box.name = "ia:title-01"
        paragraph = text_box.text_frame.paragraphs[0]
        for value, bold in run_values:
            run = paragraph.add_run()
            run.text = value
            run.font.bold = bold
            if font_name:
                self._set_run_font(run, font_name)
        return presentation

    def test_expected_bold_run_rejects_nonbold_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "missing-bold.pptx"
            presentation = self._presentation([(self.TEXT, False)])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(
                output,
                1,
                self._contract([{"start": 0, "end": len(self.TEXT), "font_weight": 700}]),
            )

        self.assertFalse(result["valid"])
        self.assertIn("TEXT_RUN_FONT_WEIGHT_MISMATCH", result["errors"])

    def test_partial_bold_range_must_match_spec(self) -> None:
        expected = [
            {"start": 0, "end": 2, "font_weight": 700},
            {"start": 2, "end": len(self.TEXT), "font_weight": 400},
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "wrong-bold-range.pptx"
            presentation = self._presentation([(self.TEXT[:2], False), (self.TEXT[2:], True)])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._contract(expected))

        self.assertFalse(result["valid"])
        self.assertIn("TEXT_RUN_FONT_WEIGHT_MISMATCH", result["errors"])

    def test_matching_partial_bold_range_passes(self) -> None:
        expected = [
            {"start": 0, "end": 2, "font_weight": 700},
            {"start": 2, "end": len(self.TEXT), "font_weight": 400},
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "matching-bold-range.pptx"
            presentation = self._presentation([(self.TEXT[:2], True), (self.TEXT[2:], False)])
            presentation.save(output)
            result = VALIDATOR.validate_pptx(output, 1, self._contract(expected))

        self.assertTrue(result["valid"], result)

    def test_selected_font_must_match_pptx_internal_declaration(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "wrong-font.pptx"
            presentation = self._presentation(
                [(self.TEXT, True)],
                font_name="Arial",
            )
            presentation.save(output)
            result = VALIDATOR.validate_pptx(
                output,
                1,
                self._contract(
                    [{"start": 0, "end": len(self.TEXT), "font_weight": 700}],
                    selected_font="Noto Sans CJK SC",
                ),
            )

        self.assertFalse(result["valid"])
        self.assertIn("TEXT_RUN_FONT_DECLARATION_MISMATCH", result["errors"])

    def test_every_text_run_must_use_selected_font(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "mixed-fonts.pptx"
            presentation = self._presentation(
                [(self.TEXT[:2], True), (self.TEXT[2:], True)]
            )
            runs = presentation.slides[0].shapes[0].text_frame.paragraphs[0].runs
            self._set_run_font(runs[0], "Noto Sans CJK SC")
            self._set_run_font(runs[1], "Arial")
            presentation.save(output)
            result = VALIDATOR.validate_pptx(
                output,
                1,
                self._contract(
                    [
                        {"start": 0, "end": 2, "font_weight": 700},
                        {"start": 2, "end": len(self.TEXT), "font_weight": 700},
                    ],
                    selected_font="Noto Sans CJK SC",
                ),
            )

        self.assertFalse(result["valid"])
        self.assertIn("TEXT_RUN_FONT_DECLARATION_MISMATCH", result["errors"])


class ValidatePptxRoundedRectangleTest(unittest.TestCase):
    def _validate_shape(self, shape_type: MSO_SHAPE, adjustment: float | None = None) -> dict:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "shape.pptx"
            presentation = Presentation()
            presentation.slide_width = 12_192_000
            presentation.slide_height = 6_858_000
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            shape = slide.shapes.add_shape(shape_type, 900_000, 900_000, 2_400_000, 800_000)
            if adjustment is not None:
                shape.adjustments[0] = adjustment
            presentation.save(output)
            return VALIDATOR.validate_pptx(output, expected_slides=1)

    def test_default_round_rect_without_adjustment_is_rejected(self) -> None:
        result = self._validate_shape(MSO_SHAPE.ROUNDED_RECTANGLE)

        self.assertFalse(result["valid"])
        self.assertIn("ROUND_RECT_ADJUSTMENT_MISSING", result["errors"])

    def test_round_rect_with_explicit_adjustment_passes(self) -> None:
        result = self._validate_shape(MSO_SHAPE.ROUNDED_RECTANGLE, adjustment=0.25)

        self.assertTrue(result["valid"], result)

    def test_round_rect_with_out_of_range_adjustment_is_rejected(self) -> None:
        result = self._validate_shape(MSO_SHAPE.ROUNDED_RECTANGLE, adjustment=0.75)

        self.assertFalse(result["valid"])
        self.assertIn("ROUND_RECT_ADJUSTMENT_INVALID", result["errors"])

    def test_plain_rectangle_does_not_require_adjustment(self) -> None:
        result = self._validate_shape(MSO_SHAPE.RECTANGLE)

        self.assertTrue(result["valid"], result)


class ValidatePptxElementBindingTest(unittest.TestCase):
    X = 900_000
    Y = 900_000
    W = 3_500_000
    H = 700_000

    def _presentation(self) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = 12_192_000
        presentation.slide_height = 6_858_000
        presentation.slides.add_slide(presentation.slide_layouts[6])
        return presentation

    def _text_spec(self, expected_text: str = "可编辑标题") -> dict:
        return {
            "elements": [
                {
                    "element_id": "title",
                    "kind": "text",
                    "slide_bbox": [self.X, self.Y, self.W, self.H],
                    "editable": True,
                    "content": {"text": expected_text},
                }
            ]
        }

    def test_spec_rejects_missing_named_element_object(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "unnamed.pptx"
            presentation = self._presentation()
            text_box = presentation.slides[0].shapes.add_textbox(
                self.X, self.Y, self.W, self.H
            )
            text_box.text = "可编辑标题"
            presentation.save(output)

            result = VALIDATOR.validate_pptx(output, 1, self._text_spec())

        self.assertFalse(result["valid"])
        self.assertIn("ELEMENT_OBJECT_MISSING", result["errors"])

    def test_spec_rejects_wrong_named_text(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "wrong-text.pptx"
            presentation = self._presentation()
            text_box = presentation.slides[0].shapes.add_textbox(
                self.X, self.Y, self.W, self.H
            )
            text_box.name = "ia:title"
            text_box.text = "错误标题"
            presentation.save(output)

            result = VALIDATOR.validate_pptx(output, 1, self._text_spec())

        self.assertFalse(result["valid"])
        self.assertIn("ELEMENT_TEXT_MISMATCH", result["errors"])

    def test_spec_rejects_icon_with_wrong_embedded_media_hash(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            expected_asset = root / "orange.png"
            embedded_asset = root / "blue.png"
            Image.new("RGB", (32, 32), "#F15A22").save(expected_asset)
            Image.new("RGB", (32, 32), "#0066CC").save(embedded_asset)
            output = root / "wrong-icon.pptx"
            presentation = self._presentation()
            slide = presentation.slides[0]
            text_box = slide.shapes.add_textbox(self.X, self.Y, self.W, self.H)
            text_box.name = "ia:title"
            text_box.text = "可编辑标题"
            picture = slide.shapes.add_picture(
                str(embedded_asset), 5_000_000, 1_000_000, 300_000, 300_000
            )
            picture.name = "ia:status-icon"
            presentation.save(output)
            spec = self._text_spec()
            spec["elements"].append(
                {
                    "element_id": "status-icon",
                    "kind": "icon",
                    "slide_bbox": [5_000_000, 1_000_000, 300_000, 300_000],
                    "editable": False,
                    "content": {},
                }
            )
            spec["modules"] = {
                "icons": {
                    "icons": [
                        {
                            "element_id": "status-icon",
                            "asset_sha256": hashlib.sha256(expected_asset.read_bytes()).hexdigest(),
                        }
                    ]
                }
            }

            result = VALIDATOR.validate_pptx(output, 1, spec)

        self.assertFalse(result["valid"])
        self.assertIn("ELEMENT_MEDIA_HASH_MISMATCH", result["errors"])

    def test_offslide_empty_shape_does_not_bypass_full_slide_picture_gate(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image = root / "background.png"
            Image.new("RGB", (1600, 900), "#27506B").save(image)
            output = root / "offslide-bypass.pptx"
            presentation = self._presentation()
            slide = presentation.slides[0]
            slide.shapes.add_picture(
                str(image), 0, 0, presentation.slide_width, presentation.slide_height
            )
            slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                presentation.slide_width + 100_000,
                100_000,
                200_000,
                200_000,
            )
            presentation.save(output)

            result = VALIDATOR.validate_pptx(output, 1)

        self.assertFalse(result["valid"])
        self.assertIn("FULL_SLIDE_PICTURE_ONLY", result["errors"])

    def test_named_multipart_shape_element_passes(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "multipart-shape.pptx"
            presentation = self._presentation()
            slide = presentation.slides[0]
            for part in ("fill", "border"):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, self.X, self.Y, self.W, self.H
                )
                shape.name = f"ia:card:{part}"
            presentation.save(output)
            spec = {
                "elements": [
                    {
                        "element_id": "card",
                        "kind": "shape",
                        "slide_bbox": [self.X, self.Y, self.W, self.H],
                        "editable": True,
                        "content": {},
                    }
                ]
            }

            result = VALIDATOR.validate_pptx(output, 1, spec)

        self.assertTrue(result["valid"], result)


class ValidatePptxCliOutputTest(unittest.TestCase):
    def test_cli_output_matches_stdout_json(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            pptx = root / "single-slide.pptx"
            report = root / "reports" / "structure.json"
            presentation = Presentation()
            presentation.slide_width = 12_192_000
            presentation.slide_height = 6_858_000
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(
                914_400, 914_400, 4_572_000, 914_400
            )
            text_box.text = "可编辑标题"
            presentation.save(pptx)

            completed = subprocess.run(
                [
                    sys.executable,
                    str(VALIDATOR_PATH),
                    str(pptx),
                    "--expected-slides",
                    "1",
                    "--summary",
                    "--output",
                    str(report),
                ],
                check=False,
                capture_output=True,
                text=True,
            )

            self.assertEqual(0, completed.returncode, completed.stderr)
            self.assertTrue(report.is_file())
            self.assertEqual(json.loads(completed.stdout), json.loads(report.read_text(encoding="utf-8")))


if __name__ == "__main__":
    unittest.main()
