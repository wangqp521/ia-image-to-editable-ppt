"""Behavior contracts for the schema v2 single-page compiler."""

from __future__ import annotations

import copy
import hashlib
import importlib.util
import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from unittest import mock
from xml.etree import ElementTree as ET
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from PIL import Image

from tests.fixture_specs import make_minimal_spec, write_valid_fixture


SCRIPTS_ROOT = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS_ROOT) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_ROOT))

from lib.error_codes import ToolError
from lib.capabilities import capability_manifest_sha256
from lib.hashing import canonical_json_sha256, file_sha256
from lib.representation_contracts import validate_representation_plan
from lib.schema_contracts import CANONICAL_VALUES
from lib.spec_identity import content_spec_sha256


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _load_script(filename: str, module_name: str):
    path = SCRIPTS_ROOT / filename
    module_spec = importlib.util.spec_from_file_location(module_name, path)
    if module_spec is None or module_spec.loader is None:
        raise AssertionError(f"cannot load {path}")
    module = importlib.util.module_from_spec(module_spec)
    module_spec.loader.exec_module(module)
    return module


def _compile_single_page():
    try:
        module = _load_script(
            "build_pptx_from_spec.py", "test_build_pptx_from_spec_script"
        )
    except (FileNotFoundError, ModuleNotFoundError) as exc:
        raise AssertionError("schema compiler is not implemented") from exc
    try:
        return module.compile_single_page
    except AttributeError as exc:
        raise AssertionError("schema compiler is not implemented") from exc


def compile_fixture(root: Path, spec: dict[str, Any]) -> tuple[Path, dict[str, Any]]:
    """Compile one passing fixture through the real public transaction."""
    root.mkdir(parents=True, exist_ok=True)
    spec_path, prebuild_path = write_valid_fixture(root, spec)
    output = root / "page.pptx"
    report_path = root / "build-report.json"
    report = _compile_single_page()(spec_path, prebuild_path, output, report_path)
    return output, report


def compile_with_claimed_prebuild(
    root: Path, spec: dict[str, Any]
) -> tuple[Path, dict[str, Any]]:
    """Reach the compiler's secondary gate with a stale claimed-pass report."""
    root.mkdir(parents=True, exist_ok=True)
    spec_path = root / "page-reconstruction.json"
    prebuild_path = root / "prebuild-report.json"
    output = root / "page.pptx"
    report_path = root / "build-report.json"
    spec_path.write_text(
        json.dumps(spec, ensure_ascii=False, indent=2, sort_keys=True),
        encoding="utf-8",
    )
    prebuild_path.write_text(
        json.dumps(
            {
                "valid": True,
                "stage": "prebuild",
                "errors": [],
                "spec_sha256": canonical_json_sha256(spec),
            },
            ensure_ascii=False,
            indent=2,
            sort_keys=True,
        ),
        encoding="utf-8",
    )
    report = _compile_single_page()(spec_path, prebuild_path, output, report_path)
    return output, report


def compile_structured_fixture(
    root: Path, spec: dict[str, Any]
) -> tuple[Path, dict[str, Any]]:
    """Keep missing structured renderers in the assertion-failure RED lane."""
    unsupported_path: str | None = None
    try:
        return compile_fixture(root, spec)
    except ToolError as exc:
        if exc.code == "UNSUPPORTED_KIND":
            unsupported_path = exc.path
        else:
            raise
    raise AssertionError(
        f"structured renderer is not implemented: {unsupported_path}"
    )


def make_native_list_spec(root: Path) -> dict[str, Any]:
    """Return a passing two-item native-list variant of the shared fixture."""
    spec = make_minimal_spec(root)
    element = spec["elements"][0]
    contract = spec["modules"]["typography"]["items"][0]
    text = "第一项第二项"
    element["content"]["text"] = text
    contract["text"] = text
    contract["runs"] = [
        {
            "start": 0,
            "end": len(text),
            "font_size": 20,
            "font_weight": 400,
            "color": "#203040",
            "letter_spacing": 0,
            "italic": False,
            "underline": False,
            "strike": False,
            "baseline": 0,
        }
    ]
    list_style = {
        "is_list": True,
        "level": 0,
        "bullet_type": "char",
        "bullet": "•",
        "bullet_font": "follow_text",
        "bullet_size_mode": "follow_text",
        "bullet_size_value": None,
        "bullet_color": "follow_text",
    }
    contract["paragraphs"] = [
        {
            "start": 0,
            "end": 3,
            "alignment": "left",
            "line_spacing": 1.0,
            "space_before": 0,
            "space_after": 0,
            "margin_left": 320040,
            "indent": -228600,
            "list": dict(list_style),
        },
        {
            "start": 3,
            "end": len(text),
            "alignment": "left",
            "line_spacing": 1.0,
            "space_before": 0,
            "space_after": 0,
            "margin_left": 320040,
            "indent": -228600,
            "list": dict(list_style),
        },
    ]
    contract["text_box"]["paragraph_breaks"] = [3]
    return spec


def _append_primitive(
    spec: dict[str, Any],
    *,
    element_id: str,
    kind: str,
    style: dict[str, Any],
    content: dict[str, Any],
    source_bbox: list[int] | None = None,
    slide_bbox: list[int] | None = None,
    selected_mode: str = "native",
) -> dict[str, Any]:
    """Append one real primitive and its representation binding."""
    source_bbox = source_bbox or [80, 120, 240, 120]
    slide_bbox = slide_bbox or [609600, 914400, 1828800, 914400]
    element = {
        "element_id": element_id,
        "kind": kind,
        "source_bbox": source_bbox,
        "slide_bbox": slide_bbox,
        "layer": 2,
        "editable": kind not in {"picture", "icon"},
        "confidence": "high",
        "style": style,
        "content": content,
    }
    spec["elements"].append(element)
    spec["regions"][0]["element_ids"].append(element_id)
    spec["reading_order"].append(element_id)
    evidence = Path(spec["modules"]["representation_plan"]["items"][0]["evidence"][0])
    spec["modules"]["representation_plan"]["items"].append(
        {
            "source_fact_id": f"fact-{element_id}",
            "semantic_role": kind,
            "source_bbox": list(source_bbox),
            "required": True,
            "selected_mode": selected_mode,
            "required_editability": "none" if selected_mode == "asset" else "full",
            "fallback_policy": "allow_minimal_asset" if selected_mode == "asset" else "forbid",
            "bound_element_ids": [element_id],
            "reason": "primitive compiler fixture",
            "coverage_status": "covered",
            "evidence": [str(evidence)],
        }
    )
    return element


def make_shape_spec(root: Path, shape_type: str) -> dict[str, Any]:
    spec = make_minimal_spec(root)
    style: dict[str, Any] = {
        "shape_type": shape_type,
        "fill": "noFill",
        "line": {"color": "#264653", "width": 12700, "dash": "solid", "opacity": 1.0},
        "effects": "none",
        "rotation": 0,
    }
    if shape_type == "roundRect":
        style["adjustments"] = [0.25]
    _append_primitive(
        spec,
        element_id="card",
        kind="shape",
        style=style,
        content={},
    )
    return spec


def make_picture_spec(root: Path, *, mode: str = "none") -> dict[str, Any]:
    spec = make_minimal_spec(root)
    asset = root / "photo.png"
    Image.new("RGB", (40, 20), (32, 96, 192)).save(asset)
    _append_primitive(
        spec,
        element_id="photo",
        kind="picture",
        style={"rotation": 15, "opacity": 0.6},
        content={
            "asset": {
                "path": str(asset.resolve()),
                "asset_sha256": hashlib.sha256(asset.read_bytes()).hexdigest(),
                "pixel_size": [40, 20],
            },
            "mode": mode,
            "crop": {"left": 0.1, "top": 0.0, "right": 0.1, "bottom": 0.0},
        },
        source_bbox=[80, 120, 240, 240],
        slide_bbox=[609600, 914400, 1828800, 1828800],
        selected_mode="asset",
    )
    return spec


def make_icon_spec(root: Path) -> dict[str, Any]:
    spec = make_minimal_spec(root)
    source_bbox = [30, 30, 16, 16]
    slide_bbox = [228600, 228600, 121920, 121920]
    icons_dir = root / "assets" / "icons"
    icons_dir.mkdir(parents=True, exist_ok=True)
    asset = icons_dir / "status.png"
    icon = Image.new("RGBA", (20, 20), (255, 255, 255, 0))
    icon.putalpha(Image.new("L", (20, 20), 0))
    alpha = icon.getchannel("A")
    for y in range(2, 18):
        for x in range(2, 18):
            alpha.putpixel((x, y), 255)
    icon.putalpha(alpha)
    icon.save(asset)
    asset_sha256 = hashlib.sha256(asset.read_bytes()).hexdigest()
    alpha_sha256 = hashlib.sha256(alpha.tobytes()).hexdigest()
    _append_primitive(
        spec,
        element_id="status-icon",
        kind="icon",
        style={"rotation": 0, "opacity": 1.0},
        content={
            "asset": {
                "path": str(asset.resolve()),
                "asset_sha256": asset_sha256,
                "pixel_size": [20, 20],
            },
            "mode": "contain",
            "crop": {"left": 0.0, "top": 0.0, "right": 0.0, "bottom": 0.0},
        },
        source_bbox=source_bbox,
        slide_bbox=slide_bbox,
        selected_mode="asset",
    )
    reference = spec["clean_visual_reference"]
    spec["activated_modules"].append("icons")
    spec["modules"]["icons"] = {
        "schema_version": 2,
        "page_id": spec["page_id"],
        "slide_coordinate_unit": "EMU",
        "clean_visual_reference": reference["path"],
        "clean_visual_sha256": reference["sha256"],
        "icons": [
            {
                "icon_id": "status-icon",
                "element_id": "status-icon",
                "category": "status",
                "instance_count": 1,
                "repeat_group": None,
                "semantic_scope": "icon_only",
                "source_bbox": source_bbox,
                "slide_bbox": slide_bbox,
                "layer": 2,
                "source_path": reference["path"],
                "source_sha256": reference["sha256"],
                "crop_mode": "alpha_isolation",
                "padding": 2,
                "background_handling": "alpha_only",
                "asset_path": str(asset.resolve()),
                "asset_sha256": asset_sha256,
                "alpha_mask_sha256": alpha_sha256,
                "final_width": 20,
                "final_height": 20,
                "sharpness": "preserved",
                "validation": "passed",
                "native_redraw": False,
                "selectable_picture_verified": False,
                "object_type": "picture",
            }
        ],
    }
    return spec


def make_merged_table_spec(root: Path) -> dict[str, Any]:
    spec = make_minimal_spec(root)
    _append_primitive(
        spec,
        element_id="table",
        kind="table",
        style={"rotation": 0},
        content={
            "rows": [320040, 594360],
            "columns": [609600, 1219200],
            "cells": [
                {
                    "row": 0,
                    "column": 0,
                    "row_span": 1,
                    "column_span": 2,
                    "text": "合并标题",
                    "fill": "#DDEEFF",
                    "margins": {
                        "left": 101,
                        "right": 102,
                        "top": 103,
                        "bottom": 104,
                    },
                    "alignment": "center",
                    "vertical_alignment": "middle",
                    "font": {
                        "name": "Arial",
                        "size": 14,
                        "weight": 700,
                        "color": "#112233",
                        "italic": True,
                    },
                    "borders": {
                        "bottom": {"color": "#445566", "width": 12700}
                    },
                },
                {
                    "row": 1,
                    "column": 0,
                    "row_span": 1,
                    "column_span": 1,
                    "text": "A",
                    "fill": "#FFFFFF",
                    "margins": {"left": 0, "right": 0, "top": 0, "bottom": 0},
                    "alignment": "left",
                    "vertical_alignment": "top",
                    "font": {
                        "name": "Arial",
                        "size": 12,
                        "weight": 400,
                        "color": "#000000",
                        "italic": False,
                    },
                    "borders": {},
                },
                {
                    "row": 1,
                    "column": 1,
                    "row_span": 1,
                    "column_span": 1,
                    "text": "B",
                    "fill": "noFill",
                    "margins": {"left": 1, "right": 2, "top": 3, "bottom": 4},
                    "alignment": "right",
                    "vertical_alignment": "bottom",
                    "font": {
                        "name": "Arial",
                        "size": 11,
                        "weight": 400,
                        "color": "#334455",
                        "italic": False,
                    },
                    "borders": {
                        "left": {"color": "#778899", "width": 25400}
                    },
                },
            ],
        },
    )
    return spec


def _status_part(part_id: str, bbox: list[int], text: str, fill: str) -> dict[str, Any]:
    return {
        "part_id": part_id,
        "slide_bbox": bbox,
        "style": {
            "shape_type": "roundRect",
            "adjustments": [0.2],
            "fill": fill,
            "line": {
                "color": "#245566",
                "width": 12700,
                "dash": "solid",
                "opacity": 1.0,
            },
            "effects": "none",
            "rotation": 12.5,
            "text_style": {
                "font_name": "Arial",
                "font_size": 12,
                "font_weight": 600,
                "color": "#FFFFFF",
                "italic": False,
                "alignment": "center",
                "vertical_alignment": "middle",
                "margins": {
                    "left": 101,
                    "right": 102,
                    "top": 103,
                    "bottom": 104,
                },
                "wrap": True,
            },
        },
        "content": {"text": text},
    }


def make_status_spec(root: Path, *, repeat: bool) -> dict[str, Any]:
    spec = make_minimal_spec(root)
    left = _status_part("segment-0", [609600, 914400, 914400, 914400], "已完成", "#22AA66")
    right = _status_part("segment-1", [1524000, 914400, 914400, 914400], "待处理", "#8899AA")
    if repeat:
        common_style = copy.deepcopy(left["style"])
        del common_style["fill"]
        content: dict[str, Any] = {
            "part_defaults": {"style": common_style},
            "repeat_sequence": [
                {
                    "part_id": left["part_id"],
                    "slide_bbox": left["slide_bbox"],
                    "style": {"fill": left["style"]["fill"]},
                    "content": left["content"],
                },
                {
                    "part_id": right["part_id"],
                    "slide_bbox": right["slide_bbox"],
                    "style": {"fill": right["style"]["fill"]},
                    "content": right["content"],
                },
            ],
        }
    else:
        content = {"part_defaults": {}, "parts": [left, right]}
    _append_primitive(
        spec,
        element_id="status",
        kind="status",
        style={"rotation": 0},
        content=content,
        slide_bbox=[609600, 914400, 1828800, 914400],
        selected_mode="composite",
    )
    return spec


class BuildTextTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self.temporary.cleanup)
        self.root = Path(self.temporary.name)
        self.output = self.root / "page.pptx"
        self.report = self.root / "build-report.json"

    def test_compiler_builds_native_background_without_asset_fallback(self) -> None:
        spec = make_minimal_spec(self.root)

        pptx, report = compile_fixture(self.root / "compile", spec)

        self.assertTrue(pptx.is_file())
        self.assertEqual(
            {"native": 1, "background_picture": 0},
            report["background_summary"],
        )
        self.assertEqual([], report["background_pictures"])
        self.assertNotIn(
            "background-base",
            {item["source_fact_id"] for item in report["asset_fallbacks"]},
        )

    def test_passing_authoring_report_is_not_a_compiler_credential(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root,
            make_minimal_spec(self.root),
        )
        authoring_report = json.loads(prebuild_path.read_text(encoding="utf-8"))
        authoring_report["stage"] = "authoring"
        prebuild_path.write_text(
            json.dumps(authoring_report, ensure_ascii=False),
            encoding="utf-8",
        )

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(
                spec_path,
                prebuild_path,
                self.output,
                self.report,
            )

        self.assertEqual("BUILD_OUTPUT_INCOMPLETE", raised.exception.code)
        self.assertEqual(str(prebuild_path.resolve()), raised.exception.path)
        self.assertEqual(
            "compiler requires a passing prebuild report",
            raised.exception.detail,
        )
        self.assertFalse(self.output.exists())
        self.assertFalse(self.report.exists())

    def test_text_spec_builds_named_editable_textbox(self) -> None:
        spec = make_minimal_spec(self.root)
        element_id = spec["elements"][0]["element_id"]
        spec_path, prebuild_path = write_valid_fixture(self.root, spec)

        report = _compile_single_page()(
            spec_path, prebuild_path, self.output, self.report
        )

        presentation = Presentation(self.output)
        names = {shape.name for shape in presentation.slides[0].shapes}
        self.assertIn(f"ia:{element_id}", names)
        self.assertEqual(report["elements"][element_id]["object_type"], "sp")

    def test_text_contract_writes_explicit_frame_and_run_properties(self) -> None:
        spec = make_minimal_spec(self.root)
        contract = spec["modules"]["typography"]["items"][0]
        contract["runs"][0].update(
            {
                "italic": True,
                "underline": True,
                "strike": True,
                "baseline": 25000,
                "letter_spacing": 1.5,
            }
        )
        contract["text_box"]["margins"] = {
            "left": 101,
            "right": 102,
            "top": 103,
            "bottom": 104,
        }
        contract["text_box"]["vertical_alignment"] = "middle"
        contract["text_box"]["wrap"] = True

        output, _ = compile_fixture(self.root, spec)

        with zipfile.ZipFile(output) as archive:
            slide = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
        text_shape = next(
            shape
            for shape in slide.findall(".//p:sp", NS)
            if shape.find("p:nvSpPr/p:cNvPr", NS).get("name") == "ia:element-001"
        )
        body_pr = text_shape.find("p:txBody/a:bodyPr", NS)
        self.assertIsNotNone(body_pr)
        assert body_pr is not None
        self.assertIsNotNone(body_pr.find("a:noAutofit", NS))
        self.assertEqual(body_pr.get("lIns"), "101")
        self.assertEqual(body_pr.get("rIns"), "102")
        self.assertEqual(body_pr.get("tIns"), "103")
        self.assertEqual(body_pr.get("bIns"), "104")
        self.assertEqual(body_pr.get("anchor"), "ctr")
        self.assertEqual(body_pr.get("wrap"), "square")
        run_properties = text_shape.find("p:txBody/a:p/a:r/a:rPr", NS)
        self.assertIsNotNone(run_properties)
        assert run_properties is not None
        self.assertEqual(run_properties.get("b"), "1")
        self.assertEqual(run_properties.get("i"), "1")
        self.assertEqual(run_properties.get("u"), "sng")
        self.assertEqual(run_properties.get("strike"), "sngStrike")
        self.assertEqual(run_properties.get("baseline"), "25000")
        self.assertEqual(run_properties.get("spc"), "150")
        east_asian_font = run_properties.find("a:ea", NS)
        self.assertIsNotNone(east_asian_font)
        assert east_asian_font is not None
        self.assertEqual(
            east_asian_font.get("typeface"),
            contract["selected_font"],
        )

    def test_stale_prebuild_report_publishes_nothing(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        payload = json.loads(prebuild_path.read_text(encoding="utf-8"))
        payload["spec_sha256"] = "0" * 64
        prebuild_path.write_text(json.dumps(payload), encoding="utf-8")

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(
                spec_path, prebuild_path, self.output, self.report
            )

        self.assertEqual(raised.exception.code, "SPEC_HASH_MISMATCH")
        self.assertFalse(self.output.exists())
        self.assertFalse(self.report.exists())

    def test_native_list_is_normalized_before_report_hash(self) -> None:
        spec = make_native_list_spec(self.root)
        spec_path, prebuild_path = write_valid_fixture(self.root, spec)

        report = _compile_single_page()(
            spec_path, prebuild_path, self.output, self.report
        )

        validator = _load_script("validate_pptx.py", "test_compiler_validate_pptx")
        normalized = validator.validate_pptx(self.output, 1, spec)
        self.assertTrue(normalized["valid"], normalized)
        self.assertEqual(report["pptx_sha256"], file_sha256(self.output))
        self.assertEqual(
            report["capability_manifest_sha256"], capability_manifest_sha256()
        )
        self.assertTrue(report["normalization"]["applied"])
        self.assertEqual(report["normalization"]["valid"], True)

    def test_native_list_normalizer_consumes_frozen_spec_snapshot(self) -> None:
        spec = make_native_list_spec(self.root)
        spec_path, prebuild_path = write_valid_fixture(self.root, spec)
        original_spec_sha256 = canonical_json_sha256(spec)
        original_content_sha256 = content_spec_sha256(spec)
        compiler = _load_script(
            "build_pptx_from_spec.py", "test_frozen_native_list_compiler"
        )
        real_normalize = compiler.normalize_pptx
        observed_snapshot_paths: list[Path] = []

        def replace_original_before_normalization(
            input_pptx: Path,
            snapshot_path: Path,
            output_pptx: Path,
            normalization_report: Path,
        ) -> dict[str, Any]:
            mutated = json.loads(spec_path.read_text(encoding="utf-8"))
            for paragraph in mutated["modules"]["typography"]["items"][0][
                "paragraphs"
            ]:
                paragraph["list"]["bullet_color"] = "#FF0000"
            spec_path.write_text(
                json.dumps(mutated, ensure_ascii=False, indent=2, sort_keys=True),
                encoding="utf-8",
            )
            observed_snapshot_paths.append(Path(snapshot_path))
            return real_normalize(
                input_pptx,
                snapshot_path,
                output_pptx,
                normalization_report,
            )

        with mock.patch.object(
            compiler,
            "normalize_pptx",
            side_effect=replace_original_before_normalization,
        ):
            report = compiler.compile_single_page(
                spec_path,
                prebuild_path,
                self.output,
                self.report,
            )

        self.assertEqual(len(observed_snapshot_paths), 1)
        self.assertNotEqual(observed_snapshot_paths[0], spec_path)
        with zipfile.ZipFile(self.output) as archive:
            slide_xml = archive.read("ppt/slides/slide1.xml")
        self.assertIn(b"buClrTx", slide_xml)
        self.assertNotIn(b"FF0000", slide_xml)
        self.assertEqual(report["schema_sha256"], original_spec_sha256)
        self.assertEqual(report["content_spec_sha256"], original_content_sha256)
        self.assertEqual(report["input_spec_sha256"], original_spec_sha256)

    def test_compiler_and_prebuild_report_identical_background_binding_errors(
        self,
    ) -> None:
        validator = _load_script(
            "validate_reconstruction_spec.py", "test_background_error_parity"
        )
        cases: list[tuple[str, dict[str, Any]]] = []

        overlap = make_minimal_spec(self.root / "overlap")
        overlap_fact = copy.deepcopy(
            overlap["modules"]["representation_plan"]["items"][0]
        )
        overlap_fact.update(
            {
                "source_fact_id": "fact-background-overlap",
                "semantic_role": "shape",
                "source_bbox": [0, 0, 1600, 900],
                "bound_element_ids": ["background-base"],
            }
        )
        overlap["modules"]["representation_plan"]["items"].append(overlap_fact)
        cases.append(("overlap", overlap))

        duplicate = make_minimal_spec(self.root / "duplicate")
        duplicate_item = copy.deepcopy(
            duplicate["modules"]["background"]["items"][0]
        )
        duplicate_item.update(
            {"background_id": "background-002", "role": "texture"}
        )
        duplicate["modules"]["background"]["items"].append(duplicate_item)
        cases.append(("duplicate", duplicate))

        for name, candidate in cases:
            with self.subTest(name=name):
                prebuild = validator.validate_spec(candidate, stage="prebuild")
                prebuild_issue = next(
                    item
                    for item in prebuild["errors"]
                    if item["code"] == "BACKGROUND_BINDING_CONFLICT"
                )
                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(
                        self.root / f"compile-{name}", candidate
                    )

                self.assertEqual(
                    (
                        raised.exception.code,
                        raised.exception.path,
                        raised.exception.detail,
                    ),
                    (
                        prebuild_issue["code"],
                        prebuild_issue["path"],
                        prebuild_issue["detail"],
                    ),
                )

    def test_invalid_table_contract_publishes_nothing(self) -> None:
        spec = make_minimal_spec(self.root)
        element = spec["elements"][0]
        element["kind"] = "table"
        element["style"] = {}
        element["content"] = {"rows": [], "columns": [], "cells": []}
        spec["modules"]["representation_plan"]["items"][0][
            "semantic_role"
        ] = "table"
        spec_path = self.root / "page-reconstruction.json"
        prebuild_path = self.root / "prebuild-report.json"
        spec_path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        prebuild_path.write_text(
            json.dumps(
                {
                    "valid": True,
                    "stage": "prebuild",
                    "errors": [],
                    "spec_sha256": canonical_json_sha256(spec),
                }
            ),
            encoding="utf-8",
        )

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(
                spec_path, prebuild_path, self.output, self.report
            )

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertFalse(self.output.exists())
        self.assertFalse(self.report.exists())

    def test_prebuild_and_compiler_reject_unbound_shape_with_same_issue(self) -> None:
        spec = make_shape_spec(self.root, "rectangle")
        spec["modules"]["representation_plan"]["items"].pop()
        expected_issue = {
            "code": "REPRESENTATION_INCOMPLETE",
            "path": "modules.representation_plan.items",
            "detail": "every element must have exactly one selected mode",
        }

        with self.assertRaises(ToolError) as raised:
            compile_with_claimed_prebuild(self.root, spec)
        compiler_issue = {
            "code": raised.exception.code,
            "path": raised.exception.path,
            "detail": raised.exception.detail,
        }
        validator = _load_script(
            "validate_reconstruction_spec.py", "test_mode_coverage_validator"
        )
        prebuild = validator.validate_spec(spec, stage="prebuild")

        self.assertEqual(compiler_issue, expected_issue)
        self.assertFalse(prebuild["valid"], prebuild)
        self.assertEqual(prebuild["errors"], [compiler_issue])

    def test_equal_layers_follow_schema_reading_order(self) -> None:
        spec = make_minimal_spec(self.root)
        first = spec["elements"][0]
        first["layer"] = 1
        second = copy.deepcopy(first)
        second.update(
            {
                "element_id": "second",
                "source_bbox": [30, 120, 800, 60],
                "slide_bbox": [228600, 914400, 6096000, 457200],
                "content": {"text": "第二项"},
            }
        )
        spec["elements"].append(second)
        spec["regions"][0]["element_ids"].append("second")
        spec["reading_order"] = [
            "background-base",
            "second",
            first["element_id"],
        ]
        typography = copy.deepcopy(spec["modules"]["typography"]["items"][0])
        typography.update({"element_id": "second", "text": "第二项"})
        typography["runs"][0].update({"start": 0, "end": 3})
        typography["paragraphs"][0].update({"start": 0, "end": 3})
        typography["text_box"].update(
            {"x": 228600, "y": 914400, "w": 6096000, "h": 457200}
        )
        spec["modules"]["typography"]["items"].append(typography)
        plan = copy.deepcopy(spec["modules"]["representation_plan"]["items"][0])
        plan.update(
            {
                "source_fact_id": "fact-second",
                "source_bbox": [30, 120, 800, 60],
                "bound_element_ids": ["second"],
            }
        )
        spec["modules"]["representation_plan"]["items"].append(plan)

        output, _ = compile_fixture(self.root, spec)

        names = [shape.name for shape in Presentation(output).slides[0].shapes]
        self.assertEqual(
            names,
            ["ia:background-base", "ia:second", f"ia:{first['element_id']}"],
        )

    def test_cli_publishes_the_same_build_report(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        command = [
            sys.executable,
            str(SCRIPTS_ROOT / "build_pptx_from_spec.py"),
            "--spec",
            str(spec_path),
            "--prebuild-report",
            str(prebuild_path),
            "--output",
            str(self.output),
            "--build-report",
            str(self.report),
        ]

        completed = subprocess.run(command, check=False, capture_output=True, text=True)

        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertTrue(self.output.is_file())
        self.assertTrue(self.report.is_file())
        self.assertEqual(
            json.loads(completed.stdout),
            json.loads(self.report.read_text(encoding="utf-8")),
        )

    def test_api_rejects_dangling_output_symlink_before_resolving_it(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        symlink_target = self.root / "symlink-target.pptx"
        self.output.symlink_to(symlink_target)

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(
                spec_path, prebuild_path, self.output, self.report
            )

        self.assertEqual(
            raised.exception.as_dict(),
            {
                "code": "BUILD_OUTPUT_INCOMPLETE",
                "path": str(self.output),
                "detail": "output path already exists",
            },
        )
        self.assertTrue(self.output.is_symlink())
        self.assertFalse(symlink_target.exists())
        self.assertFalse(self.report.exists())

    def test_cli_rejects_dangling_output_symlink_before_resolving_it(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        symlink_target = self.root / "symlink-target.pptx"
        self.output.symlink_to(symlink_target)

        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPTS_ROOT / "build_pptx_from_spec.py"),
                "--spec",
                str(spec_path),
                "--prebuild-report",
                str(prebuild_path),
                "--output",
                str(self.output),
                "--build-report",
                str(self.report),
            ],
            check=False,
            capture_output=True,
            text=True,
        )

        self.assertEqual(completed.returncode, 2, completed.stderr)
        self.assertEqual(completed.stderr, "")
        self.assertEqual(
            json.loads(completed.stdout),
            {
                "valid": False,
                "errors": [
                    {
                        "code": "BUILD_OUTPUT_INCOMPLETE",
                        "path": str(self.output),
                        "detail": "output path already exists",
                    }
                ],
            },
        )
        self.assertTrue(self.output.is_symlink())
        self.assertFalse(symlink_target.exists())
        self.assertFalse(self.report.exists())

    def test_api_rejects_dangling_report_symlink_before_resolving_it(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        symlink_target = self.root / "symlink-target.json"
        self.report.symlink_to(symlink_target)

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(
                spec_path, prebuild_path, self.output, self.report
            )

        self.assertEqual(
            raised.exception.as_dict(),
            {
                "code": "BUILD_OUTPUT_INCOMPLETE",
                "path": str(self.report),
                "detail": "output path already exists",
            },
        )
        self.assertTrue(self.report.is_symlink())
        self.assertFalse(symlink_target.exists())
        self.assertFalse(self.output.exists())

    def test_api_rejects_symlinked_output_parent_before_resolving_it(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        real_parent = self.root / "real-output"
        real_parent.mkdir()
        linked_parent = self.root / "linked-output"
        linked_parent.symlink_to(real_parent, target_is_directory=True)
        output = linked_parent / "page.pptx"

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(spec_path, prebuild_path, output, self.report)

        self.assertEqual(
            raised.exception.as_dict(),
            {
                "code": "BUILD_OUTPUT_INCOMPLETE",
                "path": str(linked_parent),
                "detail": "output parent must be an existing real directory",
            },
        )
        self.assertTrue(linked_parent.is_symlink())
        self.assertFalse((real_parent / "page.pptx").exists())
        self.assertFalse(self.report.exists())

    def test_api_rejects_nested_output_ancestor_symlink(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        real_parent = self.root / "real-output"
        (real_parent / "subdirectory").mkdir(parents=True)
        linked_parent = self.root / "linked-output"
        linked_parent.symlink_to(real_parent, target_is_directory=True)
        output = linked_parent / "subdirectory" / "page.pptx"

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(spec_path, prebuild_path, output, self.report)

        self.assertEqual(
            raised.exception.as_dict(),
            {
                "code": "BUILD_OUTPUT_INCOMPLETE",
                "path": str(linked_parent),
                "detail": "output parent must be an existing real directory",
            },
        )
        self.assertTrue(linked_parent.is_symlink())
        self.assertFalse((real_parent / "subdirectory" / "page.pptx").exists())
        self.assertFalse(self.report.exists())

    def test_cli_rejects_relative_nested_output_ancestor_symlink(self) -> None:
        spec_path, prebuild_path = write_valid_fixture(
            self.root, make_minimal_spec(self.root)
        )
        real_parent = self.root / "real-output"
        (real_parent / "subdirectory").mkdir(parents=True)
        linked_parent = self.root / "linked-output"
        linked_parent.symlink_to(real_parent, target_is_directory=True)
        relative_output = Path("linked-output/subdirectory/page.pptx")

        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPTS_ROOT / "build_pptx_from_spec.py"),
                "--spec",
                str(spec_path),
                "--prebuild-report",
                str(prebuild_path),
                "--output",
                str(relative_output),
                "--build-report",
                str(self.report),
            ],
            cwd=self.root,
            check=False,
            capture_output=True,
            text=True,
        )

        self.assertEqual(completed.returncode, 2, completed.stderr)
        self.assertEqual(completed.stderr, "")
        self.assertEqual(
            json.loads(completed.stdout),
            {
                "valid": False,
                "errors": [
                    {
                        "code": "BUILD_OUTPUT_INCOMPLETE",
                        "path": str(Path("linked-output")),
                        "detail": "output parent must be an existing real directory",
                    }
                ],
            },
        )
        self.assertTrue(linked_parent.is_symlink())
        self.assertFalse((real_parent / "subdirectory" / "page.pptx").exists())
        self.assertFalse(self.report.exists())

    def test_capability_gate_rejects_invalid_typography_before_rendering(self) -> None:
        cases = ("missing-italic", "invalid-vertical-alignment")
        for case in cases:
            with self.subTest(case=case):
                root = self.root / case
                spec = make_minimal_spec(root)
                contract = spec["modules"]["typography"]["items"][0]
                if case == "missing-italic":
                    del contract["runs"][0]["italic"]
                else:
                    contract["text_box"]["vertical_alignment"] = "diagonal"
                spec_path = root / "page-reconstruction.json"
                prebuild_path = root / "prebuild-report.json"
                output = root / "page.pptx"
                report_path = root / "build-report.json"
                spec_path.write_text(
                    json.dumps(spec, ensure_ascii=False), encoding="utf-8"
                )
                prebuild_path.write_text(
                    json.dumps(
                        {
                            "valid": True,
                            "stage": "prebuild",
                            "errors": [],
                            "spec_sha256": canonical_json_sha256(spec),
                        }
                    ),
                    encoding="utf-8",
                )

                error = None
                try:
                    _compile_single_page()(
                        spec_path, prebuild_path, output, report_path
                    )
                except Exception as exc:  # behavior assertion below owns the type
                    error = exc

                self.assertIsInstance(error, ToolError)
                assert isinstance(error, ToolError)
                self.assertEqual(error.code, "UNSUPPORTED_CAPABILITY")
                self.assertFalse(output.exists())
                self.assertFalse(report_path.exists())


class BuildPrimitiveTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self.temporary.cleanup)
        self.root = Path(self.temporary.name)

    def test_all_canonical_shapes_build_with_stable_names(self) -> None:
        for shape_type in sorted(CANONICAL_VALUES["shape_type"]):
            with self.subTest(shape_type=shape_type):
                root = self.root / shape_type
                pptx, report = compile_fixture(root, make_shape_spec(root, shape_type))

                shapes = {shape.name: shape for shape in Presentation(pptx).slides[0].shapes}
                self.assertIn("ia:card", shapes)
                self.assertEqual(report["elements"]["card"]["object_type"], "sp")
                self.assertEqual(report["elements"]["card"]["semantic_kind"], "shape")
                self.assertEqual(report["elements"]["card"]["selected_mode"], "native")
                if shape_type == "roundRect":
                    adjustment = shapes["ia:card"]._element.spPr.prstGeom.avLst[0]
                    self.assertEqual(adjustment.get("fmla"), "val 25000")
                if shape_type == "rightArrow":
                    self.assertEqual(
                        shapes["ia:card"].auto_shape_type,
                        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                    )
                    self.assertEqual(
                        shapes["ia:card"]._element.spPr.prstGeom.get("prst"),
                        "rightArrow",
                    )

    def test_right_arrow_rejects_adjustments_without_publication(self) -> None:
        spec = make_shape_spec(self.root, "rightArrow")
        spec["elements"][-1]["style"]["adjustments"] = [0.25]

        with self.assertRaises(ToolError) as raised:
            compile_with_claimed_prebuild(self.root, spec)

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertEqual(raised.exception.path, "elements.card.style.adjustments")
        self.assertEqual(
            raised.exception.detail,
            "adjustments are supported only for roundRect",
        )
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_round_rect_rejects_adjustment_that_quantizes_to_zero_without_publication(self) -> None:
        spec = make_shape_spec(self.root, "roundRect")
        spec["elements"][-1]["style"]["adjustments"] = [0.000001]
        spec_path = self.root / "page-reconstruction.json"
        prebuild_path = self.root / "prebuild-report.json"
        output = self.root / "page.pptx"
        report_path = self.root / "build-report.json"
        spec_path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        prebuild_path.write_text(
            json.dumps(
                {
                    "valid": True,
                    "stage": "prebuild",
                    "errors": [],
                    "spec_sha256": canonical_json_sha256(spec),
                }
            ),
            encoding="utf-8",
        )

        with self.assertRaises(ToolError) as raised:
            _compile_single_page()(spec_path, prebuild_path, output, report_path)

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertEqual(raised.exception.path, "elements.card.style.adjustments")
        self.assertEqual(raised.exception.capability, "shape.roundRect.adjustment")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_shape_line_width_error_preserves_shape_capability(self) -> None:
        spec = make_shape_spec(self.root, "rectangle")
        spec["elements"][-1]["style"]["line"]["width"] = 20116801

        with self.assertRaises(ToolError) as raised:
            compile_with_claimed_prebuild(self.root, spec)

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertEqual(raised.exception.path, "elements.card.style.line.width")
        self.assertEqual(raised.exception.capability, "shape.line")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_shape_emits_linear_gradient_explicit_line_and_outer_shadow(self) -> None:
        spec = make_shape_spec(self.root, "rectangle")
        spec["elements"][-1]["style"].update(
            {
                "fill": {
                    "type": "linear_gradient",
                    "angle": 45,
                    "stops": [
                        {"position": 0.0, "color": "#112233", "opacity": 1.0},
                        {"position": 1.0, "color": "#DDEEFF", "opacity": 0.5},
                    ],
                },
                "effects": {
                    "outer_shadow": {
                        "color": "#000000",
                        "opacity": 0.4,
                        "blur_radius": 63500,
                        "distance": 25400,
                        "angle": 45,
                    }
                },
            }
        )

        pptx, _ = compile_fixture(self.root, spec)

        shape = next(item for item in Presentation(pptx).slides[0].shapes if item.name == "ia:card")
        xml = ET.fromstring(shape._element.xml)
        self.assertEqual(len(xml.findall(".//a:gradFill/a:gsLst/a:gs", NS)), 2)
        self.assertIsNotNone(xml.find(".//a:effectLst/a:outerShdw", NS))
        line = xml.find(".//p:spPr/a:ln", NS)
        self.assertEqual(line.get("w") if line is not None else None, "12700")
        child_tags = [child.tag.rsplit("}", 1)[-1] for child in xml.find("p:spPr", NS)]
        self.assertLess(child_tags.index("xfrm"), child_tags.index("prstGeom"))
        self.assertLess(child_tags.index("prstGeom"), child_tags.index("gradFill"))
        self.assertLess(child_tags.index("gradFill"), child_tags.index("ln"))

    def test_shape_solid_fill_and_no_effects_are_explicit(self) -> None:
        spec = make_shape_spec(self.root, "ellipse")
        spec["elements"][-1]["style"]["fill"] = {
            "type": "solid", "color": "#2A9D8F", "opacity": 0.75
        }

        pptx, _ = compile_fixture(self.root, spec)

        shape = next(item for item in Presentation(pptx).slides[0].shapes if item.name == "ia:card")
        xml = ET.fromstring(shape._element.xml)
        self.assertEqual(xml.find(".//a:solidFill/a:srgbClr", NS).get("val"), "2A9D8F")
        self.assertEqual(xml.find(".//a:solidFill/a:srgbClr/a:alpha", NS).get("val"), "75000")
        self.assertIsNone(xml.find(".//a:effectLst", NS))
        self.assertEqual(xml.find(".//p:style/a:effectRef", NS).get("idx"), "0")

    def test_shape_rejects_360_degree_gradient_and_shadow_without_publication(self) -> None:
        cases = {
            "gradient": (
                "shape.fill.linear_gradient",
                lambda style: style.update(
                    {
                        "fill": {
                            "type": "linear_gradient",
                            "angle": 360,
                            "stops": [
                                {"position": 0.0, "color": "#112233", "opacity": 1.0},
                                {"position": 1.0, "color": "#DDEEFF", "opacity": 1.0},
                            ],
                        }
                    }
                ),
                "elements.card.style.fill.angle",
            ),
            "shadow": (
                "shape.effect.shadow",
                lambda style: style.update(
                    {
                        "effects": {
                            "outer_shadow": {
                                "color": "#000000",
                                "opacity": 0.4,
                                "blur_radius": 63500,
                                "distance": 25400,
                                "angle": 360,
                            }
                        }
                    }
                ),
                "elements.card.style.effects.outer_shadow.angle",
            ),
            "gradient-quantized-full-circle": (
                "shape.fill.linear_gradient",
                lambda style: style.update(
                    {
                        "fill": {
                            "type": "linear_gradient",
                            "angle": 359.999999,
                            "stops": [
                                {"position": 0.0, "color": "#112233", "opacity": 1.0},
                                {"position": 1.0, "color": "#DDEEFF", "opacity": 1.0},
                            ],
                        }
                    }
                ),
                "elements.card.style.fill.angle",
            ),
            "shadow-quantized-full-circle": (
                "shape.effect.shadow",
                lambda style: style.update(
                    {
                        "effects": {
                            "outer_shadow": {
                                "color": "#000000",
                                "opacity": 0.4,
                                "blur_radius": 63500,
                                "distance": 25400,
                                "angle": 359.999999,
                            }
                        }
                    }
                ),
                "elements.card.style.effects.outer_shadow.angle",
            ),
        }
        for name, (capability, mutate, expected_path) in cases.items():
            with self.subTest(name=name):
                root = self.root / name
                spec = make_shape_spec(root, "rectangle")
                mutate(spec["elements"][-1]["style"])

                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(root, spec)

                self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
                self.assertEqual(raised.exception.path, expected_path)
                self.assertEqual(raised.exception.capability, capability)
                self.assertFalse((root / "page.pptx").exists())
                self.assertFalse((root / "build-report.json").exists())


    def test_line_builds_connector_with_dash_transparency_and_arrowheads(self) -> None:
        spec = make_minimal_spec(self.root)
        _append_primitive(
            spec,
            element_id="arrow",
            kind="line",
            style={
                "line": {"color": "#E76F51", "width": 25400, "dash": "dashDot", "opacity": 0.5},
                "head_arrow": "oval",
                "tail_arrow": "triangle",
                "rotation": 0,
            },
            content={},
            slide_bbox=[609600, 914400, 1828800, 914400],
        )

        pptx, report = compile_fixture(self.root, spec)

        arrow = next(item for item in Presentation(pptx).slides[0].shapes if item.name == "ia:arrow")
        xml = ET.fromstring(arrow._element.xml)
        self.assertEqual(report["elements"]["arrow"]["object_type"], "cxnSp")
        self.assertEqual(xml.find(".//a:headEnd", NS).get("type"), "oval")
        self.assertEqual(xml.find(".//a:tailEnd", NS).get("type"), "triangle")
        self.assertEqual(xml.find(".//a:prstDash", NS).get("val"), "dashDot")
        self.assertEqual(xml.find(".//a:solidFill/a:srgbClr/a:alpha", NS).get("val"), "50000")

    def test_line_rejects_width_above_drawingml_limit_as_tool_error(self) -> None:
        spec = make_minimal_spec(self.root)
        _append_primitive(
            spec,
            element_id="arrow",
            kind="line",
            style={
                "line": {
                    "color": "#E76F51",
                    "width": 20116801,
                    "dash": "solid",
                    "opacity": 1.0,
                }
            },
            content={},
        )

        with self.assertRaises(ToolError) as raised:
            compile_with_claimed_prebuild(self.root, spec)

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertEqual(raised.exception.path, "elements.arrow.style.line.width")
        self.assertEqual(raised.exception.capability, "line.stroke")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_prebuild_rejects_stale_picture_asset_hash_without_output(self) -> None:
        spec = make_picture_spec(self.root)
        spec["elements"][-1]["content"]["asset"]["asset_sha256"] = "0" * 64

        issues = validate_representation_plan(spec)

        self.assertTrue(issues)
        self.assertEqual(issues[0].code, "ASSET_HASH_MISMATCH")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_prebuild_rejects_picture_asset_through_symlinked_parent(self) -> None:
        spec = make_picture_spec(self.root)
        asset = Path(spec["elements"][-1]["content"]["asset"]["path"])
        linked_parent = self.root.resolve() / "linked-assets"
        linked_parent.symlink_to(asset.parent, target_is_directory=True)
        spec["elements"][-1]["content"]["asset"]["path"] = str(linked_parent / asset.name)

        issues = validate_representation_plan(spec)

        self.assertTrue(issues)
        self.assertEqual(issues[0].code, "UNSUPPORTED_CAPABILITY")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_picture_places_original_media_with_crop_rotation_and_opacity(self) -> None:
        spec = make_picture_spec(self.root, mode="none")
        expected_sha256 = spec["elements"][-1]["content"]["asset"]["asset_sha256"]

        pptx, report = compile_fixture(self.root, spec)

        picture = next(item for item in Presentation(pptx).slides[0].shapes if item.name == "ia:photo")
        xml = ET.fromstring(picture._element.xml)
        src_rect = xml.find(".//a:srcRect", NS)
        self.assertEqual(
            {side: src_rect.get(side) for side in ("l", "t", "r", "b")},
            {"l": "10000", "t": "0", "r": "10000", "b": "0"},
        )
        self.assertEqual(xml.find(".//a:blip/a:alphaModFix", NS).get("amt"), "60000")
        self.assertEqual(picture.rotation, 15)
        self.assertEqual(report["elements"]["photo"]["objects"][0]["media_sha256"], expected_sha256)

    def test_picture_contain_and_cover_preserve_bbox_with_controlled_crop(self) -> None:
        expected = {
            "contain": {"l": "0", "t": "-50000", "r": "0", "b": "-50000"},
            "cover": {"l": "25000", "t": "0", "r": "25000", "b": "0"},
        }
        for mode in ("contain", "cover"):
            with self.subTest(mode=mode):
                root = self.root / mode
                spec = make_picture_spec(root, mode=mode)
                spec["elements"][-1]["content"]["crop"] = {
                    "left": 0.0, "top": 0.0, "right": 0.0, "bottom": 0.0
                }

                pptx, report = compile_fixture(root, spec)

                picture = next(item for item in Presentation(pptx).slides[0].shapes if item.name == "ia:photo")
                src_rect = ET.fromstring(picture._element.xml).find(".//a:srcRect", NS)
                self.assertEqual({side: src_rect.get(side) for side in ("l", "t", "r", "b")}, expected[mode])
                self.assertEqual(report["elements"]["photo"]["objects"][0]["bbox"], [609600, 914400, 1828800, 1828800])

    def test_picture_rejects_crop_that_quantizes_to_full_image_without_publication(self) -> None:
        cases = {
            "explicit-side": (
                "none",
                {"left": 0.999999, "top": 0.0, "right": 0.0, "bottom": 0.0},
            ),
            "cover-sum": (
                "cover",
                {"left": 0.499999, "top": 0.0, "right": 0.0, "bottom": 0.0},
            ),
        }
        for name, (mode, crop) in cases.items():
            with self.subTest(name=name):
                root = self.root / name
                spec = make_picture_spec(root, mode=mode)
                spec["elements"][-1]["content"]["crop"] = crop

                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(root, spec)

                self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
                self.assertEqual(raised.exception.path, "elements.photo.content.crop")
                self.assertEqual(raised.exception.capability, "picture.crop.explicit")
                self.assertFalse((root / "page.pptx").exists())
                self.assertFalse((root / "build-report.json").exists())

    def test_nested_primitive_contracts_fail_closed_without_publication(self) -> None:
        cases = []
        shape = make_shape_spec(self.root / "shape-unknown", "rectangle")
        shape["elements"][-1]["style"]["fill"] = {
            "type": "solid", "color": "#112233", "opacity": 1.0, "future": True
        }
        cases.append(("shape-unknown", shape))
        line = make_minimal_spec(self.root / "line-unknown")
        _append_primitive(
            line, element_id="arrow", kind="line",
            style={"line": {"color": "#112233", "width": 12700, "dash": "solid", "opacity": 1.0, "future": True}},
            content={},
        )
        cases.append(("line-unknown", line))
        arrow_value = make_minimal_spec(self.root / "arrow-value")
        _append_primitive(
            arrow_value, element_id="arrow", kind="line",
            style={
                "line": {"color": "#112233", "width": 12700, "dash": "solid", "opacity": 1.0},
                "head_arrow": [],
            },
            content={},
        )
        cases.append(("arrow-value", arrow_value))
        picture = make_picture_spec(self.root / "picture-unknown")
        picture["elements"][-1]["content"]["crop"]["future"] = 0.0
        cases.append(("picture-unknown", picture))
        for name, spec in cases:
            with self.subTest(name=name):
                root = self.root / name
                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(root, spec)
                self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
                self.assertFalse((root / "page.pptx").exists())
                self.assertFalse((root / "build-report.json").exists())

    def test_icon_consumes_matching_alpha_isolation_asset_record(self) -> None:
        spec = make_icon_spec(self.root)

        pptx, report = compile_fixture(self.root, spec)

        self.assertIn("ia:status-icon", {shape.name for shape in Presentation(pptx).slides[0].shapes})
        self.assertEqual(report["elements"]["status-icon"]["object_type"], "pic")
        self.assertEqual(
            report["elements"]["status-icon"]["objects"][0]["media_sha256"],
            spec["elements"][-1]["content"]["asset"]["asset_sha256"],
        )

    def test_picture_renderer_keeps_icon_mode_isolated(self) -> None:
        from pptx_builder.common import RenderContext
        from pptx_builder.pictures import ICON_RENDERER, PICTURE_RENDERER
        from pptx_builder.registry import ObjectRegistry

        icon_spec = make_icon_spec(self.root / "icon")
        icon = next(
            item for item in icon_spec["elements"] if item["kind"] == "icon"
        )
        icon_asset_context = RenderContext(
            slide=None,
            spec=icon_spec,
            representation_modes={icon["element_id"]: "asset"},
            typography={},
            registry=ObjectRegistry(),
        )
        icon_background_context = RenderContext(
            slide=None,
            spec=icon_spec,
            representation_modes={icon["element_id"]: "background_picture"},
            typography={},
            registry=ObjectRegistry(),
        )
        picture_spec = make_picture_spec(self.root / "picture")
        picture = next(
            item for item in picture_spec["elements"] if item["kind"] == "picture"
        )
        picture_background_context = RenderContext(
            slide=None,
            spec=picture_spec,
            representation_modes={picture["element_id"]: "background_picture"},
            typography={},
            registry=ObjectRegistry(),
        )

        self.assertEqual(
            ICON_RENDERER.validate_contract(icon, icon_asset_context), []
        )
        icon_issues = ICON_RENDERER.validate_contract(
            icon, icon_background_context
        )
        self.assertEqual(len(icon_issues), 1)
        self.assertEqual(
            (icon_issues[0].code, icon_issues[0].path),
            ("UNSUPPORTED_CAPABILITY", "elements.status-icon.representation"),
        )
        self.assertEqual(
            PICTURE_RENDERER.validate_contract(
                picture, picture_background_context
            ),
            [],
        )

    def test_icon_rejects_unknown_current_asset_record_fields(self) -> None:
        spec = make_icon_spec(self.root)
        spec["modules"]["icons"]["icons"][0]["future_crop_mode"] = "invented"

        with self.assertRaises(ToolError) as raised:
            compile_with_claimed_prebuild(self.root, spec)

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())


class BuildStructuredTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self.temporary.cleanup)
        self.root = Path(self.temporary.name)

    def _compile_claimed_valid(self, root: Path, spec: dict[str, Any]):
        root.mkdir(parents=True, exist_ok=True)
        spec_path = root / "page-reconstruction.json"
        prebuild_path = root / "prebuild-report.json"
        output = root / "page.pptx"
        report_path = root / "build-report.json"
        spec_path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        prebuild_path.write_text(
            json.dumps(
                {
                    "valid": True,
                    "stage": "prebuild",
                    "errors": [],
                    "spec_sha256": canonical_json_sha256(spec),
                }
            ),
            encoding="utf-8",
        )
        error: Exception | None = None
        try:
            _compile_single_page()(spec_path, prebuild_path, output, report_path)
        except Exception as exc:  # behavior assertions own the exact type
            error = exc
        return error, output, report_path

    def _compile_schema_number_literal(
        self,
        root: Path,
        spec: dict[str, Any],
        literal: str,
    ):
        root.mkdir(parents=True, exist_ok=True)
        spec_path = root / "page-reconstruction.json"
        prebuild_path = root / "prebuild-report.json"
        output = root / "page.pptx"
        report_path = root / "build-report.json"
        marker = '"__NON_FINITE_NUMBER__"'
        serialized = json.dumps(spec, ensure_ascii=False, allow_nan=False)
        self.assertEqual(serialized.count(marker), 1)
        spec_path.write_text(serialized.replace(marker, literal), encoding="utf-8")
        prebuild_path.write_text(
            json.dumps(
                {
                    "valid": True,
                    "stage": "prebuild",
                    "errors": [],
                    "spec_sha256": "0" * 64,
                }
            ),
            encoding="utf-8",
        )
        error: Exception | None = None
        try:
            _compile_single_page()(spec_path, prebuild_path, output, report_path)
        except Exception as exc:  # behavior assertions own the exact type
            error = exc
        return error, output, report_path

    def test_merged_table_preserves_sizes_styles_and_local_borders(self) -> None:
        pptx, report = compile_structured_fixture(
            self.root, make_merged_table_spec(self.root)
        )

        table_shape = next(
            shape
            for shape in Presentation(pptx).slides[0].shapes
            if shape.has_table
        )
        table = table_shape.table
        self.assertEqual(report["elements"]["table"]["object_type"], "graphicFrame")
        self.assertEqual(table_shape.name, "ia:table")
        self.assertEqual([row.height for row in table.rows], [320040, 594360])
        self.assertEqual([column.width for column in table.columns], [609600, 1219200])
        self.assertTrue(table.cell(0, 0).is_merge_origin)
        self.assertTrue(table.cell(0, 1).is_spanned)
        header = table.cell(0, 0)
        self.assertEqual(
            [header.margin_left, header.margin_right, header.margin_top, header.margin_bottom],
            [101, 102, 103, 104],
        )
        self.assertEqual(header.text_frame.paragraphs[0].alignment, 2)
        run = header.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "Arial")
        self.assertEqual(run.font.size.pt, 14)
        self.assertTrue(run.font.bold)
        self.assertTrue(run.font.italic)
        self.assertEqual(str(header.fill.fore_color.rgb), "DDEEFF")
        self.assertEqual(header.vertical_anchor, MSO_ANCHOR.MIDDLE)
        self.assertEqual(str(run.font.color.rgb), "112233")
        tc_pr = header._tc.tcPr
        for side in ("lnL", "lnR", "lnT"):
            node = tc_pr.find(f"{{{NS['a']}}}{side}")
            self.assertIsNotNone(node)
            self.assertIsNotNone(node.find("a:noFill", NS))
        bottom = tc_pr.find("a:lnB", NS)
        self.assertIsNotNone(bottom)
        self.assertEqual(bottom.get("w"), "12700")
        self.assertEqual(bottom.find("a:solidFill/a:srgbClr", NS).get("val"), "445566")
        declared_borders = {(0, 0): {"lnB"}, (1, 1): {"lnL"}}
        for row_index in range(2):
            for column_index in range(2):
                properties = table.cell(row_index, column_index)._tc.tcPr
                for side in ("lnL", "lnR", "lnT", "lnB"):
                    node = properties.find(f"{{{NS['a']}}}{side}")
                    self.assertIsNotNone(
                        node,
                        f"physical cell {row_index},{column_index} lacks {side}",
                    )
                    if side not in declared_borders.get(
                        (row_index, column_index), set()
                    ):
                        self.assertIsNotNone(
                            node.find("a:noFill", NS),
                            f"undeclared {side} on {row_index},{column_index} must be noFill",
                        )

    def test_table_accepts_a_legitimate_empty_cell(self) -> None:
        spec = make_merged_table_spec(self.root)
        spec["elements"][-1]["content"]["cells"][1]["text"] = ""
        error: Exception | None = None

        try:
            pptx, _ = compile_fixture(self.root, spec)
        except Exception as exc:  # behavior assertion below owns the failure type
            error = exc

        self.assertIsNone(error, f"empty table cell must compile: {error!r}")
        table = next(
            shape.table
            for shape in Presentation(pptx).slides[0].shapes
            if shape.has_table
        )
        self.assertEqual(table.cell(1, 0).text, "")

    def test_repeat_status_matches_explicit_parts(self) -> None:
        explicit_root = self.root / "explicit"
        repeat_root = self.root / "repeat"

        explicit = compile_structured_fixture(
            explicit_root, make_status_spec(explicit_root, repeat=False)
        )
        repeated = compile_structured_fixture(
            repeat_root, make_status_spec(repeat_root, repeat=True)
        )

        self.assertEqual(
            explicit[1]["elements"]["status"]["objects"],
            repeated[1]["elements"]["status"]["objects"],
        )
        shapes = {
            shape.name: shape
            for shape in Presentation(repeated[0]).slides[0].shapes
        }
        self.assertEqual(
            set(shapes) & {"ia:status:segment-0", "ia:status:segment-1"},
            {"ia:status:segment-0", "ia:status:segment-1"},
        )
        self.assertEqual(shapes["ia:status:segment-0"].text, "已完成")
        self.assertEqual(
            shapes["ia:status:segment-0"].text_frame.paragraphs[0].runs[0].font.name,
            "Arial",
        )
        left = shapes["ia:status:segment-0"]
        right = shapes["ia:status:segment-1"]
        self.assertEqual(str(left.fill.fore_color.rgb), "22AA66")
        self.assertEqual(str(right.fill.fore_color.rgb), "8899AA")
        self.assertEqual(str(left.line.color.rgb), "245566")
        self.assertEqual(left.line.width, 12700)
        self.assertEqual(left.rotation, 12.5)
        self.assertEqual(
            [
                left.text_frame.margin_left,
                left.text_frame.margin_right,
                left.text_frame.margin_top,
                left.text_frame.margin_bottom,
            ],
            [101, 102, 103, 104],
        )
        self.assertEqual(left.text_frame.word_wrap, True)
        self.assertEqual(left.text_frame.vertical_anchor, MSO_ANCHOR.MIDDLE)
        paragraph = left.text_frame.paragraphs[0]
        self.assertEqual(paragraph.alignment, PP_ALIGN.CENTER)
        self.assertEqual(paragraph.runs[0].font.size.pt, 12)
        self.assertEqual(str(paragraph.runs[0].font.color.rgb), "FFFFFF")

    def test_multipart_error_preserves_the_schema_part_path(self) -> None:
        spec = make_status_spec(self.root, repeat=True)
        spec["elements"][-1]["content"]["part_defaults"]["style"]["line"][
            "width"
        ] = 20_116_801

        with self.assertRaises(ToolError) as raised:
            compile_with_claimed_prebuild(self.root, spec)

        self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
        self.assertEqual(
            raised.exception.path,
            "elements.status.content.part_defaults.style.line.width",
        )
        self.assertEqual(raised.exception.capability, "shape.line")
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_multipart_alignment_type_error_is_a_stable_tool_error(self) -> None:
        spec = make_status_spec(self.root, repeat=True)
        spec["elements"][-1]["content"]["part_defaults"]["style"]["text_style"][
            "alignment"
        ] = []
        spec_path = self.root / "page-reconstruction.json"
        prebuild_path = self.root / "prebuild-report.json"
        output = self.root / "page.pptx"
        report_path = self.root / "build-report.json"
        spec_path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        prebuild_path.write_text(
            json.dumps(
                {
                    "valid": True,
                    "stage": "prebuild",
                    "errors": [],
                    "spec_sha256": canonical_json_sha256(spec),
                }
            ),
            encoding="utf-8",
        )
        error: Exception | None = None

        try:
            _compile_single_page()(spec_path, prebuild_path, output, report_path)
        except Exception as exc:  # behavior assertion below owns the failure type
            error = exc

        self.assertIsInstance(error, ToolError)
        assert isinstance(error, ToolError)
        self.assertEqual(error.code, "UNSUPPORTED_CAPABILITY")
        self.assertEqual(
            error.path,
            "elements.status.content.part_defaults.style.text_style.alignment",
        )
        self.assertFalse(output.exists())
        self.assertFalse(report_path.exists())

    def test_matrix_compiles_through_the_shared_multipart_renderer(self) -> None:
        spec = make_status_spec(self.root, repeat=True)
        element = spec["elements"][-1]
        element["element_id"] = "matrix"
        element["kind"] = "matrix"
        spec["regions"][0]["element_ids"][-1] = "matrix"
        spec["reading_order"][-1] = "matrix"
        plan = spec["modules"]["representation_plan"]["items"][-1]
        plan["source_fact_id"] = "fact-matrix"
        plan["semantic_role"] = "matrix"
        plan["bound_element_ids"] = ["matrix"]

        pptx, report = compile_fixture(self.root, spec)

        self.assertEqual(report["elements"]["matrix"]["object_type"], "sp")
        self.assertEqual(
            {
                shape.name
                for shape in Presentation(pptx).slides[0].shapes
                if shape.name.startswith("ia:matrix:")
            },
            {"ia:matrix:segment-0", "ia:matrix:segment-1"},
        )

    def test_table_numeric_overflow_fails_as_tool_error_without_publication(self) -> None:
        cases = (
            ("font", ("font", "size"), 4000.01, "font.size"),
            ("margin", ("margins", "left"), 2_147_483_648, "margins.left"),
        )
        for name, path, value, suffix in cases:
            with self.subTest(name=name):
                root = self.root / f"table-{name}"
                spec = make_merged_table_spec(root)
                cell = spec["elements"][-1]["content"]["cells"][0]
                cell[path[0]][path[1]] = value

                error, output, report = self._compile_claimed_valid(root, spec)

                self.assertIsInstance(error, ToolError)
                assert isinstance(error, ToolError)
                self.assertEqual(error.code, "UNSUPPORTED_CAPABILITY")
                self.assertEqual(
                    error.path,
                    f"elements.table.content.cells[0].{suffix}",
                )
                self.assertFalse(output.exists())
                self.assertFalse(report.exists())

    def test_multipart_numeric_overflow_fails_as_tool_error_without_publication(self) -> None:
        cases = (
            ("font", ("font_size",), 4000.01, "font_size"),
            ("margin", ("margins", "left"), 2_147_483_648, "margins.left"),
        )
        for name, path, value, suffix in cases:
            with self.subTest(name=name):
                root = self.root / f"multipart-{name}"
                spec = make_status_spec(root, repeat=True)
                style = spec["elements"][-1]["content"]["part_defaults"]["style"][
                    "text_style"
                ]
                if len(path) == 1:
                    style[path[0]] = value
                else:
                    style[path[0]][path[1]] = value

                error, output, report = self._compile_claimed_valid(root, spec)

                self.assertIsInstance(error, ToolError)
                assert isinstance(error, ToolError)
                self.assertEqual(error.code, "UNSUPPORTED_CAPABILITY")
                self.assertEqual(
                    error.path,
                    f"elements.status.content.part_defaults.style.text_style.{suffix}",
                )
                self.assertFalse(output.exists())
                self.assertFalse(report.exists())

    def test_table_and_multipart_reject_rotations_that_collapse_in_drawingml(self) -> None:
        cases = (
            (
                "table",
                make_merged_table_spec,
                lambda spec: spec["elements"][-1]["style"].update(
                    {"rotation": 359.999999}
                ),
                "elements.table.style.rotation",
            ),
            (
                "multipart",
                lambda root: make_status_spec(root, repeat=True),
                lambda spec: spec["elements"][-1]["content"]["part_defaults"][
                    "style"
                ].update({"rotation": 359.999999}),
                "elements.status.content.part_defaults.style.rotation",
            ),
        )
        for name, make_spec, mutate, expected_path in cases:
            with self.subTest(name=name):
                root = self.root / f"rotation-{name}"
                spec = make_spec(root)
                mutate(spec)

                error, output, report = self._compile_claimed_valid(root, spec)

                self.assertIsInstance(error, ToolError)
                assert isinstance(error, ToolError)
                self.assertEqual(error.code, "UNSUPPORTED_CAPABILITY")
                self.assertEqual(error.path, expected_path)
                self.assertFalse(output.exists())
                self.assertFalse(report.exists())

    def test_compiler_rejects_overflowed_json_numbers_without_publication(self) -> None:
        cases = (
            (
                "table",
                make_merged_table_spec,
                lambda spec: spec["elements"][-1]["content"]["cells"][0]["font"].update(
                    {"size": "__NON_FINITE_NUMBER__"}
                ),
                "elements[2].content.cells[0].font.size",
            ),
            (
                "multipart",
                lambda root: make_status_spec(root, repeat=True),
                lambda spec: spec["elements"][-1]["content"]["part_defaults"]["style"][
                    "text_style"
                ].update({"font_size": "__NON_FINITE_NUMBER__"}),
                "elements[2].content.part_defaults.style.text_style.font_size",
            ),
        )
        for name, make_spec, mutate, expected_path in cases:
            with self.subTest(name=name):
                root = self.root / f"non-finite-{name}"
                spec = make_spec(root)
                mutate(spec)

                error, output, report = self._compile_schema_number_literal(
                    root, spec, "1e400"
                )

                self.assertIsInstance(error, ToolError)
                assert isinstance(error, ToolError)
                self.assertEqual(error.code, "SPEC_NUMBER_NON_FINITE")
                self.assertEqual(error.path, expected_path)
                self.assertEqual(error.detail, "number must be finite")
                self.assertFalse(output.exists())
                self.assertFalse(report.exists())

    def test_compiler_rejects_nonstandard_json_number_tokens(self) -> None:
        root = self.root / "nonstandard-json-number"
        spec = make_status_spec(root, repeat=True)
        spec["elements"][-1]["content"]["part_defaults"]["style"]["text_style"].update(
            {"font_size": "__NON_FINITE_NUMBER__"}
        )

        error, output, report = self._compile_schema_number_literal(
            root, spec, "NaN"
        )

        self.assertIsInstance(error, ToolError)
        assert isinstance(error, ToolError)
        self.assertEqual(error.code, "SCHEMA_JSON_INVALID")
        self.assertEqual(error.path, str((root / "page-reconstruction.json").resolve()))
        self.assertFalse(output.exists())
        self.assertFalse(report.exists())


class BuildTextContractFieldTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self.temporary.cleanup)
        self.root = Path(self.temporary.name)

    def test_capability_gate_rejects_unknown_typography_fields_without_publication(self) -> None:
        cases = (
            ("item", "future_effect", "glow", lambda item: item),
            ("text-box", "future_effect", "glow", lambda item: item["text_box"]),
            ("paragraph", "future_effect", "glow", lambda item: item["paragraphs"][0]),
            ("run", "future_effect", "glow", lambda item: item["runs"][0]),
            ("list", "future_effect", "glow", lambda item: item["paragraphs"][0]["list"]),
            ("legacy-decoration", "decoration", "none", lambda item: item["runs"][0]),
        )
        for case, field, value, select_contract in cases:
            with self.subTest(case=case):
                root = self.root / case
                spec = make_minimal_spec(root)
                typography_item = spec["modules"]["typography"]["items"][0]
                select_contract(typography_item)[field] = value
                spec_path = root / "page-reconstruction.json"
                prebuild_path = root / "prebuild-report.json"
                output = root / "page.pptx"
                report_path = root / "build-report.json"
                spec_path.write_text(
                    json.dumps(spec, ensure_ascii=False), encoding="utf-8"
                )
                prebuild_path.write_text(
                    json.dumps(
                        {
                            "valid": True,
                            "stage": "prebuild",
                            "errors": [],
                            "spec_sha256": canonical_json_sha256(spec),
                        }
                    ),
                    encoding="utf-8",
                )

                with self.assertRaises(ToolError) as raised:
                    _compile_single_page()(
                        spec_path, prebuild_path, output, report_path
                    )

                self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")
                self.assertFalse(output.exists())
                self.assertFalse(report_path.exists())


if __name__ == "__main__":
    unittest.main()
