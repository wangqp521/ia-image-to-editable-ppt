"""Machine-authorable schema-v2 contract and generated-schema regressions."""

from __future__ import annotations

import copy
import importlib
import json
import os
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS_ROOT = ROOT / "scripts"
GENERATED_SCHEMA = ROOT / "schemas" / "page-reconstruction-v2.schema.json"


class SchemaContractTests(unittest.TestCase):
    def _api(self):
        scripts_root = str(SCRIPTS_ROOT)
        if scripts_root not in sys.path:
            sys.path.insert(0, scripts_root)
        try:
            return importlib.import_module("lib.schema_contracts")
        except ModuleNotFoundError as exc:
            if exc.name == "lib.schema_contracts":
                self.fail("shared schema contracts are not implemented")
            raise

    def test_manifest_is_deterministic_with_exact_authoring_envelopes(self) -> None:
        contracts = self._api()

        first = contracts.schema_contract_manifest()
        second = contracts.schema_contract_manifest()

        self.assertEqual(first, second)
        self.assertEqual(first["contract_id"], "page-reconstruction-v2")
        self.assertEqual(first["schema_version"], 2)
        records = first["records"]
        expected_root = {
            "schema_version",
            "page_id",
            "verification_profile",
            "delivery_status",
            "session_reuse",
            "content_reference",
            "clean_visual_reference",
            "canvas",
            "activated_modules",
            "modules",
            "regions",
            "elements",
            "reading_order",
            "visual_gate",
            "editability_gate",
            "runtime_preflight",
        }
        expected_required_root = expected_root - {
            "verification_profile",
            "delivery_status",
            "runtime_preflight",
        }
        expected_element = {
            "element_id",
            "kind",
            "source_bbox",
            "slide_bbox",
            "layer",
            "editable",
            "confidence",
            "style",
            "content",
        }
        expected_icon_module = {
            "schema_version",
            "page_id",
            "slide_coordinate_unit",
            "clean_visual_reference",
            "clean_visual_sha256",
            "icons",
        }
        expected_icon_item = {
            "icon_id",
            "element_id",
            "category",
            "instance_count",
            "repeat_group",
            "semantic_scope",
            "source_bbox",
            "slide_bbox",
            "layer",
            "source_path",
            "source_sha256",
            "crop_mode",
            "padding",
            "background_handling",
            "asset_path",
            "asset_sha256",
            "alpha_mask_sha256",
            "final_width",
            "final_height",
            "sharpness",
            "validation",
            "native_redraw",
            "selectable_picture_verified",
            "object_type",
        }
        expected_background_provenance = {
            "kind",
            "source_path",
            "source_sha256",
        }
        expected_background_item = {
            "background_id",
            "role",
            "source_bbox",
            "selected_mode",
            "bound_element_id",
            "source_provenance",
            "reason",
            "evidence",
            "contains_foreground_semantics",
        }
        self.assertEqual(set(records["PageReconstruction"]["properties"]), expected_root)
        self.assertEqual(
            set(records["PageReconstruction"]["required"]),
            expected_required_root,
        )
        self.assertFalse(records["PageReconstruction"]["additionalProperties"])
        self.assertEqual(set(records["Element"]["properties"]), expected_element)
        self.assertEqual(set(records["Element"]["required"]), expected_element)
        self.assertEqual(set(records["IconsModule"]["properties"]), expected_icon_module)
        self.assertEqual(set(records["IconsModule"]["required"]), expected_icon_module)
        self.assertEqual(set(records["IconItem"]["properties"]), expected_icon_item)
        self.assertEqual(set(records["IconItem"]["required"]), expected_icon_item)
        self.assertFalse(records["IconItem"]["additionalProperties"])
        self.assertEqual(
            set(records["BackgroundProvenance"]["properties"]),
            expected_background_provenance,
        )
        self.assertEqual(
            set(records["BackgroundProvenance"]["required"]),
            expected_background_provenance,
        )
        self.assertEqual(
            set(records["BackgroundItem"]["properties"]),
            expected_background_item,
        )
        self.assertEqual(
            set(records["BackgroundItem"]["required"]), expected_background_item
        )
        self.assertEqual(
            set(records["BackgroundModule"]["properties"]), {"items"}
        )
        modules = records["Modules"]
        self.assertEqual(
            modules["properties"]["background"],
            {"$ref": "#/$defs/BackgroundModule"},
        )
        self.assertIn("background", modules["required"])

    def test_shared_fields_match_validator_and_compiler_consumers(self) -> None:
        contracts = self._api()
        from lib import capabilities, element_contracts, representation_contracts

        self.assertEqual(element_contracts.ELEMENT_FIELDS, contracts.ELEMENT_FIELDS)
        self.assertEqual(
            representation_contracts.REQUIRED_FIELDS,
            contracts.REPRESENTATION_ITEM_FIELDS,
        )
        self.assertEqual(
            capabilities.TEXT_CONTRACT_ALLOWED_FIELDS,
            contracts.TEXT_CONTRACT_ALLOWED_FIELDS,
        )
        self.assertEqual(
            capabilities.TEXT_RUN_MODERN_ALLOWED_FIELDS,
            contracts.TEXT_RUN_FIELDS,
        )

    def test_generated_schema_is_checked_in_byte_for_byte_with_identity_hash(self) -> None:
        contracts = self._api()
        expected = contracts.canonical_json_schema()

        self.assertTrue(GENERATED_SCHEMA.is_file(), "generated JSON Schema is missing")
        self.assertEqual(GENERATED_SCHEMA.read_text(encoding="utf-8"), expected)
        document = json.loads(expected)
        self.assertEqual(
            document["x-schema-contract-sha256"],
            contracts.schema_contract_sha256(),
        )
        self.assertEqual(document["$ref"], "#/$defs/PageReconstruction")
        self.assertFalse(document["$defs"]["IconItem"]["additionalProperties"])
        self.assertIn("examples", document["$defs"]["TypographyItem"])
        self.assertIn("examples", document["$defs"]["Element"])

        shape_style = document["$defs"]["ShapeStyle"]
        self.assertEqual(shape_style["required"], ["shape_type"])
        self.assertEqual(
            shape_style["properties"]["shape_type"]["enum"],
            ["chevron", "ellipse", "rectangle", "rightArrow", "roundRect"],
        )
        self.assertEqual(
            shape_style["properties"]["line"], {"$ref": "#/$defs/Stroke"}
        )
        self.assertEqual(
            shape_style["properties"]["effects"]["oneOf"][0], {"const": "none"}
        )
        picture_content = document["$defs"]["PictureContent"]
        self.assertEqual(picture_content["required"], ["asset", "crop", "mode"])
        self.assertEqual(
            picture_content["properties"]["mode"]["enum"],
            ["contain", "cover", "none"],
        )
        self.assertEqual(
            picture_content["properties"]["crop"], {"$ref": "#/$defs/Crop"}
        )
        self.assertFalse(document["$defs"]["Stroke"]["additionalProperties"])

    def test_write_json_schema_top_level_entry_writes_exact_bytes(self) -> None:
        contracts = self._api()
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "schema.json"
            path_before = list(sys.path)

            contracts.write_json_schema(output)

            self.assertEqual(
                contracts.canonical_json_schema().encode("utf-8"),
                output.read_bytes(),
            )
            self.assertEqual(path_before, sys.path)

    def test_write_json_schema_package_entry_ignores_foreign_lib_from_other_cwd(
        self,
    ) -> None:
        script = """
import sys
import types
from pathlib import Path

foreign_lib = types.ModuleType("lib")
foreign_lib.marker = object()
sys.modules["lib"] = foreign_lib
from scripts.lib import schema_contracts

path_before = list(sys.path)
root = Path(sys.argv[1])
try:
    schema_contracts.write_json_schema(root / "missing" / "schema.json")
except Exception as exc:
    if type(exc) is not FileNotFoundError:
        raise AssertionError(f"expected FileNotFoundError, got {type(exc).__name__}") from exc
else:
    raise AssertionError("missing output directory did not fail")

output = root / "schema.json"
schema_contracts.write_json_schema(output)
if output.read_bytes() != schema_contracts.canonical_json_schema().encode("utf-8"):
    raise AssertionError("generated schema bytes differ")
if sys.path != path_before:
    raise AssertionError("write_json_schema mutated sys.path")
if sys.modules["lib"] is not foreign_lib:
    raise AssertionError("package call replaced the unrelated lib module")
"""
        with tempfile.TemporaryDirectory() as directory:
            environment = dict(os.environ)
            prior_pythonpath = environment.get("PYTHONPATH")
            environment["PYTHONPATH"] = (
                str(ROOT)
                if not prior_pythonpath
                else str(ROOT) + os.pathsep + prior_pythonpath
            )

            completed = subprocess.run(
                [sys.executable, "-c", script, directory],
                cwd=directory,
                env=environment,
                capture_output=True,
                text=True,
                check=False,
            )

        self.assertEqual(
            0,
            completed.returncode,
            f"stdout:\n{completed.stdout}\nstderr:\n{completed.stderr}",
        )

    def test_record_constructor_rejects_stale_aliases_and_preserves_exact_keys(self) -> None:
        contracts = self._api()
        values = {
            "region_id": "region-001",
            "source_bbox": [0, 0, 1600, 900],
            "slide_bbox": [0, 0, 12192000, 6858000],
            "layer": 0,
            "padding": {"left": 0, "right": 0, "top": 0, "bottom": 0},
            "element_ids": [],
        }

        record = contracts.construct_record("Region", **values)

        self.assertEqual(record, values)
        for stale in ("instance", "repeat", "pixel_bbox"):
            with self.subTest(stale=stale), self.assertRaises(
                contracts.ContractConstructionError
            ):
                contracts.construct_record("Region", **values, **{stale: None})

    def test_generated_shape_example_passes_real_prebuild_and_compiler(self) -> None:
        from tests.fixture_specs import make_minimal_spec
        from tests.test_build_pptx_from_spec import _append_primitive, compile_fixture

        contracts = self._api()
        shape = copy.deepcopy(
            contracts.json_schema_document()["$defs"]["Element"]["examples"][1]
        )
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            spec = make_minimal_spec(root)
            _append_primitive(
                spec,
                element_id=shape["element_id"],
                kind=shape["kind"],
                style=shape["style"],
                content=shape["content"],
                source_bbox=shape["source_bbox"],
                slide_bbox=shape["slide_bbox"],
            )

            pptx, report = compile_fixture(root / "compile", spec)
            pptx_exists = pptx.is_file()

        self.assertTrue(pptx_exists)
        self.assertEqual(report["elements"][shape["element_id"]]["object_type"], "sp")

    def test_root_and_region_unknown_fields_match_all_production_gates(self) -> None:
        contracts = self._api()
        from lib.error_codes import ToolError
        from tests.fixture_specs import make_minimal_spec
        from tests.test_build_pptx_from_spec import compile_with_claimed_prebuild

        validator_spec = importlib.util.spec_from_file_location(
            "schema_contract_root_region_validator",
            SCRIPTS_ROOT / "validate_reconstruction_spec.py",
        )
        self.assertIsNotNone(validator_spec)
        self.assertIsNotNone(validator_spec.loader if validator_spec else None)
        assert validator_spec is not None and validator_spec.loader is not None
        validator = importlib.util.module_from_spec(validator_spec)
        validator_spec.loader.exec_module(validator)
        cases = (
            (
                "PageReconstruction",
                lambda spec: spec.update({"stale_root_alias": True}),
                {
                    "code": "UNSUPPORTED_CAPABILITY",
                    "path": "$",
                    "detail": "unknown fields: stale_root_alias",
                },
            ),
            (
                "Region",
                lambda spec: spec["regions"][0].update(
                    {"pixel_bbox": [0, 0, 1600, 900]}
                ),
                {
                    "code": "UNSUPPORTED_CAPABILITY",
                    "path": "regions.region-001",
                    "detail": "unknown fields: pixel_bbox",
                },
            ),
        )

        for record_name, mutate, expected in cases:
            with self.subTest(record_name=record_name), tempfile.TemporaryDirectory() as directory:
                root = Path(directory)
                spec = make_minimal_spec(root)
                mutate(spec)

                prebuild = validator.validate_spec(spec, stage="prebuild")
                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(root / "compile", spec)

                self.assertEqual(
                    contracts.unknown_field_detail(
                        record_name,
                        spec if record_name == "PageReconstruction" else spec["regions"][0],
                    ),
                    expected["detail"],
                )
                self.assertIn(expected, prebuild["errors"])
                self.assertEqual(raised.exception.as_dict(), expected)
                self.assertFalse(
                    contracts.json_schema_document()["$defs"][record_name][
                        "additionalProperties"
                    ]
                )

    def test_nested_exact_envelopes_reject_stale_fields_at_both_gates(self) -> None:
        contracts = self._api()
        from lib.error_codes import ToolError
        from tests.fixture_specs import make_minimal_spec
        from tests.test_build_pptx_from_spec import compile_with_claimed_prebuild

        validator_spec = importlib.util.spec_from_file_location(
            "schema_contract_nested_envelope_validator",
            SCRIPTS_ROOT / "validate_reconstruction_spec.py",
        )
        self.assertIsNotNone(validator_spec)
        self.assertIsNotNone(validator_spec.loader if validator_spec else None)
        assert validator_spec is not None and validator_spec.loader is not None
        validator = importlib.util.module_from_spec(validator_spec)
        validator_spec.loader.exec_module(validator)
        cases = (
            ("SessionReuse", ("session_reuse",), "session_reuse"),
            ("Reference", ("content_reference",), "content_reference"),
            ("Canvas", ("canvas",), "canvas"),
            ("Modules", ("modules",), "modules"),
            ("PageLayoutModule", ("modules", "page_layout"), "modules.page_layout"),
            (
                "CoordinateOverlayEvidence",
                ("modules", "page_layout", "coordinate_overlay_evidence"),
                "modules.page_layout.coordinate_overlay_evidence",
            ),
        )

        for record_name, keys, path in cases:
            with self.subTest(record_name=record_name), tempfile.TemporaryDirectory() as directory:
                root = Path(directory)
                spec = make_minimal_spec(root)
                target = spec
                for key in keys:
                    target = target[key]
                target["stale_alias"] = True
                expected = {
                    "code": "UNSUPPORTED_CAPABILITY",
                    "path": path,
                    "detail": "unknown fields: stale_alias",
                }

                prebuild = validator.validate_spec(spec, stage="prebuild")
                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(root / "compile", spec)

                self.assertEqual(prebuild["errors"], [expected])
                self.assertEqual(raised.exception.as_dict(), expected)
                self.assertEqual(
                    contracts.unknown_field_detail(record_name, target),
                    expected["detail"],
                )
                self.assertFalse(
                    contracts.json_schema_document()["$defs"][record_name][
                        "additionalProperties"
                    ]
                )

    def test_all_exact_envelopes_reject_missing_fields_identically_at_both_gates(self) -> None:
        self._api()
        from lib.error_codes import ToolError
        from tests.fixture_specs import make_minimal_spec
        from tests.test_build_pptx_from_spec import compile_with_claimed_prebuild

        validator_spec = importlib.util.spec_from_file_location(
            "schema_contract_missing_envelope_validator",
            SCRIPTS_ROOT / "validate_reconstruction_spec.py",
        )
        self.assertIsNotNone(validator_spec)
        self.assertIsNotNone(validator_spec.loader if validator_spec else None)
        assert validator_spec is not None and validator_spec.loader is not None
        validator = importlib.util.module_from_spec(validator_spec)
        validator_spec.loader.exec_module(validator)
        cases = (
            ("PageReconstruction", (), "visual_gate", "$"),
            ("SessionReuse", ("session_reuse",), "reason", "session_reuse"),
            ("Reference", ("content_reference",), "sha256", "content_reference"),
            ("Canvas", ("canvas",), "background", "canvas"),
            ("Modules", ("modules",), "page_layout", "modules"),
            (
                "PageLayoutModule",
                ("modules", "page_layout"),
                "density_targets",
                "modules.page_layout",
            ),
            (
                "CoordinateOverlayEvidence",
                ("modules", "page_layout", "coordinate_overlay_evidence"),
                "inspection",
                "modules.page_layout.coordinate_overlay_evidence",
            ),
            ("Region", ("regions", 0), "padding", "regions.region-001"),
        )

        for record_name, keys, missing_field, path in cases:
            with self.subTest(record_name=record_name), tempfile.TemporaryDirectory() as directory:
                root = Path(directory)
                spec = make_minimal_spec(root)
                target = spec
                for key in keys:
                    target = target[key]
                del target[missing_field]
                expected = {
                    "code": "UNSUPPORTED_CAPABILITY",
                    "path": path,
                    "detail": f"missing fields: {missing_field}",
                }
                compile_root = root / "compile"

                self.assertEqual(
                    self._api().schema_envelope_issues(spec),
                    [(path, expected["detail"])],
                )
                prebuild = validator.validate_spec(spec, stage="prebuild")
                with self.assertRaises(ToolError) as raised:
                    compile_with_claimed_prebuild(compile_root, spec)

                self.assertEqual(prebuild["errors"], [expected])
                self.assertEqual(raised.exception.as_dict(), expected)
                self.assertFalse((compile_root / "page.pptx").exists())
                self.assertFalse((compile_root / "build-report.json").exists())

    def test_complete_text_style_matches_schema_prebuild_compiler_and_render(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR
        from tests.fixture_specs import make_minimal_spec
        from tests.test_build_pptx_from_spec import compile_fixture
        from pptx_builder.text import TEXT_RENDERER

        contracts = self._api()
        complete_style = {
            "fill": "noFill",
            "line": {
                "color": "#445566",
                "width": 12700,
                "dash": "solid",
                "opacity": 1,
            },
            "margins": {"left": 0, "right": 0, "top": 0, "bottom": 0},
            "vertical_alignment": "top",
            "wrap": False,
            "rotation": 7,
            "effects": "none",
        }
        text_style = contracts.json_schema_document()["$defs"]["TextStyle"]
        self.assertEqual(
            set(text_style["properties"]), contracts.KIND_STYLE_FIELDS["text"]
        )
        self.assertEqual(
            TEXT_RENDERER.supported_fields,
            contracts.KIND_STYLE_FIELDS["text"]
            | contracts.KIND_CONTENT_FIELDS["text"],
        )
        self.assertEqual(
            TEXT_RENDERER.required_fields,
            contracts.KIND_REQUIRED_STYLE_FIELDS["text"]
            | contracts.KIND_REQUIRED_CONTENT_FIELDS["text"],
        )

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            spec = make_minimal_spec(root)
            spec["elements"][0]["style"] = complete_style

            pptx, _report = compile_fixture(root / "compile", spec)
            presentation = Presentation(pptx)
            shape = next(
                item
                for item in presentation.slides[0].shapes
                if item.name == "ia:element-001"
            )
            actual = {
                "rotation": shape.rotation,
                "line_color": shape.line.color.rgb,
                "vertical_alignment": shape.text_frame.vertical_anchor,
                "wrap": shape.text_frame.word_wrap,
            }

        self.assertEqual(actual["rotation"], 7)
        self.assertEqual(actual["line_color"], RGBColor(0x44, 0x55, 0x66))
        self.assertEqual(actual["vertical_alignment"], MSO_ANCHOR.TOP)
        self.assertFalse(actual["wrap"])

    def test_icon_unknown_fields_have_identical_prebuild_and_compiler_contract(self) -> None:
        from tests.test_build_pptx_from_spec import (
            compile_with_claimed_prebuild,
            make_icon_spec,
        )
        from lib.error_codes import ToolError

        validator_spec = importlib.util.spec_from_file_location(
            "schema_contract_icon_validator",
            SCRIPTS_ROOT / "validate_reconstruction_spec.py",
        )
        self.assertIsNotNone(validator_spec)
        self.assertIsNotNone(validator_spec.loader if validator_spec else None)
        assert validator_spec is not None and validator_spec.loader is not None
        validator = importlib.util.module_from_spec(validator_spec)
        validator_spec.loader.exec_module(validator)
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            spec = make_icon_spec(root)
            spec["modules"]["icons"]["icons"][0].update(
                {"instance": 1, "repeat": False, "pixel_bbox": [1, 2, 3, 4]}
            )
            expected = {
                "code": "UNSUPPORTED_CAPABILITY",
                "path": "modules.icons.icons.status-icon",
                "detail": "unknown fields: instance, pixel_bbox, repeat",
            }

            prebuild = validator.validate_spec(spec, stage="prebuild")
            with self.assertRaises(ToolError) as raised:
                compile_with_claimed_prebuild(root / "compile", spec)

        self.assertIn(expected, prebuild["errors"])
        self.assertEqual(raised.exception.as_dict(), expected)

    def test_supported_icon_record_passes_prebuild_and_compiler(self) -> None:
        from tests.test_build_pptx_from_spec import compile_fixture, make_icon_spec

        validator_spec = importlib.util.spec_from_file_location(
            "schema_contract_supported_icon_validator",
            SCRIPTS_ROOT / "validate_reconstruction_spec.py",
        )
        self.assertIsNotNone(validator_spec)
        self.assertIsNotNone(validator_spec.loader if validator_spec else None)
        assert validator_spec is not None and validator_spec.loader is not None
        validator = importlib.util.module_from_spec(validator_spec)
        validator_spec.loader.exec_module(validator)
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            spec = make_icon_spec(root)

            prebuild = validator.validate_spec(spec, stage="prebuild")
            pptx, report = compile_fixture(root / "compile", spec)
            pptx_exists = pptx.is_file()

        self.assertTrue(prebuild["valid"], prebuild)
        self.assertTrue(pptx_exists)
        self.assertEqual(report["elements"]["status-icon"]["object_type"], "pic")


if __name__ == "__main__":
    unittest.main()
