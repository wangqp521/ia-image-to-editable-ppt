"""End-to-end regressions for schema-v2 prebuild, compile, and validation."""

from __future__ import annotations

import copy
import hashlib
import importlib.util
import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock

from PIL import Image
from pptx import Presentation
from tests.fixture_specs import make_asset_fallback_spec, make_minimal_spec
from tests.test_build_pptx_from_spec import (
    _append_primitive,
    make_icon_spec,
    make_merged_table_spec,
    make_native_list_spec,
    make_picture_spec,
    make_shape_spec,
    make_status_spec,
)
from tests.test_build_report_validation import (
    remove_named_shape_from_pptx,
    replace_embedded_media,
)


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS_ROOT = ROOT / "scripts"
if str(SCRIPTS_ROOT) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_ROOT))

from lib.atomic_write import atomic_write_json
from lib.error_codes import ToolError
from lib.hashing import canonical_json_sha256, file_sha256
from lib.schema_contracts import construct_record


def _load_script(filename: str, module_name: str):
    path = SCRIPTS_ROOT / filename
    module_spec = importlib.util.spec_from_file_location(module_name, path)
    if module_spec is None or module_spec.loader is None:
        raise AssertionError(f"cannot load {path}")
    module = importlib.util.module_from_spec(module_spec)
    module_spec.loader.exec_module(module)
    return module


SPEC_VALIDATOR = _load_script(
    "validate_reconstruction_spec.py", "round_trip_spec_validator"
)
PPTX_VALIDATOR = _load_script("validate_pptx.py", "round_trip_pptx_validator")


def _write_prebuild(
    root: Path, spec: dict
) -> tuple[Path, Path, dict]:
    root.mkdir(parents=True, exist_ok=True)
    spec_path = root / "page-reconstruction.json"
    prebuild_path = root / "prebuild-report.json"
    atomic_write_json(spec_path, spec)
    prebuild = SPEC_VALIDATOR.validate_spec(spec, stage="prebuild")
    if prebuild["valid"] is not True:
        raise AssertionError(f"round-trip prebuild must pass: {prebuild['errors']!r}")
    atomic_write_json(prebuild_path, prebuild)
    return spec_path, prebuild_path, prebuild


def _compile_and_validate(
    root: Path, spec: dict
) -> tuple[Path, Path, dict, dict]:
    spec_path, prebuild_path, _prebuild = _write_prebuild(root, spec)
    pptx_path = root / "page.pptx"
    build_report_path = root / "build-report.json"
    completed = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_ROOT / "build_pptx_from_spec.py"),
            "--spec",
            str(spec_path),
            "--prebuild-report",
            str(prebuild_path),
            "--output",
            str(pptx_path),
            "--build-report",
            str(build_report_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    if completed.returncode != 0:
        raise AssertionError(
            f"compiler CLI failed: stdout={completed.stdout!r}, stderr={completed.stderr!r}"
        )
    build_report = json.loads(build_report_path.read_text(encoding="utf-8"))
    validation = PPTX_VALIDATOR.validate_pptx(
        pptx_path,
        expected_slides=1,
        reconstruction_spec=spec,
        build_report=build_report_path,
    )
    return pptx_path, build_report_path, build_report, validation


def _make_primitives_spec(root: Path) -> dict:
    spec = make_picture_spec(root)
    _append_primitive(
        spec,
        element_id="card",
        kind="shape",
        style={
            "shape_type": "rectangle",
            "fill": {"type": "solid", "color": "#DDEEFF", "opacity": 1.0},
            "line": {
                "color": "#264653",
                "width": 12700,
                "dash": "solid",
                "opacity": 1.0,
            },
            "effects": "none",
            "rotation": 0,
        },
        content={},
        source_bbox=[400, 120, 240, 120],
        slide_bbox=[3048000, 914400, 1828800, 914400],
    )
    _append_primitive(
        spec,
        element_id="arrow",
        kind="line",
        style={
            "line": {
                "color": "#E76F51",
                "width": 25400,
                "dash": "dashDot",
                "opacity": 0.5,
            },
            "head_arrow": "oval",
            "tail_arrow": "triangle",
            "rotation": 0,
        },
        content={},
        source_bbox=[80, 400, 240, 120],
        slide_bbox=[609600, 3048000, 1828800, 914400],
    )
    return spec


class RoundTripSuccessTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self.temporary.cleanup)
        self.root = Path(self.temporary.name)

    def _assert_common_success(
        self, root: Path, spec: dict
    ) -> tuple[dict, dict]:
        pptx, report_path, report, validation = _compile_and_validate(root, spec)

        self.assertTrue(validation["valid"], validation["errors"])
        self.assertEqual(validation["errors"], [])
        self.assertGreater(validation["build_report_objects_checked"], 0)
        self.assertGreater(validation["representation_facts_checked"], 0)
        self.assertEqual(report["pptx_sha256"], file_sha256(pptx))
        self.assertEqual(validation["pptx_sha256"], file_sha256(pptx))
        self.assertEqual(
            json.loads(report_path.read_text(encoding="utf-8")), report
        )
        return report, validation

    def test_text_and_native_bullets_complete_round_trip(self) -> None:
        _report, validation = self._assert_common_success(
            self.root, make_native_list_spec(self.root)
        )

        self.assertGreater(validation["text_shape_count"], 0)
        self.assertGreater(validation["native_list_paragraphs"], 0)
        self.assertGreater(validation["native_list_contracts_checked"], 0)

    def test_shape_line_and_picture_complete_round_trip(self) -> None:
        _report, validation = self._assert_common_success(
            self.root, _make_primitives_spec(self.root)
        )

        self.assertGreater(validation["editable_object_count"], 0)
        self.assertGreater(validation["picture_count"], 0)
        self.assertGreater(validation["asset_fallbacks_checked"], 0)

    def test_contract_factories_build_dense_text_shape_icon_regions(self) -> None:
        spec = make_icon_spec(self.root)
        evidence = str((self.root / "representation-evidence.json").resolve())
        shape = construct_record(
            "Element",
            element_id="dense-card",
            kind="shape",
            source_bbox=[400, 120, 240, 120],
            slide_bbox=[3048000, 914400, 1828800, 914400],
            layer=1,
            editable=True,
            confidence="high",
            style={
                "shape_type": "rectangle",
                "fill": {"type": "solid", "color": "#DDEEFF", "opacity": 1},
                "line": {
                    "color": "#264653",
                    "width": 12700,
                    "dash": "solid",
                    "opacity": 1,
                },
                "effects": "none",
                "rotation": 0,
            },
            content={},
        )
        text = "Second editable"
        text_bbox = [609600, 4572000, 2438400, 457200]
        second_text = construct_record(
            "Element",
            element_id="dense-label",
            kind="text",
            source_bbox=[80, 600, 320, 60],
            slide_bbox=text_bbox,
            layer=3,
            editable=True,
            confidence="high",
            style={"fill": "noFill"},
            content={"text": text},
        )
        typography = construct_record(
            "TypographyItem",
            element_id="dense-label",
            text=text,
            source_font_guess="unknown",
            selected_font="Noto Sans CJK SC",
            fallback_reason="source_font_uncertain",
            fallback_trace=None,
            runs=[
                construct_record(
                    "TextRun",
                    start=0,
                    end=len(text),
                    font_size=18,
                    font_weight=400,
                    color="#112233",
                    letter_spacing=0,
                    italic=False,
                    underline=False,
                    strike=False,
                    baseline=0,
                )
            ],
            paragraphs=[
                construct_record(
                    "Paragraph",
                    start=0,
                    end=len(text),
                    alignment="left",
                    line_spacing=1,
                    space_before=0,
                    space_after=0,
                    indent=0,
                    list={"is_list": False, "level": 0, "bullet": None},
                )
            ],
            text_box=construct_record(
                "TextBox",
                x=text_bbox[0],
                y=text_bbox[1],
                w=text_bbox[2],
                h=text_bbox[3],
                margins={"left": 0, "right": 0, "top": 0, "bottom": 0},
                alignment="left",
                vertical_alignment="top",
                wrap=False,
                overflow=False,
                soft_breaks=[],
                paragraph_breaks=[],
            ),
            internal_font_declaration="Noto Sans CJK SC",
            font_declaration_verified=False,
        )
        spec["elements"].extend([shape, second_text])
        spec["reading_order"].extend(["dense-card", "dense-label"])
        spec["modules"]["typography"]["items"].append(typography)
        spec["modules"]["representation_plan"]["items"].extend(
            [
                construct_record(
                    "RepresentationItem",
                    source_fact_id="fact-dense-card",
                    semantic_role="shape",
                    source_bbox=[400, 120, 240, 120],
                    required=True,
                    selected_mode="native",
                    required_editability="full",
                    fallback_policy="forbid",
                    bound_element_ids=["dense-card"],
                    reason="card is a native editable shape",
                    coverage_status="covered",
                    evidence=[evidence],
                ),
                construct_record(
                    "RepresentationItem",
                    source_fact_id="fact-dense-label",
                    semantic_role="label",
                    source_bbox=[80, 600, 320, 60],
                    required=True,
                    selected_mode="native",
                    required_editability="full",
                    fallback_policy="forbid",
                    bound_element_ids=["dense-label"],
                    reason="label is native editable text",
                    coverage_status="covered",
                    evidence=[evidence],
                ),
            ]
        )
        spec["regions"] = [
            construct_record(
                "Region",
                region_id="upper",
                source_bbox=[0, 0, 1600, 500],
                slide_bbox=[0, 0, 12192000, 3810000],
                layer=0,
                padding={"left": 0, "right": 0, "top": 0, "bottom": 0},
                element_ids=[
                    "background-base",
                    "element-001",
                    "status-icon",
                    "dense-card",
                ],
            ),
            construct_record(
                "Region",
                region_id="lower",
                source_bbox=[0, 500, 1600, 400],
                slide_bbox=[0, 3810000, 12192000, 3048000],
                layer=0,
                padding={"left": 0, "right": 0, "top": 0, "bottom": 0},
                element_ids=["dense-label"],
            ),
        ]

        report, validation = self._assert_common_success(self.root, spec)

        self.assertEqual(len(report["elements"]), 5)
        self.assertEqual(validation["representation_facts_checked"], 4)
        self.assertGreaterEqual(validation["text_shape_count"], 2)
        self.assertGreater(validation["picture_count"], 0)

    def test_webp_picture_preserves_original_media_through_round_trip(self) -> None:
        spec = make_picture_spec(self.root)
        webp = self.root / "photo.webp"
        Image.new("RGB", (40, 20), "#336699").save(
            webp, format="WEBP", lossless=True
        )
        original = webp.read_bytes()
        asset = spec["elements"][-1]["content"]["asset"]
        asset["path"] = str(webp.resolve())
        asset["asset_sha256"] = hashlib.sha256(original).hexdigest()

        _report, validation = self._assert_common_success(self.root, spec)

        picture = next(
            item
            for item in validation["picture_objects"]
            if item["object_name"] == "ia:photo"
        )
        self.assertEqual(picture["media_sha256"], asset["asset_sha256"])
        self.assertTrue(picture["media_basename"].endswith(".webp"))
        with zipfile.ZipFile(self.root / "page.pptx") as archive:
            self.assertEqual(archive.read(picture["media_part"]), original)

    def test_native_merged_table_completes_round_trip(self) -> None:
        _report, validation = self._assert_common_success(
            self.root, make_merged_table_spec(self.root)
        )

        self.assertGreater(validation["graphic_frame_count"], 0)

    def test_explicit_and_repeat_multipart_complete_round_trip(self) -> None:
        for repeat in (False, True):
            with self.subTest(repeat=repeat):
                root = self.root / ("repeat" if repeat else "explicit")
                _report, validation = self._assert_common_success(
                    root, make_status_spec(root, repeat=repeat)
                )

                self.assertGreater(validation["multipart_contracts_checked"], 0)

    def test_labels_only_asset_closes_native_label_and_picture_bindings(self) -> None:
        _report, validation = self._assert_common_success(
            self.root,
            make_asset_fallback_spec(
                self.root, required_editability="labels_only"
            ),
        )

        self.assertGreater(validation["text_shape_count"], 0)
        self.assertEqual(validation["asset_fallbacks_checked"], 1)


class RoundTripFailureTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.addCleanup(self.temporary.cleanup)
        self.root = Path(self.temporary.name)

    def _assert_prebuild_failure(
        self, root: Path, spec: dict, expected_code: str
    ) -> dict:
        before = canonical_json_sha256(spec)
        report = SPEC_VALIDATOR.validate_spec(spec, stage="prebuild")

        self.assertFalse(report["valid"])
        self.assertIn(expected_code, [item["code"] for item in report["errors"]])
        self.assertEqual(canonical_json_sha256(spec), before)
        self.assertFalse((root / "page.pptx").exists())
        self.assertFalse((root / "build-report.json").exists())
        return report

    def test_unknown_element_field_fails_prebuild_without_outputs(self) -> None:
        spec = make_minimal_spec(self.root)
        spec["elements"][0]["future_field"] = True

        self._assert_prebuild_failure(
            self.root, spec, "UNSUPPORTED_CAPABILITY"
        )

    def test_unsupported_rect_fails_prebuild_without_outputs(self) -> None:
        spec = make_shape_spec(self.root, "rectangle")
        spec["elements"][-1]["style"]["shape_type"] = "rect"

        self._assert_prebuild_failure(
            self.root, spec, "UNSUPPORTED_CAPABILITY"
        )

    def test_stale_prebuild_fails_compile_without_partial_pair(self) -> None:
        spec = make_minimal_spec(self.root)
        spec_path, prebuild_path, _prebuild = _write_prebuild(self.root, spec)
        changed = copy.deepcopy(spec)
        changed["modules"]["representation_plan"]["items"][0][
            "reason"
        ] = "changed after prebuild"
        atomic_write_json(spec_path, changed)
        spec_hash_before = file_sha256(spec_path)
        prebuild_hash_before = file_sha256(prebuild_path)

        completed = subprocess.run(
            [
                sys.executable,
                str(SCRIPTS_ROOT / "build_pptx_from_spec.py"),
                "--spec",
                str(spec_path),
                "--prebuild-report",
                str(prebuild_path),
                "--output",
                str(self.root / "page.pptx"),
                "--build-report",
                str(self.root / "build-report.json"),
            ],
            check=False,
            capture_output=True,
            text=True,
        )
        failure = json.loads(completed.stdout)

        self.assertEqual(completed.returncode, 2)
        self.assertEqual(completed.stderr, "")
        self.assertEqual(failure["errors"][0]["code"], "SPEC_HASH_MISMATCH")
        self.assertEqual(file_sha256(spec_path), spec_hash_before)
        self.assertEqual(file_sha256(prebuild_path), prebuild_hash_before)
        self.assertFalse((self.root / "page.pptx").exists())
        self.assertFalse((self.root / "build-report.json").exists())

    def test_forbidden_asset_fails_prebuild_without_outputs(self) -> None:
        spec = make_asset_fallback_spec(
            self.root, required_editability="full"
        )

        self._assert_prebuild_failure(
            self.root, spec, "REPRESENTATION_FALLBACK_FORBIDDEN"
        )

    def test_near_full_page_asset_fails_prebuild_without_outputs(self) -> None:
        spec = make_asset_fallback_spec(
            self.root, required_editability="none"
        )
        bbox = [0, 0, 1599, 899]
        fact = spec["modules"]["representation_plan"]["items"][-1]
        fact["source_bbox"] = list(bbox)
        spec["elements"][-1]["source_bbox"] = list(bbox)
        spec["elements"][-1]["slide_bbox"] = [
            0,
            0,
            12_184_380,
            6_850_380,
        ]

        self._assert_prebuild_failure(
            self.root, spec, "REPRESENTATION_FALLBACK_FORBIDDEN"
        )

    def test_disguised_media_format_fails_prebuild_without_outputs(self) -> None:
        spec = make_picture_spec(self.root)
        asset = spec["elements"][-1]["content"]["asset"]
        path = Path(asset["path"])
        Image.new("RGB", (40, 20), "#336699").save(path, format="GIF")
        asset["asset_sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()

        self._assert_prebuild_failure(
            self.root, spec, "UNSUPPORTED_CAPABILITY"
        )

    def test_compiler_cli_rejects_json_escaped_nul_asset_paths(self) -> None:
        for name, raw_path in (
            ("leading", "\x00/private/tmp/bad.png"),
            ("middle", "/private/tmp/bad\x00name.png"),
            ("trailing", "/private/tmp/bad.png\x00"),
            ("multiple", "/private/\x00tmp/bad\x00name.png"),
        ):
            with self.subTest(name=name):
                root = self.root / name
                root.mkdir(parents=True)
                spec = make_picture_spec(root)
                spec["elements"][-1]["content"]["asset"]["path"] = raw_path
                spec_path = root / "page-reconstruction.json"
                prebuild_path = root / "prebuild-report.json"
                output = root / "page.pptx"
                build_report = root / "build-report.json"
                atomic_write_json(spec_path, spec)
                self.assertIn("\\u0000", spec_path.read_text(encoding="utf-8"))
                atomic_write_json(
                    prebuild_path,
                    {
                        "valid": True,
                        "stage": "prebuild",
                        "errors": [],
                        "spec_sha256": canonical_json_sha256(spec),
                    },
                )

                completed = subprocess.run(
                    [
                        sys.executable,
                        str(SCRIPTS_ROOT / "build_pptx_from_spec.py"),
                        "--spec",
                        str(spec_path),
                        "--prebuild-report",
                        str(prebuild_path),
                        "--output",
                        str(output),
                        "--build-report",
                        str(build_report),
                    ],
                    check=False,
                    capture_output=True,
                    text=True,
                )

                self.assertEqual(completed.returncode, 2, completed.stderr)
                self.assertEqual(completed.stderr, "")
                failure = json.loads(completed.stdout)
                self.assertEqual(
                    failure["errors"][0]["code"], "UNSUPPORTED_CAPABILITY"
                )
                self.assertEqual(
                    failure["errors"][0]["path"],
                    "elements.photo.content.asset.path",
                )
                self.assertEqual(
                    failure["errors"][0]["capability"],
                    "picture.asset.local_hash",
                )
                self.assertEqual(
                    failure["errors"][0]["detail"],
                    "asset path must not contain NUL characters",
                )
                self.assertFalse(output.exists())
                self.assertFalse(build_report.exists())

    def test_prebuild_rejects_nul_asset_paths_without_native_exception(self) -> None:
        for name, raw_path in (
            ("leading", "\x00/private/tmp/bad.png"),
            ("middle", "/private/tmp/bad\x00name.png"),
            ("trailing", "/private/tmp/bad.png\x00"),
            ("multiple", "/private/\x00tmp/bad\x00name.png"),
        ):
            with self.subTest(name=name):
                root = self.root / f"prebuild-{name}"
                spec = make_picture_spec(root)
                spec["elements"][-1]["content"]["asset"]["path"] = raw_path
                before = canonical_json_sha256(spec)
                raised = None
                report = None
                try:
                    report = SPEC_VALIDATOR.validate_spec(
                        spec, stage="prebuild"
                    )
                except Exception as exc:  # assertion below owns public behavior
                    raised = exc

                self.assertIsNone(
                    raised, "prebuild must not leak a native path exception"
                )
                assert report is not None
                error = next(
                    item
                    for item in report["errors"]
                    if item["code"] == "UNSUPPORTED_CAPABILITY"
                )
                self.assertEqual(
                    error["path"], "elements.photo.content.asset.path"
                )
                self.assertEqual(
                    error["capability"], "picture.asset.local_hash"
                )
                self.assertEqual(canonical_json_sha256(spec), before)
                self.assertFalse((root / "page.pptx").exists())
                self.assertFalse((root / "build-report.json").exists())

    def test_webp_read_preserves_existing_tool_error(self) -> None:
        from pptx_builder.pictures import _add_webp_picture

        sentinel = ToolError(
            "SENTINEL",
            "sentinel.path",
            "sentinel detail",
            "sentinel.capability",
        )
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        asset_path = self.root / "sentinel.webp"
        with mock.patch.object(
            Path, "read_bytes", side_effect=sentinel
        ):
            raised = None
            try:
                _add_webp_picture(
                    slide,
                    asset_path,
                    0,
                    0,
                    914_400,
                    914_400,
                    "elements.photo.content.asset.path",
                )
            except ToolError as exc:
                raised = exc

        self.assertIsNotNone(raised)
        assert raised is not None
        self.assertEqual(raised.as_dict(), sentinel.as_dict())

    def test_missing_multipart_part_fails_validation_without_mutation(self) -> None:
        spec = make_status_spec(self.root, repeat=False)
        pptx, _report_path, report, _validation = _compile_and_validate(
            self.root, spec
        )
        remove_named_shape_from_pptx(pptx, "ia:status:segment-1")
        pptx_before = file_sha256(pptx)
        spec_before = canonical_json_sha256(spec)
        report_before = canonical_json_sha256(report)

        result = PPTX_VALIDATOR.validate_pptx(pptx, 1, spec, report)

        self.assertFalse(result["valid"])
        self.assertIn("BUILD_OUTPUT_INCOMPLETE", result["errors"])
        self.assertEqual(file_sha256(pptx), pptx_before)
        self.assertEqual(canonical_json_sha256(spec), spec_before)
        self.assertEqual(canonical_json_sha256(report), report_before)

    def test_tampered_build_report_fails_validation_without_mutation(self) -> None:
        spec = make_native_list_spec(self.root)
        pptx, _report_path, report, _validation = _compile_and_validate(
            self.root, spec
        )
        tampered = copy.deepcopy(report)
        tampered["schema_sha256"] = "0" * 64
        pptx_before = file_sha256(pptx)
        spec_before = canonical_json_sha256(spec)
        report_before = canonical_json_sha256(tampered)

        result = PPTX_VALIDATOR.validate_pptx(pptx, 1, spec, tampered)

        self.assertFalse(result["valid"])
        self.assertIn("BUILD_REPORT_MISMATCH", result["errors"])
        self.assertEqual(file_sha256(pptx), pptx_before)
        self.assertEqual(canonical_json_sha256(spec), spec_before)
        self.assertEqual(canonical_json_sha256(tampered), report_before)

    def test_embedded_media_hash_fails_validation_without_extra_mutation(self) -> None:
        spec = make_picture_spec(self.root)
        pptx, _report_path, report, _validation = _compile_and_validate(
            self.root, spec
        )
        replace_embedded_media(pptx, b"different")
        pptx_before = file_sha256(pptx)
        spec_before = canonical_json_sha256(spec)
        report_before = canonical_json_sha256(report)

        result = PPTX_VALIDATOR.validate_pptx(pptx, 1, spec, report)

        self.assertFalse(result["valid"])
        self.assertIn("ASSET_HASH_MISMATCH", result["errors"])
        self.assertEqual(file_sha256(pptx), pptx_before)
        self.assertEqual(canonical_json_sha256(spec), spec_before)
        self.assertEqual(canonical_json_sha256(report), report_before)


if __name__ == "__main__":
    unittest.main()
