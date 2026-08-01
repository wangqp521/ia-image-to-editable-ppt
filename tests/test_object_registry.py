"""Public contracts for the schema v2 renderer foundation."""

from __future__ import annotations

import importlib
import sys
import tempfile
import unittest
from dataclasses import FrozenInstanceError
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE

from tests.fixture_specs import make_minimal_spec


SCRIPTS_ROOT = Path(__file__).resolve().parents[1] / "scripts"


class FakeShape:
    left = 10
    top = 20
    width = 30
    height = 40
    rotation = 0
    name = ""


class FakeRenderer:
    kind = "shape"
    supported_fields = frozenset({"shape_type"})
    supported_values = {"shape_type": frozenset({"rectangle"})}
    required_fields = frozenset({"shape_type"})
    capability_ids = frozenset({"shape.rectangle"})

    def validate_contract(self, element, context):
        return []

    def render(self, element, context):
        return None


class ObjectRegistryTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary_directory = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary_directory.name)
        common, _, _ = self._api()
        self.registered_renderers = dict(common.RENDERERS)

    def tearDown(self) -> None:
        common, _, _ = self._api()
        common.RENDERERS.clear()
        common.RENDERERS.update(self.registered_renderers)
        self.temporary_directory.cleanup()

    def _api(self):
        """Import the target in the test body so absence is an assertion failure."""
        scripts_root = str(SCRIPTS_ROOT)
        if scripts_root not in sys.path:
            sys.path.insert(0, scripts_root)
        try:
            common = importlib.import_module("pptx_builder.common")
            registry = importlib.import_module("pptx_builder.registry")
            ooxml = importlib.import_module("pptx_builder.ooxml")
        except ModuleNotFoundError as exc:
            if exc.name is not None and exc.name.startswith("pptx_builder"):
                self.fail("object registry is not implemented")
            raise
        return common, registry, ooxml

    def test_register_assigns_stable_names(self) -> None:
        _, registry_module, _ = self._api()
        registry = registry_module.ObjectRegistry()

        registry.register(
            "status",
            FakeShape(),
            "sp",
            semantic_kind="status",
            selected_mode="composite",
            part_id="fill",
        )

        self.assertEqual(registry.records["status"][0]["ooxml_name"], "ia:status:fill")

    def test_register_accepts_real_python_pptx_emu_geometry(self) -> None:
        _, registry_module, _ = self._api()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(10, 20, 30, 40)
        registry = registry_module.ObjectRegistry()

        error = None
        try:
            registry.register(
                "title",
                shape,
                "sp",
                semantic_kind="text",
                selected_mode="native",
            )
        except registry_module.ToolError as exc:
            error = exc

        self.assertIsNone(error, f"real python-pptx geometry was rejected: {error}")

        self.assertEqual(
            registry.records["title"][0]["ooxml_name"], "ia:title"
        )

    def test_render_context_carries_a_read_only_typography_index(self) -> None:
        common, registry_module, _ = self._api()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        typography = {"title": {"text": "标题"}}
        try:
            context = common.RenderContext(
                slide=slide,
                spec={},
                representation_modes={"title": "native"},
                typography=typography,
                registry=registry_module.ObjectRegistry(),
            )
        except TypeError as exc:
            self.fail(f"RenderContext typography index is not implemented: {exc}")

        self.assertEqual(dict(context.typography), typography)
        with self.assertRaises(TypeError):
            context.typography["other"] = {"text": "其他"}
        with self.assertRaises(FrozenInstanceError):
            context.typography = {}

    def test_duplicate_part_fails_closed(self) -> None:
        _, registry_module, _ = self._api()
        registry = registry_module.ObjectRegistry()
        registry.register(
            "status",
            FakeShape(),
            "sp",
            semantic_kind="status",
            selected_mode="composite",
            part_id="fill",
        )

        with self.assertRaises(registry_module.ToolError) as raised:
            registry.register(
                "status",
                FakeShape(),
                "sp",
                semantic_kind="status",
                selected_mode="composite",
                part_id="fill",
            )

        self.assertEqual(raised.exception.code, "BUILD_OBJECT_NAME_COLLISION")

    def test_finalize_rejects_unregistered_element(self) -> None:
        _, registry_module, _ = self._api()
        spec = make_minimal_spec(self.root)

        with self.assertRaises(registry_module.ToolError) as raised:
            registry_module.ObjectRegistry().finalize(
                spec,
                {"element-001": "native", "background-base": "native"},
            )

        self.assertEqual(raised.exception.code, "BUILD_OUTPUT_INCOMPLETE")

    def test_finalize_returns_sorted_json_only_element_report(self) -> None:
        _, registry_module, _ = self._api()
        spec = make_minimal_spec(self.root)
        first = dict(spec["elements"][0])
        first.update({"element_id": "a", "slide_bbox": [10, 20, 30, 40]})
        second = dict(spec["elements"][0])
        second.update({"element_id": "z", "slide_bbox": [50, 60, 70, 80]})
        spec["elements"] = [second, first]
        for item, element_id in zip(
            spec["modules"]["representation_plan"]["items"], ("z",)
        ):
            item["bound_element_ids"] = [element_id]
        copied_plan_item = dict(spec["modules"]["representation_plan"]["items"][0])
        copied_plan_item.update({"source_fact_id": "fact-a", "bound_element_ids": ["a"]})
        spec["modules"]["representation_plan"]["items"].append(copied_plan_item)

        registry = registry_module.ObjectRegistry()
        last_shape = FakeShape()
        last_shape.left, last_shape.top, last_shape.width, last_shape.height = 50, 60, 70, 80
        registry.register("z", last_shape, "sp", semantic_kind="text", selected_mode="native")
        first_shape = FakeShape()
        first_shape.left, first_shape.top, first_shape.width, first_shape.height = 10, 20, 30, 40
        registry.register("a", first_shape, "sp", semantic_kind="text", selected_mode="native")

        report = registry.finalize(spec, {"a": "native", "z": "native"})

        self.assertEqual(list(report), ["a", "z"])
        self.assertEqual(
            report["a"],
            {
                "semantic_kind": "text",
                "selected_mode": "native",
                "object_type": "sp",
                "objects": [
                    {
                        "ooxml_name": "ia:a",
                        "object_type": "sp",
                        "bbox": [10, 20, 30, 40],
                        "rotation": 0,
                        "part_id": None,
                        "media_sha256": None,
                        "text_summary": None,
                        "font_declarations": [],
                    }
                ],
            },
        )
        self._assert_json_only(report)

    def test_finalize_rejects_wrong_actual_bbox(self) -> None:
        _, registry_module, _ = self._api()
        spec = make_minimal_spec(self.root)
        registry = registry_module.ObjectRegistry()
        registry.register(
            "element-001", FakeShape(), "sp", semantic_kind="text", selected_mode="native"
        )

        with self.assertRaises(registry_module.ToolError) as raised:
            registry.finalize(
                spec,
                {"element-001": "native", "background-base": "native"},
            )

        self.assertEqual(raised.exception.code, "BUILD_OUTPUT_INCOMPLETE")

    def test_finalize_rejects_multipart_objects_without_parent_union(self) -> None:
        _, registry_module, _ = self._api()
        spec = make_minimal_spec(self.root)
        element = spec["elements"][0]
        element.update(
            {
                "kind": "status",
                "slide_bbox": [0, 0, 100, 100],
                "content": {
                    "part_defaults": {"style": {}, "content": {}},
                    "parts": [
                        {"part_id": "left", "slide_bbox": [0, 0, 50, 50]},
                        {"part_id": "right", "slide_bbox": [50, 0, 50, 50]},
                    ],
                },
            }
        )
        spec["modules"]["representation_plan"]["items"][0]["semantic_role"] = "status"
        registry = registry_module.ObjectRegistry()
        for part_id, bbox in (("left", [0, 0, 50, 50]), ("right", [50, 0, 50, 50])):
            shape = FakeShape()
            shape.left, shape.top, shape.width, shape.height = bbox
            registry.register(
                "element-001",
                shape,
                "sp",
                semantic_kind="status",
                selected_mode="native",
                part_id=part_id,
            )

        with self.assertRaises(registry_module.ToolError) as raised:
            registry.finalize(
                spec,
                {"element-001": "native", "background-base": "native"},
            )

        self.assertEqual(raised.exception.code, "BUILD_OUTPUT_INCOMPLETE")

    def test_register_renderer_rejects_metadata_outside_capability_registry(self) -> None:
        common, _, _ = self._api()

        invalid = FakeRenderer()
        invalid.supported_values = {"shape_type": frozenset({"star5"})}
        invalid.capability_ids = frozenset({"shape.star5"})

        with self.assertRaises(ValueError):
            common.register_renderer("shape", invalid)

    def test_register_renderer_rejects_capability_value_mismatch(self) -> None:
        common, _, _ = self._api()

        invalid = FakeRenderer()
        invalid.capability_ids = frozenset({"shape.ellipse"})

        with self.assertRaises(ValueError):
            common.register_renderer("shape", invalid)

    def test_register_renderer_rejects_values_without_declared_field(self) -> None:
        common, _, _ = self._api()
        invalid = FakeRenderer()
        invalid.supported_fields = frozenset()
        invalid.required_fields = frozenset()

        with self.assertRaises(ValueError):
            common.register_renderer("shape", invalid)

    def test_register_renderer_rejects_non_contract_metadata_containers(self) -> None:
        common, _, _ = self._api()
        cases = (
            ("supported_fields", {"shape_type"}),
            ("supported_values", [("shape_type", frozenset({"rectangle"}))]),
            ("required_fields", {"shape_type"}),
            ("capability_ids", {"shape.rectangle"}),
        )
        for field, value in cases:
            with self.subTest(field=field):
                common.RENDERERS.clear()
                invalid = FakeRenderer()
                setattr(invalid, field, value)

                with self.assertRaises(ValueError):
                    common.register_renderer("shape", invalid)

    def test_register_renderer_requires_explicit_unique_static_kind(self) -> None:
        common, _, _ = self._api()
        common.RENDERERS.clear()
        renderer = FakeRenderer()

        common.register_renderer("shape", renderer)

        self.assertIs(common.RENDERERS["shape"], renderer)
        with self.assertRaises(RuntimeError):
            common.register_renderer("shape", renderer)

    def test_ooxml_helpers_apply_native_properties(self) -> None:
        _, _, ooxml = self._api()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(10, 20, 100, 100)
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.text = "项目"
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 10, 20, 100, 100)
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, 10, 20, 100, 100)
        cell = slide.shapes.add_table(1, 1, 10, 20, 100, 100).table.cell(0, 0)

        ooxml.set_native_bullet(
            paragraph,
            {"bullet_type": "char", "bullet": "•", "level": 0, "indent": -100},
            "elements.a",
        )
        ooxml.set_round_rect_adjustment(shape, [0.25], "elements.a.style.adjustments")
        ooxml.set_line_arrowheads(line, {"head_arrow": "triangle", "tail_arrow": "none"}, "elements.a.style")
        ooxml.set_table_cell_border(cell._tc.tcPr, "left", {"color": "FF0000", "width": 12700}, "elements.a.style")
        ooxml.neutralize_shape_effects(shape)

        self.assertEqual(paragraph._p.pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar").get("char"), "•")
        self.assertAlmostEqual(shape.adjustments[0], 0.25)
        self.assertIsNotNone(line._element.spPr.ln.find("{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd"))
        self.assertIsNotNone(cell._tc.tcPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}lnL"))
        self.assertIsNone(shape._element.spPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst"))

    def test_line_arrowheads_reject_unknown_contract_fields_and_values(self) -> None:
        _, _, ooxml = self._api()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, 10, 20, 100, 100)
        cases = (
            {"head_arrow": "not-a-drawingml-arrow"},
            {"head_arrow": "triangle", "unknown": "value"},
        )
        for contract in cases:
            with self.subTest(contract=contract), self.assertRaises(ooxml.ToolError) as raised:
                ooxml.set_line_arrowheads(line, contract, "elements.a.style")

            self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")

    def test_table_border_helper_rejects_malformed_contracts(self) -> None:
        _, _, ooxml = self._api()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        cell = slide.shapes.add_table(1, 1, 10, 20, 100, 100).table.cell(0, 0)
        cases = (
            {"color": "#GG0000", "width": 12700},
            {"color": "#FF0000", "width": 20_116_801},
            {"color": "#FF0000", "width": 12700, "rogue": True},
        )

        for contract in cases:
            with self.subTest(contract=contract):
                with self.assertRaises(ooxml.ToolError) as raised:
                    ooxml.set_table_cell_border(
                        cell._tc.tcPr,
                        "left",
                        contract,
                        "elements.a.content.cells[0].borders.left",
                    )

                self.assertEqual(raised.exception.code, "UNSUPPORTED_CAPABILITY")

    def _assert_json_only(self, value) -> None:
        if isinstance(value, dict):
            for key, nested in value.items():
                self.assertIsInstance(key, str)
                self._assert_json_only(nested)
        elif isinstance(value, list):
            for nested in value:
                self._assert_json_only(nested)
        else:
            self.assertTrue(value is None or type(value) in {str, int, float, bool})


if __name__ == "__main__":
    unittest.main()
