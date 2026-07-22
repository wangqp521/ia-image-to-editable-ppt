#!/usr/bin/env python3
"""Build one editable 16:9 PPTX page from a current schema-v2 specification."""

from __future__ import annotations

import argparse
import importlib.metadata
import json
from pathlib import Path
from typing import Any

from pptx import Presentation

from lib.atomic_write import atomic_write_json
from lib.error_codes import ToolError
from lib.hashing import canonical_json_sha256, file_sha256
from lib.schema_io import load_schema_v2
from pptx_builder.common import ObjectRegistry, rgb
from pptx_builder.pictures import add_picture_element
from pptx_builder.shapes import add_line_element, add_shape_element
from pptx_builder.tables import add_multipart_element, add_table_element
from pptx_builder.text import add_text_element


SUPPORTED_KINDS = {"text", "shape", "line", "table", "matrix", "status", "picture", "icon"}


def require_current_prebuild(spec: dict[str, Any], report: dict[str, Any]) -> None:
    if report.get("valid") is not True or report.get("stage") != "prebuild":
        raise ToolError("SPEC_INVALID", "prebuild_report", "passing prebuild report required")
    if report.get("spec_sha256") != canonical_json_sha256(spec):
        raise ToolError(
            "SPEC_HASH_MISMATCH", "prebuild_report.spec_sha256", "prebuild report is stale"
        )


def _typography_map(spec: dict[str, Any]) -> dict[str, dict[str, Any]]:
    items = spec.get("modules", {}).get("typography", {}).get("items", [])
    return {
        item["element_id"]: item
        for item in items
        if isinstance(item, dict) and isinstance(item.get("element_id"), str)
    }


def _icon_map(spec: dict[str, Any]) -> dict[str, dict[str, Any]]:
    items = spec.get("modules", {}).get("icons", {}).get("icons", [])
    return {
        item["element_id"]: item
        for item in items
        if isinstance(item, dict) and isinstance(item.get("element_id"), str)
    }


def build_single_page(
    spec: dict[str, Any], prebuild_report: dict[str, Any], output: Path | str
) -> dict[str, Any]:
    require_current_prebuild(spec, prebuild_report)
    elements = spec.get("elements")
    reading_order = spec.get("reading_order")
    if not isinstance(elements, list) or not isinstance(reading_order, list):
        raise ToolError("MISSING_REQUIRED_FIELD", "elements", "elements and reading_order required")
    by_id = {element.get("element_id"): element for element in elements if isinstance(element, dict)}
    if set(reading_order) != set(by_id):
        raise ToolError("SPEC_INVALID", "reading_order", "must cover every element")

    prs = Presentation()
    canvas = spec["canvas"]
    prs.slide_width = int(canvas["slide_size_emu"][0])
    prs.slide_height = int(canvas["slide_size_emu"][1])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = canvas.get("background", "#FFFFFF")
    if isinstance(background, str) and background.startswith("#"):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(background, "canvas.background")
    registry = ObjectRegistry()
    typography = _typography_map(spec)
    icons = _icon_map(spec)
    reading_index = {element_id: index for index, element_id in enumerate(reading_order)}
    ordered = sorted(elements, key=lambda item: (item.get("layer", 0), reading_index[item["element_id"]]))
    for element in ordered:
        element_id = element["element_id"]
        kind = element.get("kind")
        if kind not in SUPPORTED_KINDS:
            raise ToolError("UNSUPPORTED_FEATURE", f"elements.{element_id}.kind", str(kind))
        if kind == "text":
            if element_id not in typography:
                raise ToolError("MISSING_REQUIRED_FIELD", f"modules.typography.{element_id}", "text contract required")
            add_text_element(slide, element, typography[element_id], registry)
        elif kind == "shape":
            add_shape_element(slide, element, registry)
        elif kind == "line":
            add_line_element(slide, element, registry)
        elif kind == "table":
            add_table_element(slide, element, registry)
        elif kind in {"matrix", "status"}:
            add_multipart_element(slide, element, registry)
        elif kind == "picture":
            asset = element.get("content", {}).get("asset")
            if not isinstance(asset, dict):
                raise ToolError("MISSING_REQUIRED_FIELD", f"elements.{element_id}.content.asset", "asset required")
            add_picture_element(slide, element, asset, registry)
        elif kind == "icon":
            if element_id not in icons:
                raise ToolError("MISSING_REQUIRED_FIELD", f"modules.icons.{element_id}", "icon asset required")
            add_picture_element(slide, element, icons[element_id], registry)

    destination = Path(output).expanduser().resolve()
    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(destination)
    script_root = Path(__file__).resolve().parent
    builder_files = [Path(__file__), *sorted((script_root / "pptx_builder").glob("*.py"))]
    report = {
        "valid": True,
        "schema_version": 2,
        "schema_sha256": canonical_json_sha256(spec),
        "builder_sha256": canonical_json_sha256(
            {str(path.name): file_sha256(path) for path in builder_files}
        ),
        "pptx_path": str(destination),
        "pptx_sha256": file_sha256(destination),
        "environment": {
            "python-pptx": importlib.metadata.version("python-pptx"),
            "Pillow": importlib.metadata.version("Pillow"),
        },
        "elements": registry.report(),
        "unsupported": [],
        "warnings": [],
    }
    if set(report["elements"]) != set(by_id):
        raise ToolError("BUILD_OUTPUT_INCOMPLETE", "elements", "not every element was registered")
    return report


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--spec", type=Path, required=True)
    parser.add_argument("--prebuild-report", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--build-report", type=Path, required=True)
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)
    try:
        spec = load_schema_v2(args.spec)
        prebuild = json.loads(args.prebuild_report.read_text(encoding="utf-8"))
        report = build_single_page(spec, prebuild, args.output)
        atomic_write_json(args.build_report, report)
        print(json.dumps(report, ensure_ascii=False, indent=2))
        return 0
    except (ToolError, OSError, json.JSONDecodeError) as exc:
        error = exc if isinstance(exc, ToolError) else ToolError("SPEC_INVALID", "$", str(exc))
        report = {"valid": False, "errors": [error.as_dict()]}
        atomic_write_json(args.build_report, report)
        print(json.dumps(report, ensure_ascii=False, indent=2))
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
