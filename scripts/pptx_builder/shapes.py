from __future__ import annotations

from typing import Any

from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

from lib.error_codes import ToolError
from .common import ObjectRegistry, rgb
from .ooxml import set_line_arrowheads, set_round_rect_adjustment


SHAPE_TYPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
}
DASHES = {
    "solid": MSO_LINE_DASH_STYLE.SOLID,
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
}


def _apply_line(line: Any, contract: dict[str, Any], path: str) -> None:
    allowed = {"color", "width_emu", "dash", "begin_arrow", "end_arrow", "transparency"}
    unknown = set(contract) - allowed
    if unknown:
        raise ToolError("UNSUPPORTED_FEATURE", path, f"unsupported line fields: {', '.join(sorted(unknown))}")
    if contract.get("transparency", 0) != 0:
        raise ToolError("UNSUPPORTED_FEATURE", f"{path}.transparency", str(contract["transparency"]))
    line.color.rgb = rgb(contract.get("color", "#000000"), f"{path}.color")
    line.width = int(contract.get("width_emu", 12700))
    dash = contract.get("dash", "solid")
    if dash not in DASHES:
        raise ToolError("UNSUPPORTED_FEATURE", f"{path}.dash", dash)
    line.dash_style = DASHES[dash]
    set_line_arrowheads(line, contract, path)


def add_shape_element(slide: Any, element: dict[str, Any], registry: ObjectRegistry) -> None:
    style = element["style"]
    unknown = set(style) - {"shape_type", "fill", "line", "adjustments"}
    if unknown:
        raise ToolError("UNSUPPORTED_FEATURE", f"elements.{element['element_id']}.style", f"unsupported fields: {', '.join(sorted(unknown))}")
    shape_type = style.get("shape_type", "rectangle")
    if shape_type not in SHAPE_TYPES:
        raise ToolError("UNSUPPORTED_FEATURE", f"elements.{element['element_id']}.shape_type", shape_type)
    shape = slide.shapes.add_shape(SHAPE_TYPES[shape_type], *element["slide_bbox"])
    fill = style.get("fill", "noFill")
    if fill == "noFill":
        shape.fill.background()
    elif isinstance(fill, dict) and isinstance(fill.get("color"), str):
        unknown_fill = set(fill) - {"color", "transparency"}
        if unknown_fill:
            raise ToolError("UNSUPPORTED_FEATURE", f"elements.{element['element_id']}.fill", f"unsupported fields: {', '.join(sorted(unknown_fill))}")
        if fill.get("transparency", 0) != 0:
            raise ToolError("UNSUPPORTED_FEATURE", f"elements.{element['element_id']}.fill.transparency", str(fill["transparency"]))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill["color"])
    else:
        raise ToolError("SPEC_INVALID", f"elements.{element['element_id']}.fill", "invalid fill")
    line = style.get("line")
    if line is None:
        shape.line.fill.background()
    elif isinstance(line, dict):
        _apply_line(shape.line, line, f"elements.{element['element_id']}.line")
    else:
        raise ToolError("SPEC_INVALID", f"elements.{element['element_id']}.line", "invalid line")
    if shape_type == "roundRect":
        set_round_rect_adjustment(
            shape,
            style.get("adjustments"),
            f"elements.{element['element_id']}.adjustments",
        )
    registry.register(element["element_id"], shape, "shape")


def add_line_element(slide: Any, element: dict[str, Any], registry: ObjectRegistry) -> None:
    style = element.get("style", {})
    unknown = set(style) - {"line"}
    if unknown:
        raise ToolError("UNSUPPORTED_FEATURE", f"elements.{element['element_id']}.style", f"unsupported fields: {', '.join(sorted(unknown))}")
    x, y, width, height = element["slide_bbox"]
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + width, y + height)
    contract = style.get("line")
    if not isinstance(contract, dict):
        raise ToolError("MISSING_REQUIRED_FIELD", f"elements.{element['element_id']}.style.line", "line style required")
    _apply_line(line.line, contract, f"elements.{element['element_id']}.line")
    registry.register(element["element_id"], line, "line")
