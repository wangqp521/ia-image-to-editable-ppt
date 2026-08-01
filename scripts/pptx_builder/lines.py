"""Editable line and arrow renderer."""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Emu

from lib.capabilities import ATOMIC_CAPABILITY_METADATA, CANONICAL_VALUES
from lib.error_codes import ContractIssue

from .common import RenderContext, register_renderer
from .ooxml import neutralize_shape_effects, set_line_arrowheads
from .shapes import _issue, _number, apply_stroke, validate_stroke


_LINE_FIELDS = frozenset({"line", "head_arrow", "tail_arrow", "rotation"})


class LineRenderer:
    kind = "line"
    supported_fields = _LINE_FIELDS
    supported_values = {
        "line_dash": CANONICAL_VALUES["line_dash"],
        "line_arrow": CANONICAL_VALUES["line_arrow"],
    }
    required_fields = frozenset({"line"})
    capability_ids = frozenset(
        f"line_dash.{value}" for value in CANONICAL_VALUES["line_dash"]
    ) | frozenset(
        f"line_arrow.{value}" for value in CANONICAL_VALUES["line_arrow"]
    ) | frozenset(
        capability for capability, field in ATOMIC_CAPABILITY_METADATA.items()
        if field in _LINE_FIELDS and capability.startswith("line.")
    )

    def validate_contract(self, element: dict[str, Any], context: RenderContext) -> list[ContractIssue]:
        path = f"elements.{element.get('element_id', '<unknown>')}"
        style = element.get("style", {})
        if "line" not in style:
            return [_issue(f"{path}.style.line", "line contract is required")]
        issues = validate_stroke(style["line"], f"{path}.style.line", "line.stroke")
        if issues:
            return issues
        for field in ("head_arrow", "tail_arrow"):
            value = style.get(field, "none")
            if not isinstance(value, str) or value not in CANONICAL_VALUES["line_arrow"]:
                return [_issue(f"{path}.style.{field}", "unsupported arrowhead", f"line_arrow.{value}")]
        rotation = style.get("rotation", 0)
        if not _number(rotation) or not -360 <= rotation <= 360:
            return [_issue(f"{path}.style.rotation", "rotation must be from -360 to 360 degrees")]
        return []

    def render(self, element: dict[str, Any], context: RenderContext) -> None:
        style = element["style"]
        x, y, width, height = element["slide_bbox"]
        line = context.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(x), Emu(y), Emu(x + width), Emu(y + height),
        )
        neutralize_shape_effects(line)
        apply_stroke(line, style["line"])
        set_line_arrowheads(
            line,
            {"head_arrow": style.get("head_arrow", "none"), "tail_arrow": style.get("tail_arrow", "none")},
            f"elements.{element['element_id']}.style",
        )
        line.rotation = style.get("rotation", 0)
        context.registry.register(
            element["element_id"], line, "cxnSp", semantic_kind="line",
            selected_mode=context.representation_modes[element["element_id"]],
        )


LINE_RENDERER = LineRenderer()
register_renderer("line", LINE_RENDERER)
