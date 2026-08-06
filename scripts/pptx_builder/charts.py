"""Native editable 2D pie, doughnut, column, bar, and line renderer."""

from __future__ import annotations

from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    XL_TICK_LABEL_POSITION,
)
from pptx.util import Emu, Pt

from lib.capabilities import ATOMIC_CAPABILITY_METADATA, CANONICAL_VALUES
from lib.element_contracts import validate_chart_contract
from lib.error_codes import ContractIssue

from .chart_ooxml import (
    set_axis_position,
    set_axis_reverse_order,
    set_chart_area_style,
    set_chart_format_line,
    set_chart_value,
    set_display_blanks_as,
    set_series_smooth,
)
from .common import RenderContext, register_renderer


_PIE_CHART_TYPES = {
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}
_CARTESIAN_CHART_TYPES = {
    ("column", "clustered"): XL_CHART_TYPE.COLUMN_CLUSTERED,
    ("column", "stacked"): XL_CHART_TYPE.COLUMN_STACKED,
    ("column", "percent_stacked"): XL_CHART_TYPE.COLUMN_STACKED_100,
    ("bar", "clustered"): XL_CHART_TYPE.BAR_CLUSTERED,
    ("bar", "stacked"): XL_CHART_TYPE.BAR_STACKED,
    ("bar", "percent_stacked"): XL_CHART_TYPE.BAR_STACKED_100,
    ("line", "standard"): XL_CHART_TYPE.LINE_MARKERS,
}
_LABEL_POSITIONS = {
    "above": XL_LABEL_POSITION.ABOVE,
    "below": XL_LABEL_POSITION.BELOW,
    "best_fit": XL_LABEL_POSITION.BEST_FIT,
    "center": XL_LABEL_POSITION.CENTER,
    "inside_base": XL_LABEL_POSITION.INSIDE_BASE,
    "inside_end": XL_LABEL_POSITION.INSIDE_END,
    "left": XL_LABEL_POSITION.LEFT,
    "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
    "right": XL_LABEL_POSITION.RIGHT,
}
_LEGEND_POSITIONS = {
    "top": XL_LEGEND_POSITION.TOP,
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
}
_MARKER_STYLES = {
    "none": XL_MARKER_STYLE.NONE,
    "circle": XL_MARKER_STYLE.CIRCLE,
    "square": XL_MARKER_STYLE.SQUARE,
    "diamond": XL_MARKER_STYLE.DIAMOND,
    "triangle": XL_MARKER_STYLE.TRIANGLE,
}
_TICK_LABEL_POSITIONS = {
    "next_to_axis": XL_TICK_LABEL_POSITION.NEXT_TO_AXIS,
    "low": XL_TICK_LABEL_POSITION.LOW,
    "high": XL_TICK_LABEL_POSITION.HIGH,
    "none": XL_TICK_LABEL_POSITION.NONE,
}
_CHART_FIELDS = frozenset(
    {
        "chart_type",
        "slices",
        "grouping",
        "categories",
        "series",
        "axes",
        "legend",
        "data_labels",
        "display_blanks_as",
        "first_slice_angle",
        "hole_size",
        "gap_width",
        "overlap",
        "chart_area",
        "plot_area",
    }
)


def _set_font(font: Any, contract: dict[str, Any]) -> None:
    font.name = contract["font_name"]
    font.size = Pt(contract["font_size"])
    font.bold = contract["font_weight"] >= 600
    font.color.rgb = RGBColor.from_string(contract["color"].removeprefix("#"))


def _add_chart(
    element: dict[str, Any],
    context: RenderContext,
    chart_type: Any,
    chart_data: CategoryChartData,
) -> Any:
    x, y, width, height = element["slide_bbox"]
    return context.slide.shapes.add_chart(
        chart_type,
        Emu(x),
        Emu(y),
        Emu(width),
        Emu(height),
        chart_data,
    )


def _render_pie_chart(element: dict[str, Any], context: RenderContext) -> Any:
    content = element["content"]
    style = element["style"]
    slices = content["slices"]
    chart_data = CategoryChartData()
    chart_data.categories = [
        "" if item["category"] is None else item["category"] for item in slices
    ]
    chart_data.add_series("", [item["value"] for item in slices])
    frame = _add_chart(element, context, _PIE_CHART_TYPES[content["chart_type"]], chart_data)
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.vary_by_categories = True
    set_chart_value(plot._element, "c:firstSliceAng", style["first_slice_angle"])
    if content["chart_type"] == "doughnut":
        set_chart_value(plot._element, "c:holeSize", style["hole_size"])

    for point, item in zip(plot.series[0].points, slices, strict=True):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(item["color"].removeprefix("#"))
        point.format.line.fill.background()

    labels = content["data_labels"]
    plot.has_data_labels = labels["enabled"]
    if labels["enabled"]:
        data_labels = plot.data_labels
        data_labels.show_category_name = labels["show_category"]
        data_labels.show_value = labels["show_value"]
        data_labels.show_percentage = labels["show_percentage"]
        data_labels.position = _LABEL_POSITIONS[labels["position"]]
        data_labels.number_format = labels["number_format"]
        data_labels.font.size = Pt(labels["font_size"])
        data_labels.font.bold = labels["font_weight"] >= 600
        data_labels.font.color.rgb = RGBColor.from_string(labels["color"].removeprefix("#"))
    return frame


def _apply_axes(chart: Any, contract: dict[str, Any]) -> None:
    category_contract = contract["category"]
    category_axis = chart.category_axis
    category_axis.visible = category_contract["visible"]
    set_axis_reverse_order(category_axis, category_contract["reverse_order"])
    category_axis.tick_label_position = _TICK_LABEL_POSITIONS[
        category_contract["label_position"]
    ]
    _set_font(category_axis.tick_labels.font, category_contract)
    set_axis_position(category_axis, category_contract["position"])
    set_chart_format_line(category_axis.format, category_contract["line"])

    value_contract = contract["value"]
    value_axis = chart.value_axis
    value_axis.visible = value_contract["visible"]
    value_axis.minimum_scale = value_contract["minimum"]
    value_axis.maximum_scale = value_contract["maximum"]
    value_axis.major_unit = value_contract["major_unit"]
    value_axis.tick_labels.number_format = value_contract["number_format"]
    value_axis.tick_labels.number_format_is_linked = False
    _set_font(value_axis.tick_labels.font, value_contract)
    set_axis_position(value_axis, value_contract["position"])
    set_chart_format_line(value_axis.format, value_contract["line"])
    gridlines = value_contract["major_gridlines"]
    value_axis.has_major_gridlines = gridlines["visible"]
    if gridlines["visible"]:
        set_chart_format_line(value_axis.major_gridlines.format, gridlines["line"])


def _apply_legend(chart: Any, contract: dict[str, Any]) -> None:
    chart.has_legend = contract["enabled"]
    if not contract["enabled"]:
        return
    legend = chart.legend
    legend.position = _LEGEND_POSITIONS[contract["position"]]
    legend.include_in_layout = contract["overlay"]
    _set_font(legend.font, contract)


def _apply_cartesian_labels(plot: Any, contract: dict[str, Any]) -> None:
    plot.has_data_labels = contract["enabled"]
    if not contract["enabled"]:
        return
    labels = plot.data_labels
    labels.show_category_name = contract["show_category"]
    labels.show_series_name = contract["show_series_name"]
    labels.show_value = contract["show_value"]
    labels.show_percentage = False
    labels.position = _LABEL_POSITIONS[contract["position"]]
    labels.number_format = contract["number_format"]
    labels.number_format_is_linked = False
    _set_font(labels.font, contract)


def _apply_line_series(series: Any, contract: dict[str, Any]) -> None:
    set_chart_format_line(
        series.format,
        {
            "color": contract["color"],
            "width": contract["line"]["width"],
            "dash": contract["line"]["dash"],
            "opacity": 1,
        },
    )
    marker_contract = contract["marker"]
    marker = series.marker
    marker.style = _MARKER_STYLES[marker_contract["style"]]
    marker.size = marker_contract["size"]
    marker.format.fill.solid()
    marker.format.fill.fore_color.rgb = RGBColor.from_string(
        marker_contract["fill"].removeprefix("#")
    )
    set_chart_format_line(
        marker.format,
        {
            "color": marker_contract["line_color"],
            "width": marker_contract["line_width"],
            "dash": "solid",
            "opacity": 1,
        },
    )
    series.smooth = False
    set_series_smooth(series, False)


def _render_cartesian_chart(element: dict[str, Any], context: RenderContext) -> Any:
    content = element["content"]
    style = element["style"]
    chart_data = CategoryChartData()
    chart_data.categories = content["categories"]
    for item in content["series"]:
        chart_data.add_series("" if item["name"] is None else item["name"], item["values"])
    frame = _add_chart(
        element,
        context,
        _CARTESIAN_CHART_TYPES[(content["chart_type"], content["grouping"])],
        chart_data,
    )
    chart = frame.chart
    chart.has_title = False
    set_chart_area_style(
        chart._chartSpace,
        style["chart_area"],
        insert_before=("c:txPr", "c:externalData", "c:printSettings", "c:userShapes"),
    )
    set_chart_area_style(chart._chartSpace.chart.plotArea, style["plot_area"])
    set_display_blanks_as(chart, content["display_blanks_as"])
    plot = chart.plots[0]
    plot.vary_by_categories = False
    if content["chart_type"] in {"column", "bar"}:
        plot.gap_width = style["gap_width"]
        plot.overlap = style["overlap"]
        set_chart_value(plot._element, "c:overlap", style["overlap"])
        for series, item in zip(plot.series, content["series"], strict=True):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(
                item["color"].removeprefix("#")
            )
            series.format.line.fill.background()
    else:
        for series, item in zip(plot.series, content["series"], strict=True):
            _apply_line_series(series, item)
    _apply_axes(chart, content["axes"])
    _apply_legend(chart, content["legend"])
    _apply_cartesian_labels(plot, content["data_labels"])
    return frame


class ChartRenderer:
    kind = "chart"
    supported_fields = _CHART_FIELDS
    supported_values = {"chart_type": CANONICAL_VALUES["chart_type"]}
    required_fields = frozenset({"chart_type", "data_labels"})
    capability_ids = frozenset(
        f"chart.{value}" for value in CANONICAL_VALUES["chart_type"]
    ) | frozenset(
        capability
        for capability, field in ATOMIC_CAPABILITY_METADATA.items()
        if capability.startswith("chart.") and field in _CHART_FIELDS
    )

    def validate_contract(
        self, element: dict[str, Any], context: RenderContext
    ) -> list[ContractIssue]:
        return validate_chart_contract(element)

    def render(self, element: dict[str, Any], context: RenderContext) -> None:
        chart_type = element["content"]["chart_type"]
        frame = (
            _render_pie_chart(element, context)
            if chart_type in _PIE_CHART_TYPES
            else _render_cartesian_chart(element, context)
        )
        context.registry.register(
            element["element_id"],
            frame,
            "graphicFrame",
            semantic_kind="chart",
            selected_mode=context.representation_modes[element["element_id"]],
        )


CHART_RENDERER = ChartRenderer()
register_renderer("chart", CHART_RENDERER)
