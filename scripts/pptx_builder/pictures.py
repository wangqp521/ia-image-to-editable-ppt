"""Lossless picture placement and current alpha-isolation icon renderer."""

from __future__ import annotations

import hashlib
from typing import Any

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.parts.image import ImagePart
from pptx.util import Emu

from lib.capabilities import ATOMIC_CAPABILITY_METADATA, CANONICAL_VALUES
from lib.error_codes import ContractIssue, ToolError
from lib.geometry import (
    DRAWINGML_PERCENT_SCALE,
    quantize_drawingml_percentage,
)
from lib.representation_contracts import require_asset
from lib.schema_contracts import (
    ICON_ITEM_FIELDS,
    ICON_MODULE_FIELDS,
    unknown_field_detail,
)

from .common import RenderContext, register_renderer
from .shapes import _issue, _number


_CROP_FIELDS = {"left", "top", "right", "bottom"}
_ICON_MODULE_FIELDS = ICON_MODULE_FIELDS
_ICON_RECORD_FIELDS = ICON_ITEM_FIELDS


def _validate_crop(value: Any, path: str) -> list[ContractIssue]:
    if not isinstance(value, dict):
        return [_issue(path, "crop must be an object")]
    unknown = sorted(set(value) - _CROP_FIELDS)
    missing = sorted(_CROP_FIELDS - set(value))
    if unknown or missing:
        detail = f"unknown crop fields: {', '.join(unknown)}" if unknown else f"missing crop fields: {', '.join(missing)}"
        return [_issue(path, detail)]
    if any(not _number(value[side]) or not 0 <= value[side] < 1 for side in _CROP_FIELDS):
        return [_issue(path, "crop fractions must be from 0 (inclusive) to 1 (exclusive)")]
    if value["left"] + value["right"] >= 1 or value["top"] + value["bottom"] >= 1:
        return [_issue(path, "opposite crop fractions must sum to less than 1")]
    return []


def _icon_record(spec: dict[str, Any], element_id: str, path: str) -> dict[str, Any]:
    module = spec.get("modules", {}).get("icons")
    module_unknown = unknown_field_detail("IconsModule", module)
    if module_unknown is not None:
        raise ToolError("UNSUPPORTED_CAPABILITY", "modules.icons", module_unknown)
    if not isinstance(module, dict) or set(module) != _ICON_MODULE_FIELDS:
        raise ToolError("UNSUPPORTED_CAPABILITY", "modules.icons", "icons module fields must exactly match the current schema")
    icons = module.get("icons") if isinstance(module, dict) else None
    if not isinstance(icons, list):
        raise ToolError("UNSUPPORTED_CAPABILITY", "modules.icons.icons", "icon renderer requires the current icons module")
    records = [item for item in icons if isinstance(item, dict) and item.get("element_id") == element_id]
    if len(records) != 1:
        raise ToolError("UNSUPPORTED_CAPABILITY", "modules.icons.icons", "icon element must have exactly one asset record")
    record = records[0]
    record_unknown = unknown_field_detail("IconItem", record)
    if record_unknown is not None:
        raise ToolError("UNSUPPORTED_CAPABILITY", path, record_unknown)
    if set(record) != _ICON_RECORD_FIELDS:
        raise ToolError("UNSUPPORTED_CAPABILITY", path, "icon asset record fields must exactly match the current schema")
    if record.get("crop_mode") != "alpha_isolation":
        raise ToolError("UNSUPPORTED_CAPABILITY", f"{path}.crop_mode", "icon crop mode must be alpha_isolation")
    return record


def _automatic_crop(mode: str, pixel_size: tuple[int, int], bbox: list[int]) -> dict[str, float]:
    crop = {side: 0.0 for side in _CROP_FIELDS}
    if mode == "none":
        return crop
    image_ratio = pixel_size[0] / pixel_size[1]
    box_ratio = bbox[2] / bbox[3]
    if mode == "cover":
        if image_ratio > box_ratio:
            amount = (1 - box_ratio / image_ratio) / 2
            crop["left"] = crop["right"] = amount
        elif image_ratio < box_ratio:
            amount = (1 - image_ratio / box_ratio) / 2
            crop["top"] = crop["bottom"] = amount
    elif image_ratio > box_ratio:
        amount = (image_ratio / box_ratio - 1) / 2
        crop["top"] = crop["bottom"] = -amount
    elif image_ratio < box_ratio:
        amount = (box_ratio / image_ratio - 1) / 2
        crop["left"] = crop["right"] = -amount
    return crop


def _quantized_picture_crop(
    mode: str,
    pixel_size: tuple[int, int],
    bbox: list[int],
    explicit: dict[str, float],
    path: str,
) -> dict[str, int]:
    automatic = _automatic_crop(mode, pixel_size, bbox)
    quantized = {
        side: quantize_drawingml_percentage(automatic[side] + explicit[side])
        for side in _CROP_FIELDS
    }
    if (
        any(value >= DRAWINGML_PERCENT_SCALE for value in quantized.values())
        or quantized["left"] + quantized["right"] >= DRAWINGML_PERCENT_SCALE
        or quantized["top"] + quantized["bottom"] >= DRAWINGML_PERCENT_SCALE
    ):
        raise ToolError(
            "UNSUPPORTED_CAPABILITY",
            path,
            "quantized crop sides and opposite-side sums must be below 100000",
            "picture.crop.explicit",
        )
    return quantized


def _set_picture_crop(picture: Any, crop: dict[str, int]) -> None:
    for side in ("left", "top", "right", "bottom"):
        setattr(
            picture,
            f"crop_{side}",
            crop[side] / DRAWINGML_PERCENT_SCALE,
        )
    src_rect = picture._element.blipFill.srcRect
    for side, attribute in (("left", "l"), ("top", "t"), ("right", "r"), ("bottom", "b")):
        src_rect.set(attribute, str(crop[side]))


def _set_picture_opacity(picture: Any, opacity: float) -> None:
    blip = picture._element.blipFill.blip
    for node in list(blip.findall(qn("a:alphaModFix"))):
        blip.remove(node)
    if opacity != 1:
        alpha = OxmlElement("a:alphaModFix")
        alpha.set("amt", str(quantize_drawingml_percentage(opacity)))
        blip.append(alpha)


def _add_webp_picture(
    slide: Any,
    asset_path: Any,
    x: int,
    y: int,
    width: int,
    height: int,
    path: str,
) -> Any:
    """Add an original WEBP blob without converting its media container."""
    try:
        blob = asset_path.read_bytes()
    except ToolError:
        raise
    except (OSError, RuntimeError, ValueError) as exc:
        raise ToolError(
            "UNSUPPORTED_CAPABILITY",
            path,
            "WEBP asset cannot be read for embedding",
            "picture.asset.local_hash",
        ) from exc
    package = slide.part.package
    sha1 = hashlib.sha1(blob).hexdigest()
    image_part = next(
        (
            part
            for part in package.iter_parts()
            if isinstance(part, ImagePart) and part.sha1 == sha1
        ),
        None,
    )
    if image_part is None:
        image_part = ImagePart(
            package.next_image_partname("webp"),
            "image/webp",
            package,
            blob,
            asset_path.name,
        )
    r_id = slide.part.relate_to(image_part, RT.IMAGE)
    shapes = slide.shapes
    picture_element = shapes._add_pic_from_image_part(
        image_part,
        r_id,
        Emu(x),
        Emu(y),
        Emu(width),
        Emu(height),
    )
    shapes._recalculate_extents()
    return shapes._shape_factory(picture_element)


class PictureRenderer:
    supported_fields = frozenset({"asset", "mode", "crop", "rotation", "opacity"})
    supported_values = {"picture_mode": CANONICAL_VALUES["picture_mode"]}
    required_fields = frozenset({"asset", "mode", "crop"})

    def __init__(self, kind: str) -> None:
        self.kind = kind
        self.capability_ids = frozenset(
            f"picture_mode.{value}" for value in CANONICAL_VALUES["picture_mode"]
        ) | frozenset(
            capability for capability, field in ATOMIC_CAPABILITY_METADATA.items()
            if field in self.supported_fields and capability.startswith("picture.")
        )

    def validate_contract(self, element: dict[str, Any], context: RenderContext) -> list[ContractIssue]:
        element_id = element.get("element_id", "<unknown>")
        path = f"elements.{element_id}"
        content = element.get("content", {})
        style = element.get("style", {})
        missing = sorted(self.required_fields - (set(content) | set(style)))
        if missing:
            return [_issue(path, f"missing picture fields: {', '.join(missing)}")]
        mode = content.get("mode")
        if mode not in CANONICAL_VALUES["picture_mode"]:
            return [_issue(f"{path}.content.mode", "unsupported picture mode", f"picture_mode.{mode}")]
        issues = _validate_crop(content.get("crop"), f"{path}.content.crop")
        if issues:
            return issues
        rotation = style.get("rotation", 0)
        opacity = style.get("opacity", 1)
        if not _number(rotation) or not -360 <= rotation <= 360:
            return [_issue(f"{path}.style.rotation", "rotation must be from -360 to 360 degrees")]
        if not _number(opacity) or not 0 <= opacity <= 1:
            return [_issue(f"{path}.style.opacity", "opacity must be from 0 to 1")]
        selected_mode = context.representation_modes.get(element_id)
        required_modes = {"asset"} if self.kind == "icon" else {
            "asset",
            "background_picture",
        }
        if selected_mode not in required_modes:
            return [_issue(f"{path}.representation", "picture renderer mode is invalid")]
        asset_path, asset_hash, pixel_size = require_asset(
            content.get("asset"), f"{path}.content.asset"
        )
        _quantized_picture_crop(
            mode,
            pixel_size,
            element["slide_bbox"],
            content["crop"],
            f"{path}.content.crop",
        )
        if self.kind == "icon":
            record_path = f"modules.icons.icons.{element_id}"
            record = _icon_record(context.spec, element_id, record_path)
            if record.get("asset_path") != str(asset_path) or str(record.get("asset_sha256", "")).lower() != asset_hash.lower():
                return [_issue(record_path, "icon asset record must match element asset path and hash")]
        return []

    def render(self, element: dict[str, Any], context: RenderContext) -> None:
        element_id = element["element_id"]
        content = element["content"]
        style = element["style"]
        asset_path, asset_hash, pixel_size = require_asset(content["asset"], f"elements.{element_id}.content.asset")
        if self.kind == "icon":
            record = _icon_record(context.spec, element_id, f"modules.icons.icons.{element_id}")
            if record.get("asset_path") != str(asset_path) or str(record.get("asset_sha256", "")).lower() != asset_hash.lower():
                raise ToolError("UNSUPPORTED_CAPABILITY", f"modules.icons.icons.{element_id}", "icon asset record must match element asset path and hash")
        bbox = element["slide_bbox"]
        crop = _quantized_picture_crop(
            content["mode"],
            pixel_size,
            bbox,
            content["crop"],
            f"elements.{element_id}.content.crop",
        )
        x, y, width, height = bbox
        if asset_path.suffix.lower() == ".webp":
            picture = _add_webp_picture(
                context.slide,
                asset_path,
                x,
                y,
                width,
                height,
                f"elements.{element_id}.content.asset.path",
            )
        else:
            picture = context.slide.shapes.add_picture(
                str(asset_path), Emu(x), Emu(y), Emu(width), Emu(height)
            )
        _set_picture_crop(picture, crop)
        picture.rotation = style.get("rotation", 0)
        _set_picture_opacity(picture, style.get("opacity", 1))
        context.registry.register(
            element_id, picture, "pic", semantic_kind=self.kind,
            selected_mode=context.representation_modes[element_id], media_sha256=asset_hash,
        )


PICTURE_RENDERER = PictureRenderer("picture")
ICON_RENDERER = PictureRenderer("icon")
register_renderer("picture", PICTURE_RENDERER)
register_renderer("icon", ICON_RENDERER)
