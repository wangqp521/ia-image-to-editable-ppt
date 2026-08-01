"""Native PowerPoint table renderer with exact cell contracts."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from lib.element_contracts import validate_element_contract
from lib.error_codes import ContractIssue

from .common import RenderContext, register_renderer
from .ooxml import set_table_cell_border


_ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
_VERTICAL_ALIGNMENTS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}
_SIDES = ("left", "right", "top", "bottom")


def _issue(path: str, detail: str) -> ContractIssue:
    return ContractIssue("UNSUPPORTED_CAPABILITY", path, detail)


def _set_font_family(run: Any, family: str) -> None:
    run.font.name = family
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
        node = properties.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            properties.append(node)
        node.set("typeface", family)


def _apply_cell(cell: Any, contract: dict[str, Any], path: str) -> None:
    fill = contract["fill"]
    if fill == "noFill":
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill[1:])

    margins = contract["margins"]
    for side in _SIDES:
        setattr(cell, f"margin_{side}", Emu(margins[side]))
    cell.vertical_anchor = _VERTICAL_ALIGNMENTS[contract["vertical_alignment"]]
    cell.text_frame.clear()
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = _ALIGNMENTS[contract["alignment"]]
    run = paragraph.add_run()
    run.text = contract["text"]
    font = contract["font"]
    _set_font_family(run, font["name"])
    run.font.size = Pt(font["size"])
    run.font.bold = font["weight"] >= 600
    run.font.italic = font["italic"]
    run.font.color.rgb = RGBColor.from_string(font["color"][1:])

    properties = cell._tc.tcPr
    borders = contract["borders"]
    for side in _SIDES:
        set_table_cell_border(
            properties,
            side,
            borders.get(side),
            f"{path}.borders.{side}",
        )


class TableRenderer:
    kind = "table"
    supported_fields = frozenset({"rotation", "rows", "columns", "cells"})
    supported_values: dict[str, frozenset[str]] = {}
    required_fields = frozenset({"rows", "columns", "cells"})
    capability_ids = frozenset({"table.merge", "table.cell.local_border"})

    def validate_contract(
        self, element: dict[str, Any], context: RenderContext
    ) -> list[ContractIssue]:
        issues = validate_element_contract(element)
        if issues:
            return issues
        element_id = element["element_id"]
        if context.representation_modes.get(element_id) != "native":
            return [_issue(
                f"elements.{element_id}.representation",
                "table renderer requires native mode",
            )]
        return []

    def render(self, element: dict[str, Any], context: RenderContext) -> None:
        content = element["content"]
        rows = content["rows"]
        columns = content["columns"]
        x, y, width, height = element["slide_bbox"]
        shape = context.slide.shapes.add_table(
            len(rows), len(columns), Emu(x), Emu(y), Emu(width), Emu(height)
        )
        table = shape.table
        for index, value in enumerate(rows):
            table.rows[index].height = Emu(value)
        for index, value in enumerate(columns):
            table.columns[index].width = Emu(value)

        for row_index in range(len(rows)):
            for column_index in range(len(columns)):
                properties = table.cell(row_index, column_index)._tc.tcPr
                for side in _SIDES:
                    set_table_cell_border(
                        properties,
                        side,
                        None,
                        f"elements.{element['element_id']}.content.cells.{side}",
                    )

        fonts: set[str] = set()
        texts: list[str] = []
        for index, contract in enumerate(content["cells"]):
            row = contract["row"]
            column = contract["column"]
            row_span = contract["row_span"]
            column_span = contract["column_span"]
            cell = table.cell(row, column)
            if row_span > 1 or column_span > 1:
                cell.merge(
                    table.cell(row + row_span - 1, column + column_span - 1)
                )
            _apply_cell(
                cell,
                contract,
                f"elements.{element['element_id']}.content.cells[{index}]",
            )
            fonts.add(contract["font"]["name"])
            texts.append(contract["text"])
        shape.rotation = element["style"].get("rotation", 0)
        context.registry.register(
            element["element_id"],
            shape,
            "graphicFrame",
            semantic_kind="table",
            selected_mode=context.representation_modes[element["element_id"]],
            text_summary="\n".join(texts),
            font_declarations=tuple(sorted(fonts)),
        )


TABLE_RENDERER = TableRenderer()
register_renderer("table", TABLE_RENDERER)
