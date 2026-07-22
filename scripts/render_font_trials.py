#!/usr/bin/env python3
"""Render explicitly requested font trials and record traceable measurements."""

from __future__ import annotations

import argparse
import hashlib
import itertools
import json
import os
import subprocess
from pathlib import Path
from typing import Any, Iterable, Sequence

from PIL import Image, ImageChops, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


SLIDE_SIZE_IN = (13.333333, 7.5)
TRIAL_ORIGIN_IN = (1.0, 1.0)
RENDER_DPI = 144


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _run(
    command: Sequence[str],
    *,
    env: dict[str, str] | None = None,
) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        list(command),
        check=True,
        capture_output=True,
        text=True,
        env=env,
        timeout=120,
    )


def parse_pdffonts(output: str) -> list[str]:
    """Return unique PDF font names in their reported order."""
    fonts: list[str] = []
    for raw_line in output.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("name ") or set(line) <= {"-", " "}:
            continue
        name = line.split()[0]
        if name not in fonts:
            fonts.append(name)
    return fonts


def _foreground_mask(image: Image.Image) -> Image.Image:
    rgb = image.convert("RGB")
    background = Image.new("RGB", rgb.size, "white")
    difference = ImageChops.difference(rgb, background).convert("L")
    return difference.point(lambda value: 255 if value >= 20 else 0, mode="1")


def measure_rendered_text(
    image: Image.Image,
    box: tuple[int, int, int, int],
) -> dict[str, Any]:
    """Measure visible non-white ink inside one explicit pixel box."""
    left, top, right, bottom = (int(value) for value in box)
    if not (0 <= left < right <= image.width and 0 <= top < bottom <= image.height):
        raise ValueError("box must stay inside rendered image")

    mask = _foreground_mask(image)
    crop = mask.crop((left, top, right, bottom))
    local_bbox = crop.getbbox()
    if local_bbox is None:
        return {
            "ink_bbox_px": None,
            "ink_size_px": [0, 0],
            "line_count": 0,
            "clipped": False,
            "foreground_pixel_ratio": 0.0,
        }

    ink_bbox = [
        left + local_bbox[0],
        top + local_bbox[1],
        left + local_bbox[2],
        top + local_bbox[3],
    ]
    occupied_rows = []
    for y in range(crop.height):
        occupied_rows.append(crop.crop((0, y, crop.width, y + 1)).getbbox() is not None)
    line_count = 0
    inside_run = False
    for occupied in occupied_rows:
        if occupied and not inside_run:
            line_count += 1
        inside_run = occupied

    clipped = (
        local_bbox[0] == 0
        or local_bbox[1] == 0
        or local_bbox[2] == crop.width
        or local_bbox[3] == crop.height
    )
    foreground = sum(1 for value in crop.get_flattened_data() if value)
    return {
        "ink_bbox_px": ink_bbox,
        "ink_size_px": [ink_bbox[2] - ink_bbox[0], ink_bbox[3] - ink_bbox[1]],
        "line_count": line_count,
        "clipped": clipped,
        "foreground_pixel_ratio": round(foreground / (crop.width * crop.height), 6),
    }


def _create_trial_pptx(
    text: str,
    font: str,
    size_pt: float,
    box_in: tuple[float, float],
    trial_dir: Path,
) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_SIZE_IN[0])
    presentation.slide_height = Inches(SLIDE_SIZE_IN[1])
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(
        Inches(TRIAL_ORIGIN_IN[0]),
        Inches(TRIAL_ORIGIN_IN[1]),
        Inches(box_in[0]),
        Inches(box_in[1]),
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(float(size_pt))
    run.font.color.rgb = RGBColor(0, 0, 0)

    output = trial_dir / "trial.pptx"
    presentation.save(output)
    return output


def _render_pdf(
    pptx: Path,
    trial_dir: Path,
    fontconfig: Path,
    soffice: str,
) -> Path:
    if not fontconfig.is_file():
        raise FileNotFoundError(f"fontconfig not found: {fontconfig}")
    profile_dir = trial_dir / "libreoffice-profile"
    profile_dir.mkdir(parents=True, exist_ok=True)
    env = os.environ.copy()
    env["FONTCONFIG_FILE"] = str(fontconfig)
    _run(
        [
            soffice,
            "--headless",
            f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(trial_dir),
            str(pptx),
        ],
        env=env,
    )
    pdf = trial_dir / "trial.pdf"
    if not pdf.is_file():
        raise RuntimeError("LibreOffice did not produce trial.pdf")
    return pdf


def _render_png(pdf: Path, trial_dir: Path, pdftoppm: str) -> Path:
    prefix = trial_dir / "render"
    _run(
        [pdftoppm, "-png", "-singlefile", "-r", str(RENDER_DPI), str(pdf), str(prefix)]
    )
    png = trial_dir / "render.png"
    if not png.is_file():
        raise RuntimeError("pdftoppm did not produce render.png")
    return png


def _trial_box_px(
    rendered_size: tuple[int, int],
    box_in: tuple[float, float],
) -> tuple[int, int, int, int]:
    scale_x = rendered_size[0] / SLIDE_SIZE_IN[0]
    scale_y = rendered_size[1] / SLIDE_SIZE_IN[1]
    left = round(TRIAL_ORIGIN_IN[0] * scale_x)
    top = round(TRIAL_ORIGIN_IN[1] * scale_y)
    right = round((TRIAL_ORIGIN_IN[0] + box_in[0]) * scale_x)
    bottom = round((TRIAL_ORIGIN_IN[1] + box_in[1]) * scale_y)
    return left, top, right, bottom


def _build_contact_sheet(trials: list[dict[str, Any]], output: Path) -> Path:
    rows: list[Image.Image] = []
    for trial in trials:
        with Image.open(trial["png"]) as opened:
            image = opened.convert("RGB")
            box = tuple(trial["measurement_box_px"])
            crop = image.crop(box)
        crop.thumbnail((960, 180), Image.Resampling.LANCZOS)
        row = Image.new("RGB", (1000, max(220, crop.height + 50)), "white")
        row.paste(crop, (20, 40))
        ImageDraw.Draw(row).text(
            (20, 12),
            f'{trial["trial_id"]} | {trial["requested_font"]} | {trial["size_pt"]} pt',
            fill="black",
        )
        rows.append(row)
    sheet = Image.new("RGB", (1000, sum(row.height for row in rows)), "white")
    y = 0
    for row in rows:
        sheet.paste(row, (0, y))
        y += row.height
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)
    return output


def render_font_trials(
    *,
    text: str,
    fonts: Iterable[str],
    sizes_pt: Iterable[float],
    box_in: tuple[float, float],
    output_dir: Path | str,
    fontconfig: Path | str,
    soffice: str = "soffice",
    pdftoppm: str = "pdftoppm",
    pdffonts: str = "pdffonts",
) -> dict[str, Any]:
    """Render only the requested font/size combinations; never choose a winner."""
    fonts = [str(font) for font in fonts]
    sizes_pt = [float(size) for size in sizes_pt]
    if not text or not fonts or not sizes_pt:
        raise ValueError("text, fonts, and sizes_pt must be non-empty")
    if len(box_in) != 2 or box_in[0] <= 0 or box_in[1] <= 0:
        raise ValueError("box_in must contain positive width and height")

    output_dir = Path(output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    fontconfig = Path(fontconfig).resolve()
    trials: list[dict[str, Any]] = []

    for index, (font, size_pt) in enumerate(
        itertools.product(fonts, sizes_pt), start=1
    ):
        trial_id = f"trial-{index:03d}"
        trial_dir = output_dir / trial_id
        trial_dir.mkdir(parents=True, exist_ok=True)
        pptx = _create_trial_pptx(text, font, size_pt, box_in, trial_dir)
        pdf = _render_pdf(pptx, trial_dir, fontconfig, soffice)
        png = _render_png(pdf, trial_dir, pdftoppm)
        resolved_fonts = parse_pdffonts(_run([pdffonts, str(pdf)]).stdout)
        with Image.open(png) as opened:
            rendered = opened.convert("RGB")
            measurement_box = _trial_box_px(rendered.size, box_in)
            metrics = measure_rendered_text(rendered, measurement_box)
        trials.append(
            {
                "trial_id": trial_id,
                "requested_font": font,
                "resolved_fonts": resolved_fonts,
                "size_pt": size_pt,
                "box_in": list(box_in),
                "measurement_box_px": list(measurement_box),
                "pptx": str(pptx),
                "pptx_sha256": _sha256(pptx),
                "pdf": str(pdf),
                "pdf_sha256": _sha256(pdf),
                "png": str(png),
                "png_sha256": _sha256(png),
                **metrics,
            }
        )

    contact_sheet = _build_contact_sheet(trials, output_dir / "contact-sheet.png")
    report: dict[str, Any] = {
        "text": text,
        "fontconfig": str(fontconfig),
        "box_in": list(box_in),
        "render_dpi": RENDER_DPI,
        "trials": trials,
        "contact_sheet": str(contact_sheet),
        "contact_sheet_sha256": _sha256(contact_sheet),
    }
    report_path = output_dir / "font-trials.json"
    report_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return report


def _parse_box(value: str) -> tuple[float, float]:
    parts = value.split(",")
    if len(parts) != 2:
        raise argparse.ArgumentTypeError("box must be WIDTH,HEIGHT in inches")
    try:
        return float(parts[0]), float(parts[1])
    except ValueError as error:
        raise argparse.ArgumentTypeError("box values must be numbers") from error


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--text", required=True)
    parser.add_argument("--font", action="append", required=True)
    parser.add_argument("--size-pt", action="append", type=float, required=True)
    parser.add_argument("--box-in", type=_parse_box, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument("--fontconfig", type=Path, required=True)
    parser.add_argument("--soffice", default="soffice")
    parser.add_argument("--pdftoppm", default="pdftoppm")
    parser.add_argument("--pdffonts", default="pdffonts")
    args = parser.parse_args()
    report = render_font_trials(
        text=args.text,
        fonts=args.font,
        sizes_pt=args.size_pt,
        box_in=args.box_in,
        output_dir=args.output_dir,
        fontconfig=args.fontconfig,
        soffice=args.soffice,
        pdftoppm=args.pdftoppm,
        pdffonts=args.pdffonts,
    )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
